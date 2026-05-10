"""
Generate the cold-start client alignment deck as native PowerPoint shapes.

Usage:
    python3 build_client_deck_v2.py

Produces: client_alignment_deck_v2.pptx alongside this script.

Architecture summary:
  - Cold-start: every catalog (incl. business_capability) discovered from this
    discipline's SOP corpus. No universal lab_stage artifact.
  - 4 agent classes: Catalog Builder, Tagger, Story Extractor, Sampled Validator.
  - Catalog Builder runs 5 times: business_capability, theme, epic, persona, test.
  - Top-down hierarchical discovery: capability ← theme ← epic.
  - HITL is a non-blocking review queue: every emission goes to corpus AND is
    logged for review; reviewer post-hoc can confirm/correct/reject; pipeline
    never waits. Corrections are applied to the corpus manually in milestone 1.
  - Phase 3 (cross-SOP synthesis) deferred to milestone 2.
  - Slug-based labels everywhere (no T01/E04/P02 nomenclature).

11 slides:
  1.  What we are building
  2.  Inputs — SOPs only
  3.  End-to-end flow (3 phases)
  4.  The 4 agents
  5.  Phase 1 — Catalog Build (cold, hierarchical, 5 catalogs)
  6.  Catalog schemas + examples (capability / theme / epic / persona / test)
  7.  Phase 2 — Per-SOP Extraction (Tagger + Story Extractor)
  8.  Story schema + 4 shapes (with examples)
  9.  Sampled Validator QC
  10. HITL — non-blocking review queue
  11. Milestone 1 ships + roadmap
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---- palette --------------------------------------------------------------
BOX_FILL = RGBColor(0xE7, 0xEF, 0xF8)
BOX_LINE = RGBColor(0x4A, 0x6E, 0xA0)
ACCENT = RGBColor(0x1F, 0x4E, 0x79)
TEXT_DARK = RGBColor(0x1F, 0x2D, 0x3D)
TEXT_GREY = RGBColor(0x55, 0x60, 0x6E)
ARROW_COLOR = RGBColor(0x4A, 0x6E, 0xA0)
SYS_FILL = RGBColor(0xFA, 0xE9, 0xCB)
SYS_LINE = RGBColor(0xC9, 0x95, 0x2C)
DECISION_FILL = RGBColor(0xFC, 0xE7, 0xC4)
DECISION_LINE = RGBColor(0xC9, 0x95, 0x2C)
START_FILL = RGBColor(0xE0, 0xEC, 0xF8)
START_LINE = RGBColor(0x4A, 0x6E, 0xA0)
END_FILL = RGBColor(0xDC, 0xE9, 0xD7)
END_LINE = RGBColor(0x5C, 0x83, 0x47)
OUT_FILL = RGBColor(0xDC, 0xE9, 0xD7)
OUT_LINE = RGBColor(0x5C, 0x83, 0x47)
PHASE_FILL = RGBColor(0xEE, 0xE7, 0xF8)
PHASE_LINE = RGBColor(0x6E, 0x4F, 0xA0)
CODE_FILL = RGBColor(0xF6, 0xF6, 0xF6)
CODE_LINE = RGBColor(0xC0, 0xC8, 0xD0)
SIDE_FILL = RGBColor(0xEC, 0xF4, 0xE8)
SIDE_LINE = RGBColor(0x5C, 0x83, 0x47)
PARK_FILL = RGBColor(0xF8, 0xE3, 0xE3)
PARK_LINE = RGBColor(0xA8, 0x4A, 0x4A)
CARD_FILL = RGBColor(0xFB, 0xFC, 0xFE)
CARD_LINE = RGBColor(0xC8, 0xD2, 0xDD)
FUTURE_FILL = RGBColor(0xF0, 0xEC, 0xF7)
FUTURE_LINE = RGBColor(0x6E, 0x4F, 0xA0)


# ---- text/shape helpers ---------------------------------------------------

def _set_text(shape, text, font_size, bold, color, align,
              anchor=MSO_ANCHOR.MIDDLE, font_name="Calibri", italic=False):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name


def add_box(slide, left, top, width, height, text, fill=BOX_FILL, line=BOX_LINE,
            font_size=12, bold=False, font_color=TEXT_DARK, align=PP_ALIGN.CENTER,
            shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    box = slide.shapes.add_shape(shape_type, Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line
    box.line.width = Pt(1.0)
    _set_text(box, text, font_size, bold, font_color, align)
    return box


def add_diamond(slide, left, top, width, height, text,
                fill=DECISION_FILL, line=DECISION_LINE,
                font_size=11, bold=True, color=TEXT_DARK):
    d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(left), Inches(top),
                               Inches(width), Inches(height))
    d.fill.solid()
    d.fill.fore_color.rgb = fill
    d.line.color.rgb = line
    d.line.width = Pt(1.25)
    _set_text(d, text, font_size, bold, color, PP_ALIGN.CENTER)
    return d


def add_node(slide, left, top, width, height, text, kind="process",
             font_size=11, bold=False):
    """kind in {start, end, process, agent, side, park, future}."""
    if kind == "start":
        return add_box(slide, left, top, width, height, text,
                       fill=START_FILL, line=START_LINE,
                       font_size=font_size, bold=True)
    if kind == "end":
        return add_box(slide, left, top, width, height, text,
                       fill=END_FILL, line=END_LINE,
                       font_size=font_size, bold=True)
    if kind == "agent":
        return add_box(slide, left, top, width, height, text,
                       fill=SYS_FILL, line=SYS_LINE,
                       font_size=font_size, bold=True)
    if kind == "side":
        return add_box(slide, left, top, width, height, text,
                       fill=SIDE_FILL, line=SIDE_LINE,
                       font_size=font_size, bold=bold)
    if kind == "park":
        return add_box(slide, left, top, width, height, text,
                       fill=PARK_FILL, line=PARK_LINE,
                       font_size=font_size, bold=bold)
    if kind == "future":
        return add_box(slide, left, top, width, height, text,
                       fill=FUTURE_FILL, line=FUTURE_LINE,
                       font_size=font_size, bold=bold)
    if kind == "card":
        return add_box(slide, left, top, width, height, text,
                       fill=CARD_FILL, line=CARD_LINE,
                       font_size=font_size, bold=bold)
    return add_box(slide, left, top, width, height, text,
                   font_size=font_size, bold=bold)


def add_text(slide, left, top, width, height, text, font_size=14,
             bold=False, italic=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font_name="Calibri"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return tb


def add_bullets(slide, left, top, width, height, items, font_size=11,
                color=TEXT_DARK, bullet="•  "):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = bullet + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
        p.space_after = Pt(3)
    return tb


def add_arrow(slide, x1, y1, x2, y2, color=ARROW_COLOR, weight=2.0,
              connector_type=MSO_CONNECTOR.STRAIGHT, label=None,
              label_offset=(0, -0.18), label_color=TEXT_GREY,
              label_size=10, label_bold=False):
    conn = slide.shapes.add_connector(connector_type,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    line_elem = conn.line._get_or_add_ln()
    tail_end = etree.SubElement(line_elem, qn("a:tailEnd"))
    tail_end.set("type", "triangle")
    tail_end.set("w", "med")
    tail_end.set("h", "med")
    if label:
        lx = (x1 + x2) / 2 + label_offset[0]
        ly = (y1 + y2) / 2 + label_offset[1]
        tb = slide.shapes.add_textbox(Inches(lx), Inches(ly),
                                      Inches(0.9), Inches(0.3))
        tf = tb.text_frame
        tf.margin_left = Inches(0.04)
        tf.margin_right = Inches(0.04)
        tf.margin_top = Inches(0.0)
        tf.margin_bottom = Inches(0.0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(label_size)
        run.font.bold = label_bold
        run.font.color.rgb = label_color
        run.font.name = "Calibri"
    return conn


def add_title(slide, text, color=ACCENT):
    add_text(slide, 0.5, 0.32, 12.5, 0.6, text,
             font_size=26, bold=True, color=color)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.5), Inches(0.95),
                                 Inches(1.6), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def add_subtitle(slide, text, y=1.05, italic=True, color=TEXT_GREY, size=12.5):
    # Subtitle box sized to hold up to 2 wrapped lines at 12.5pt without
    # bleeding into content below (which starts at ~y=1.45).
    add_text(slide, 0.5, y, 12.5, 0.42, text,
             font_size=size, italic=italic, color=color)


def add_footer(slide, slide_num, total_slides,
               deck_label="LIMS Story Generation — Client Alignment"):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.5), Inches(7.18),
                                  Inches(12.33), Inches(0.012))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xD0, 0xD6, 0xDC)
    line.line.fill.background()
    add_text(slide, 0.5, 7.22, 8.5, 0.25, deck_label,
             font_size=8, italic=True, color=RGBColor(0x90, 0x97, 0x9F),
             align=PP_ALIGN.LEFT)
    add_text(slide, 11.83, 7.22, 1.0, 0.25,
             f"{slide_num} / {total_slides}",
             font_size=8, color=RGBColor(0x90, 0x97, 0x9F),
             align=PP_ALIGN.RIGHT)


def add_code_box(slide, left, top, width, height, content, font_size=9,
                 title=None, fill=CODE_FILL, line=CODE_LINE,
                 title_color=ACCENT, title_size=10.5):
    if title:
        add_text(slide, left, top - 0.30, width, 0.28, title,
                 font_size=title_size, bold=True, color=title_color)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line
    box.line.width = Pt(0.75)
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.10)
    tf.margin_right = Inches(0.10)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    lines = content.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = ln
        run.font.size = Pt(font_size)
        run.font.color.rgb = TEXT_DARK
        run.font.name = "Consolas"
        p.space_after = Pt(0)
    return box


def add_section_header(slide, left, top, width, text, color=ACCENT, size=13):
    add_text(slide, left, top, width, 0.36, text,
             font_size=size, bold=True, color=color)


# ============================================================================
# SLIDE 1 — What we are building
# ============================================================================

def slide_one(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "What we are building")
    add_subtitle(
        slide,
        "A discipline-agnostic system that turns SOPs into structured stories. "
        "Cold-start: business capabilities, themes, epics, personas, tests are "
        "all discovered from this discipline's SOPs.",
    )

    # 3-box flow: SOPs -> System -> Outputs
    box_w, box_h, y = 3.4, 1.5, 1.85
    x1, x2, x3 = 0.6, 4.95, 9.3
    add_box(slide, x1, y, box_w, box_h,
            "SOPs\n(Standard Operating Procedures)\n— the corpus —",
            font_size=14, bold=True)
    add_box(slide, x2, y, box_w, box_h, "Our system\n4 agent classes",
            fill=SYS_FILL, line=SYS_LINE, font_size=18, bold=True)
    add_box(slide, x3, y, box_w, box_h,
            "Story corpus\n+ 5 catalogs (capability/theme/\n  epic/persona/test)\n+ audit log + review queue",
            fill=OUT_FILL, line=OUT_LINE, font_size=13, bold=True)
    add_arrow(slide, x1 + box_w + 0.05, y + box_h / 2,
              x2 - 0.05, y + box_h / 2)
    add_arrow(slide, x2 + box_w + 0.05, y + box_h / 2,
              x3 - 0.05, y + box_h / 2)

    # Three columns of why
    col_y = 4.0
    add_section_header(slide, 0.5, col_y, 12.5,
                       "Why this matters", size=15)
    bullets_y = col_y + 0.42

    col_w = 4.0
    cols = [
        (0.6, "One system, any discipline",
         [
             "Same machinery for any clinical-lab discipline — Microbiology, Histology, Hematology, Chemistry, etc.",
             "Onboarding a new discipline = pointing the system at its SOPs. No engineering project.",
         ]),
        (4.7, "Cold-start discovery",
         [
             "System discovers business capabilities, themes, epics, personas, tests from this discipline's SOPs alone — no curated catalogs up front.",
             "Top-down cascade: capabilities first, then themes within each capability, then epics within each theme. Reviewable at every level.",
         ]),
        (8.8, "Non-blocking review",
         [
             "4 agent classes: Catalog Builder, Tagger, Story Extractor, Sampled Validator.",
             "Confidence at every boundary; every emission flows to corpus and is logged for human review post-hoc — pipeline never waits.",
         ]),
    ]
    for x, header, items in cols:
        add_text(slide, x, bullets_y, col_w, 0.32, header,
                 font_size=12, bold=True, color=ACCENT)
        add_bullets(slide, x, bullets_y + 0.36, col_w, 2.2, items,
                    font_size=11)


# ============================================================================
# SLIDE 2 — Inputs
# ============================================================================

def slide_two(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Inputs — just SOPs")
    add_subtitle(
        slide,
        "Everything else — capabilities, themes, epics, personas, tests — "
        "is discovered cold from this discipline's SOPs.",
    )

    # LEFT: inputs
    lx, lw = 0.5, 6.0

    add_text(slide, lx, 1.55, lw, 0.36,
             "From the discipline (per onboarding):",
             font_size=13, bold=True, color=ACCENT)

    add_box(slide, lx, 1.95, lw, 1.05, "SOPs — the corpus",
            fill=BOX_FILL, line=BOX_LINE, font_size=15, bold=True,
            align=PP_ALIGN.LEFT)
    add_text(slide, lx + 0.20, 2.55, lw - 0.40, 0.40,
             "Typically 30–150 SOPs.  No tagging, no structuring required.",
             font_size=11, italic=True, color=TEXT_GREY)

    add_text(slide, lx, 3.30, lw, 0.36,
             "Universal defaults (not discipline-specific):",
             font_size=13, bold=True, color=ACCENT)

    add_box(slide, lx, 3.70, lw, 1.05,
            "Validator thresholds + 4-shape rubric",
            fill=END_FILL, line=END_LINE, font_size=14, bold=True,
            align=PP_ALIGN.LEFT)
    add_text(slide, lx + 0.20, 4.30, lw - 0.40, 0.40,
             "Tuned from telemetry on first run. No per-discipline calibration.",
             font_size=11, italic=True, color=TEXT_GREY)

    # bottom: small callout on what changed vs older designs
    add_text(slide, lx, 5.20, lw, 0.36,
             "Notably absent (vs prior designs):",
             font_size=11.5, bold=True, color=PARK_LINE)
    add_bullets(slide, lx, 5.55, lw, 1.30, [
        "No prior-discipline catalogs (no warm-start).",
        "No universal lab_stage_v1 artifact — workflow axis is the discovered "
        "business_capability_v1 catalog, per discipline.",
        "No SME pre-curation of personas / tests — discovered from SOP mentions.",
    ], font_size=10.5, color=TEXT_GREY)

    # RIGHT: discovered catalogs
    rx, rw = 6.85, 5.95
    add_text(slide, rx, 1.55, rw, 0.36,
             "Discovered from the corpus (per discipline):",
             font_size=13, bold=True, color=END_LINE)

    discovered = [
        ("business_capability_v1",
         "6–12 broad business outcomes the lab produces or maintains."),
        ("theme_v1",
         "12–20 subject areas, anchored under capabilities."),
        ("epic_v1",
         "80–150 fine-grained, anchored under themes."),
        ("persona_v1",
         "Actors named in SOPs (humans, systems, external systems)."),
        ("test_v1",
         "Tests / assays named in SOPs."),
    ]
    ty = 2.00
    for name, desc in discovered:
        add_box(slide, rx, ty, rw, 0.92, "",
                fill=CARD_FILL, line=CARD_LINE)
        add_text(slide, rx + 0.18, ty + 0.10, rw - 0.36, 0.30,
                 name, font_size=12.5, bold=True, color=ACCENT)
        add_text(slide, rx + 0.18, ty + 0.46, rw - 0.36, 0.40,
                 desc, font_size=10.5, color=TEXT_DARK)
        ty += 1.02


# ============================================================================
# SLIDE 3 — End-to-end flow
# ============================================================================

def slide_three(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "End-to-end flow")
    add_subtitle(
        slide,
        "Three phases. Phase 1 once (taxonomy cascade + entity extraction), Phase 2 in parallel across SOPs, Phase 3 deferred to the roadmap.",
    )

    # Three large phase blocks across the slide
    py = 1.7
    ph = 4.5
    pw = 4.0
    gap = 0.18
    px1 = 0.5
    px2 = px1 + pw + gap
    px3 = px2 + pw + gap

    # Phase 1
    add_box(slide, px1, py, pw, ph, "",
            fill=PHASE_FILL, line=PHASE_LINE)
    add_text(slide, px1 + 0.15, py + 0.10, pw - 0.3, 0.32,
             "Phase 1 — Catalog Build",
             font_size=14, bold=True, color=PHASE_LINE)
    add_text(slide, px1 + 0.15, py + 0.40, pw - 0.3, 0.28,
             "Once, on full corpus. Cold discovery.",
             font_size=10, italic=True, color=TEXT_GREY)

    inner_x = px1 + 0.25
    inner_w = pw - 0.5
    iy = py + 0.82
    add_node(slide, inner_x, iy, inner_w, 0.38,
             "SOP Ingest & Chunking", kind="process", font_size=10.5)
    iy += 0.48
    add_node(slide, inner_x, iy, inner_w, 0.38,
             "Catalog Builder × 5 runs", kind="agent", font_size=11, bold=True)

    iy += 0.48
    sub_h = 0.32
    sub_gap = 0.04
    for i, label in enumerate([
        "Capability Discovery (top)",
        "Theme Discovery (under capabilities)",
        "Epic Discovery (under themes)",
        "Persona Discovery",
        "Test Discovery",
    ]):
        add_box(slide, inner_x + 0.15, iy + i * (sub_h + sub_gap),
                inner_w - 0.30, sub_h, label,
                fill=CARD_FILL, line=CARD_LINE, font_size=9.5,
                align=PP_ALIGN.LEFT)

    add_text(slide, px1 + 0.15, py + ph - 0.62, pw - 0.3, 0.28,
             "Output: 5 catalogs (capability/theme/epic/persona/test)",
             font_size=10, bold=True, italic=True, color=END_LINE,
             align=PP_ALIGN.LEFT)
    add_text(slide, px1 + 0.15, py + ph - 0.32, pw - 0.3, 0.26,
             "↳ low-confidence admits flagged for priority review",
             font_size=9, italic=True, color=PARK_LINE,
             align=PP_ALIGN.LEFT)

    # Phase 2
    add_box(slide, px2, py, pw, ph, "",
            fill=PHASE_FILL, line=PHASE_LINE)
    add_text(slide, px2 + 0.15, py + 0.10, pw - 0.3, 0.32,
             "Phase 2 — Per-SOP Extraction",
             font_size=14, bold=True, color=PHASE_LINE)
    add_text(slide, px2 + 0.15, py + 0.40, pw - 0.3, 0.28,
             "Parallel across SOPs. Tagging + extraction.",
             font_size=10, italic=True, color=TEXT_GREY)

    inner_x2 = px2 + 0.25
    inner_w2 = pw - 0.5
    jy = py + 0.85
    add_node(slide, inner_x2, jy, inner_w2, 0.40,
             "For each SOP (in parallel):", kind="process",
             font_size=10.5, bold=True)
    jy += 0.55
    add_node(slide, inner_x2, jy, inner_w2, 0.40,
             "For each chunk (from Phase 1):",
             kind="process", font_size=10)
    jy += 0.50
    add_node(slide, inner_x2, jy, inner_w2, 0.55,
             "Tagger\n(persona + test + capability)", kind="agent", font_size=10)
    jy += 0.72
    add_node(slide, inner_x2, jy, inner_w2, 0.55,
             "Story Extractor\n(self-checks against rubric)",
             kind="agent", font_size=10)

    add_text(slide, px2 + 0.15, py + ph - 0.62, pw - 0.3, 0.28,
             "Output: story corpus (4 shapes) + audit",
             font_size=10, bold=True, italic=True, color=END_LINE,
             align=PP_ALIGN.LEFT)
    add_text(slide, px2 + 0.15, py + ph - 0.32, pw - 0.3, 0.26,
             "↳ self_check FAIL or low-conf flagged for priority review",
             font_size=9, italic=True, color=PARK_LINE,
             align=PP_ALIGN.LEFT)

    # Phase 3 (future)
    add_box(slide, px3, py, pw, ph, "",
            fill=FUTURE_FILL, line=FUTURE_LINE)
    add_text(slide, px3 + 0.15, py + 0.10, pw - 0.3, 0.32,
             "Phase 3 — Cross-SOP Synthesis",
             font_size=14, bold=True, color=FUTURE_LINE)
    add_text(slide, px3 + 0.15, py + 0.40, pw - 0.3, 0.28,
             "FUTURE ROADMAP — milestone 2.",
             font_size=10, italic=True, bold=True, color=PARK_LINE)

    inner_x3 = px3 + 0.25
    inner_w3 = pw - 0.5
    ky = py + 0.85
    add_box(slide, inner_x3, ky, inner_w3, 0.40,
            "Cluster stories across SOPs",
            fill=CARD_FILL, line=CARD_LINE, font_size=10.5)
    ky += 0.55
    add_box(slide, inner_x3, ky, inner_w3, 0.40,
            "Lift to capability stories\n(threshold ≥ 2 distinct SOPs)",
            fill=CARD_FILL, line=CARD_LINE, font_size=10)
    ky += 0.62
    add_box(slide, inner_x3, ky, inner_w3, 0.40,
            "Emit Stream B (capability)",
            fill=CARD_FILL, line=CARD_LINE, font_size=10.5)

    add_text(slide, px3 + 0.15, py + ph - 0.45, pw - 0.3, 0.32,
             "Milestone 1 ships without Phase 3.",
             font_size=10, bold=True, italic=True, color=PARK_LINE,
             align=PP_ALIGN.LEFT)

    # Cross-cutting: two labeled artifacts (Validator + HITL Queue)
    bb_y = py + ph + 0.20
    bb_h = 0.62
    cc_w = (12.33 - 0.20) / 2
    # Sampled Validator (left)
    add_box(slide, 0.5, bb_y, cc_w, bb_h, "",
            fill=SIDE_FILL, line=SIDE_LINE)
    add_text(slide, 0.65, bb_y + 0.05, cc_w - 0.30, 0.26,
             "Sampled Validator QC (Quality Control)", font_size=11.5, bold=True,
             color=SIDE_LINE)
    add_text(slide, 0.65, bb_y + 0.30, cc_w - 0.30, 0.30,
             "~10% async sampling of emitted Phase 2 stories.  Drift → priority review.",
             font_size=10, italic=True, color=TEXT_DARK)
    # HITL Queue (right)
    hitl_x = 0.5 + cc_w + 0.20
    add_box(slide, hitl_x, bb_y, cc_w, bb_h, "",
            fill=PARK_FILL, line=PARK_LINE)
    add_text(slide, hitl_x + 0.15, bb_y + 0.05, cc_w - 0.30, 0.26,
             "Review queue (HITL — Human-in-the-Loop)", font_size=11.5, bold=True,
             color=PARK_LINE)
    add_text(slide, hitl_x + 0.15, bb_y + 0.30, cc_w - 0.30, 0.30,
             "Every emission logged.  Failures + low-conf prioritized for human review (post-hoc, non-blocking).",
             font_size=10, italic=True, color=TEXT_DARK)

    # Inter-phase arrows
    arrow_y = py + ph / 2
    add_arrow(slide, px1 + pw + 0.01, arrow_y, px2 - 0.01, arrow_y)
    add_arrow(slide, px2 + pw + 0.01, arrow_y, px3 - 0.01, arrow_y,
              color=FUTURE_LINE)


# ============================================================================
# SLIDE 4 — The 4 agents
# ============================================================================

def slide_four(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "The 4 agents")
    add_subtitle(
        slide,
        "4 LLM agents, each owning one step. Chunking, embedding, and clustering "
        "are deterministic; the LLM handles naming, tagging, extraction, and validation.",
    )

    # 2x2 grid of agent cards
    cw, ch = 6.0, 2.65
    gap_x = 0.33
    gap_y = 0.30
    cx1 = 0.5
    cx2 = cx1 + cw + gap_x
    cy1 = 1.55
    cy2 = cy1 + ch + gap_y

    def agent_card(x, y, name, role, runs_in, what_it_does, why_simpler):
        add_box(slide, x, y, cw, ch, "",
                fill=CARD_FILL, line=CARD_LINE)
        # name pill
        add_box(slide, x + 0.18, y + 0.18, 2.55, 0.42, name,
                fill=SYS_FILL, line=SYS_LINE,
                font_size=12.5, bold=True, align=PP_ALIGN.CENTER)
        # role
        add_text(slide, x + 2.85, y + 0.22, cw - 3.05, 0.36, role,
                 font_size=11.5, bold=True, color=ACCENT)
        add_text(slide, x + 2.85, y + 0.55, cw - 3.05, 0.30,
                 f"Runs in: {runs_in}",
                 font_size=10, italic=True, color=TEXT_GREY)
        # body
        add_text(slide, x + 0.20, y + 1.00, cw - 0.40, 0.30,
                 "What it does:", font_size=11, bold=True, color=TEXT_DARK)
        add_text(slide, x + 0.20, y + 1.27, cw - 0.40, 0.85, what_it_does,
                 font_size=10.5, color=TEXT_DARK)
        add_text(slide, x + 0.20, y + 2.10, cw - 0.40, 0.30,
                 "Why this is simpler:", font_size=11, bold=True, color=END_LINE)
        add_text(slide, x + 0.20, y + 2.32, cw - 0.40, 0.30, why_simpler,
                 font_size=10, italic=True, color=TEXT_GREY)

    agent_card(
        cx1, cy1,
        "Catalog Builder",
        "Cold-discovers the 5 catalogs from the SOP corpus",
        "Phase 1 (5 runs, two modes)",
        "Taxonomy mode (capability → theme → epic, top-down cascade): names "
        "each cluster and admits via rubric; child levels cluster only within "
        "their parent's chunks.\nEntity mode (persona, test): extracts named "
        "mentions per chunk, dedupes across the corpus, admits each unique entity.",
        "One agent class with two modes, instead of five separate agent classes.",
    )
    agent_card(
        cx2, cy1,
        "Tagger",
        "Tags each chunk with persona / test / capability",
        "Phase 2 (per chunk, parallel across SOPs)",
        "Single structured-output call: looks up persona refs, test ref, and "
        "capability ref against persona_v1, test_v1, and business_capability_v1. "
        "theme_ref and epic_ref are inherited directly from the chunk's Phase 1 "
        "cascade buckets — no LLM lookup needed.",
        "Three LLM lookups bundled into one call; theme/epic come for free from cascade.",
    )
    agent_card(
        cx1, cy2,
        "Story Extractor",
        "Produces 4-shape stories with inline self-check",
        "Phase 2 (per tagged chunk)",
        "Generates story (capability / stage-split / config-instance / cleanup) "
        "matching the chunk + tags. Emits extraction_confidence + self_check verdict. "
        "Story always emits to corpus; FAIL or low-conf is flagged for priority review.",
        "Non-blocking: self-check on the hot path, no revise loop, no waiting for a human.",
    )
    agent_card(
        cx2, cy2,
        "Sampled Validator",
        "Async drift detection on a sample of stories",
        "Cross-cutting (~10% of emitted stories)",
        "Re-runs the full Validator rubric on a sampled fraction of stories already "
        "in the corpus. On FAIL, the story is logged to the review queue (with "
        "rubric diff) for priority review — story stays in corpus.",
        "Drift detection on a sample, not a gate — never blocks emission.",
    )


# ============================================================================
# SLIDE 5 — Phase 1: Catalog Build
# ============================================================================

def slide_five(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Phase 1 — Catalog Build (cold)")
    add_subtitle(
        slide,
        "The 3 taxonomy runs cascade top-down (capability → theme within capability → "
        "epic within theme).  The 2 entity runs (persona, test) are parallel.  All five "
        "share the same chunking; each step admits its own.",
    )

    # Left side: shared input → split (taxonomy cascade vs entity extraction) → catalog
    fx = 0.5
    full_w = 3.6
    col_w = 1.70
    col_gap = 0.20
    col_a_x = fx
    col_b_x = fx + col_w + col_gap
    nh = 0.42

    add_section_header(slide, fx, 1.45, full_w, "Phase 1 pipeline")

    # Shared top: SOPs → Chunk
    y_sops = 1.78
    add_node(slide, fx, y_sops, full_w, nh,
             "SOPs (full corpus)", kind="start", font_size=11)
    y_chunk = y_sops + nh + 0.18
    add_node(slide, fx, y_chunk, full_w, nh,
             "1. Chunk SOPs (deterministic)", kind="process", font_size=10.5)

    # ---- Cascade stages (col A) ----
    # Column meta-labels intentionally omitted; subtitle conveys taxonomy/entity split.
    # Stages start with a small gap below the chunk box to let arrows breathe.
    stage_h = 0.85
    y_stage1 = y_chunk + nh + 0.50

    def cascade_box(y, head, body):
        add_box(slide, col_a_x, y, col_w, stage_h, "",
                fill=SYS_FILL, line=SYS_LINE)
        add_text(slide, col_a_x + 0.08, y + 0.05, col_w - 0.16, 0.24,
                 head, font_size=9.5, bold=True, color=TEXT_DARK,
                 align=PP_ALIGN.LEFT)
        add_text(slide, col_a_x + 0.08, y + 0.30, col_w - 0.16, stage_h - 0.34,
                 body, font_size=8, color=TEXT_DARK,
                 align=PP_ALIGN.LEFT)

    cascade_box(y_stage1,
                "Capability run",
                "Cluster all chunks → name +\nadmit.  ~6–12 capabilities.")
    y_stage2 = y_stage1 + stage_h + 0.16
    cascade_box(y_stage2,
                "Theme run (per capability)",
                "Re-cluster within each\ncapability.  ~12–20 themes.")
    y_stage3 = y_stage2 + stage_h + 0.16
    cascade_box(y_stage3,
                "Epic run (per theme)",
                "Re-cluster within each\ntheme.  ~80–150 epics.")

    # ---- Extraction stages (col B) ----
    def extract_box(y, head, body, kind="process"):
        fill = SYS_FILL if kind == "agent" else BOX_FILL
        line = SYS_LINE if kind == "agent" else BOX_LINE
        add_box(slide, col_b_x, y, col_w, stage_h, "",
                fill=fill, line=line)
        add_text(slide, col_b_x + 0.08, y + 0.05, col_w - 0.16, 0.24,
                 head, font_size=9.5, bold=True, color=TEXT_DARK,
                 align=PP_ALIGN.LEFT)
        add_text(slide, col_b_x + 0.08, y + 0.30, col_w - 0.16, stage_h - 0.34,
                 body, font_size=8, color=TEXT_DARK,
                 align=PP_ALIGN.LEFT)

    extract_box(y_stage1,
                "Extract mentions (LLM)",
                "Per chunk, pull named\npersonas / tests.",
                kind="agent")
    extract_box(y_stage2,
                "Dedupe + admit",
                "Merge near-duplicates;\nadmit each unique entity.",
                kind="process")

    # ---- Output: Catalog YAML (full width) ----
    y_cat = y_stage3 + stage_h + 0.20
    add_node(slide, fx, y_cat, full_w, nh, "Catalog YAML  (5 catalogs)",
             kind="end", font_size=11, bold=True)

    # ---- Arrows ----
    cx_full = fx + full_w / 2
    cx_a = col_a_x + col_w / 2
    cx_b = col_b_x + col_w / 2

    # Shared top
    add_arrow(slide, cx_full, y_sops + nh, cx_full, y_chunk)
    # Chunk → top of each column
    add_arrow(slide, cx_full, y_chunk + nh, cx_a, y_stage1)
    add_arrow(slide, cx_full, y_chunk + nh, cx_b, y_stage1)
    # Cascade arrows (col A)
    add_arrow(slide, cx_a, y_stage1 + stage_h, cx_a, y_stage2)
    add_arrow(slide, cx_a, y_stage2 + stage_h, cx_a, y_stage3)
    # Extraction arrows (col B)
    add_arrow(slide, cx_b, y_stage1 + stage_h, cx_b, y_stage2)
    # Both columns merge into Catalog YAML
    add_arrow(slide, cx_a, y_stage3 + stage_h, cx_full, y_cat)
    add_arrow(slide, cx_b, y_stage2 + stage_h, cx_full, y_cat)

    # ---- Review queue: beside Catalog YAML (single arrow, no column-crossing) ----
    rev_x = fx + full_w + 0.30
    rev_w = 2.45
    rev_y = y_cat
    rev_h = nh
    add_box(slide, rev_x, rev_y, rev_w, rev_h, "",
            fill=PARK_FILL, line=PARK_LINE)
    add_text(slide, rev_x + 0.10, rev_y + 0.04, rev_w - 0.20, 0.20,
             "↘ Review queue", font_size=10, bold=True, color=PARK_LINE)
    add_text(slide, rev_x + 0.10, rev_y + 0.22, rev_w - 0.20, rev_h - 0.24,
             "Low-conf admits (any level) flagged — pipeline never waits.",
             font_size=8.5, italic=True, color=TEXT_DARK)
    # Single arrow from Catalog YAML right edge to Review queue left edge
    add_arrow(slide, fx + full_w + 0.01, y_cat + nh / 2,
              rev_x - 0.01, rev_y + rev_h / 2,
              color=PARK_LINE, weight=1.5)

    # Right side: what comes out — 5 catalogs (capability spans top row)
    rx = 7.10
    rw = 5.75
    add_section_header(slide, rx, 1.45, rw,
                       "Phase 1 outputs — 5 catalogs (3 cascade + 2 entity)")

    tile_w = (rw - 0.20) / 2
    tile_h = 1.05
    th_x1 = rx
    th_x2 = rx + tile_w + 0.20
    tr_y0 = 1.78  # capability (full width)
    tr_y1 = tr_y0 + tile_h + 0.18
    tr_y2 = tr_y1 + tile_h + 0.18

    def catalog_tile(x, y, w, name, n_text, schema_brief):
        add_box(slide, x, y, w, tile_h, "",
                fill=CARD_FILL, line=CARD_LINE)
        add_text(slide, x + 0.12, y + 0.06, w - 0.24, 0.26,
                 name, font_size=12, bold=True, color=ACCENT)
        add_text(slide, x + 0.12, y + 0.32, w - 0.24, 0.22,
                 n_text, font_size=9.5, italic=True, color=TEXT_GREY)
        add_text(slide, x + 0.12, y + 0.55, w - 0.24, 0.45,
                 schema_brief, font_size=9.5, color=TEXT_DARK)

    # capability tile spans full row width
    catalog_tile(rx, tr_y0, rw,
                 "business_capability_v1",
                 "~6–12 broad business outcomes (top of hierarchy)",
                 "key · label · desc · n_themes · admit_confidence")
    # theme + epic
    catalog_tile(th_x1, tr_y1, tile_w,
                 "theme_v1",
                 "~12–20, under capabilities",
                 "key · label · capability_ref · n_chunks · n_sops · admit_confidence")
    catalog_tile(th_x2, tr_y1, tile_w,
                 "epic_v1",
                 "~80–150, under themes",
                 "key · label · theme_ref · n_sops · admit_confidence")
    # persona + test
    catalog_tile(th_x1, tr_y2, tile_w,
                 "persona_v1",
                 "Named actors in SOPs",
                 "key · label · actor_type · n_mentions · admit_confidence")
    catalog_tile(th_x2, tr_y2, tile_w,
                 "test_v1",
                 "Named tests / assays",
                 "key · label · n_sops · admit_confidence")

    # Bottom note on unclassified bucket (taxonomy runs only — extraction has no clusters)
    add_box(slide, rx, tr_y2 + tile_h + 0.20, rw, 0.78,
            "Unclassified bucket on each taxonomy catalog (capability / theme / "
            "epic) holds chunks that did not cluster cleanly. Surfaced for "
            "review, not lost.",
            fill=PARK_FILL, line=PARK_LINE, font_size=10.5,
            align=PP_ALIGN.LEFT)


# ============================================================================
# SLIDE 6 — Catalog schemas + examples
# ============================================================================

def slide_six(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Catalog schemas + examples")
    add_subtitle(
        slide,
        "Each Phase 1 output is a small structured YAML. Slug-based keys + "
        "human-readable labels. Microbiology illustrative.",
    )

    yaml_capability = """catalog: business_capability_v1
version: 1
discipline: microbiology
business_capabilities:
  - {key: specimen_lifecycle, label: "Specimen Lifecycle Management", desc: "Intake, routing, QC (Quality Control), disposal", n_themes: 4, admit_confidence: 0.91}
  - {key: diagnostic_execution, label: "Diagnostic Test Execution", n_themes: 5, admit_confidence: 0.88}
  - {key: result_delivery, label: "Result Reporting & Delivery", n_themes: 3, admit_confidence: 0.93}
  # ... 6-12 capabilities total"""

    yaml_themes = """catalog: theme_v1
version: 1
discipline: microbiology
themes:
  - key: specimen_handling
    label: "Specimen handling"
    capability_ref: specimen_lifecycle
    desc: Intake, routing, QC (Quality Control)
    n_chunks: 142
    n_sops: 22
    admit_confidence: 0.92
  - {key: culture_setup, label: "Culture setup", capability_ref: diagnostic_execution, n_sops: 18, admit_confidence: 0.88}
  - {key: result_review, label: "Result review", capability_ref: result_delivery, n_sops: 24, admit_confidence: 0.90}
  - {key: unclassified, label: "Unclassified", n_chunks: 14}"""

    yaml_epics = """catalog: epic_v1
version: 1
discipline: microbiology
epics:
  - key: specimen_intake
    label: "Specimen intake at receiving"
    theme_ref: specimen_handling
    desc: Accessioning, barcoding
    n_sops: 14
    admit_confidence: 0.88
  - {key: gram_stain_prep, label: "Gram stain prep", theme_ref: culture_setup, n_sops: 14, admit_confidence: 0.91}
  - {key: result_sign_off, label: "Result sign-off", theme_ref: result_review, n_sops: 24, admit_confidence: 0.85}
  # ~80-150 epics total"""

    yaml_personas = """catalog: persona_v1     # named actors discovered in SOPs
version: 1
discipline: microbiology
personas:
  - key: microbiologist
    label: "Microbiologist"
    actor_type: human
    desc: Reads cultures, signs off
    n_mentions: 167
    admit_confidence: 0.95
  - {key: lab_assistant, label: "Lab assistant", actor_type: human, n_mentions: 91, admit_confidence: 0.93}
  - {key: lims, label: "LIMS (Laboratory Information Management System)", actor_type: system, n_mentions: 134, admit_confidence: 0.97}
  - {key: emr, label: "EMR (Electronic Medical Record)", actor_type: external_system, n_mentions: 28, admit_confidence: 0.86}"""

    yaml_tests = """catalog: test_v1
version: 1
discipline: microbiology
tests:
  - key: blood_culture
    label: "Blood culture"
    desc: Aerobic + anaerobic, 5d
    n_sops: 11
    admit_confidence: 0.91
  - {key: urine_culture, label: "Urine culture", n_sops: 9, admit_confidence: 0.89}
  - {key: gram_stain, label: "Gram stain", n_sops: 14, admit_confidence: 0.94}
  - {key: ast, label: "AST (Antimicrobial Susceptibility Testing)", n_sops: 7, admit_confidence: 0.82}"""

    # Layout: capability spans top, then 2x2 below for theme/epic/persona/test
    full_w = 12.33
    half_w = (full_w - 0.20) / 2  # 6.065
    cap_h = 1.10
    other_h = 1.95
    gy = 0.15

    by_cap = 1.80   # title rendered at 1.50, just below subtitle ending at 1.47
    by_row1 = by_cap + cap_h + gy
    by_row2 = by_row1 + other_h + gy

    add_code_box(slide, 0.5, by_cap, full_w, cap_h, yaml_capability,
                 font_size=8, title="business_capability_v1.yaml  (top of hierarchy)")
    add_code_box(slide, 0.5, by_row1, half_w, other_h, yaml_themes,
                 font_size=8, title="theme_v1.yaml")
    add_code_box(slide, 0.5 + half_w + 0.20, by_row1, half_w, other_h, yaml_epics,
                 font_size=8, title="epic_v1.yaml")
    add_code_box(slide, 0.5, by_row2, half_w, other_h, yaml_personas,
                 font_size=8, title="persona_v1.yaml")
    add_code_box(slide, 0.5 + half_w + 0.20, by_row2, half_w, other_h, yaml_tests,
                 font_size=8, title="test_v1.yaml")


# ============================================================================
# SLIDE 7 — Phase 2: Per-SOP Extraction
# ============================================================================

def slide_seven(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Phase 2 — Per-SOP Extraction")
    add_subtitle(
        slide,
        "For each SOP in parallel, walk its chunks (from Phase 1) and run Tagger then Story Extractor on each.",
    )

    # LEFT: flow inside a ForEach SOP container
    cont_x, cont_y = 0.5, 1.55
    cont_w, cont_h = 6.30, 5.55
    cont = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(cont_x), Inches(cont_y),
                                  Inches(cont_w), Inches(cont_h))
    cont.fill.solid()
    cont.fill.fore_color.rgb = RGBColor(0xFA, 0xFB, 0xFD)
    cont.line.color.rgb = RGBColor(0xB8, 0xC4, 0xD2)
    cont.line.width = Pt(1.0)
    # dashed effect via XML (best-effort)
    line_elem = cont.line._get_or_add_ln()
    prstDash = etree.SubElement(line_elem, qn("a:prstDash"))
    prstDash.set("val", "dash")

    add_text(slide, cont_x + 0.20, cont_y + 0.10, cont_w - 0.4, 0.30,
             "ForEach SOP (in parallel)",
             font_size=12, bold=True, color=TEXT_GREY,
             align=PP_ALIGN.LEFT)

    fx = cont_x + 0.45
    nw = cont_w - 0.90
    nh = 0.45
    sy = cont_y + 0.55

    add_node(slide, fx, sy, nw, nh, "One SOP",
             kind="start", font_size=11)
    sy2 = sy + nh + 0.20
    add_node(slide, fx, sy2, nw, nh,
             "Load this SOP's chunks (from Phase 1)",
             kind="process", font_size=10.5)

    # ForEach chunk inner container
    inner_x = fx + 0.10
    inner_y = sy2 + nh + 0.30
    inner_w = nw - 0.20
    inner_h = 1.65
    inner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(inner_x), Inches(inner_y),
                                   Inches(inner_w), Inches(inner_h))
    inner.fill.solid()
    inner.fill.fore_color.rgb = RGBColor(0xFE, 0xFB, 0xF2)
    inner.line.color.rgb = SYS_LINE
    inner.line.width = Pt(0.75)
    line_elem2 = inner.line._get_or_add_ln()
    prstDash2 = etree.SubElement(line_elem2, qn("a:prstDash"))
    prstDash2.set("val", "dash")

    add_text(slide, inner_x + 0.15, inner_y + 0.08, inner_w - 0.3, 0.26,
             "ForEach chunk", font_size=11, bold=True, italic=True,
             color=SYS_LINE, align=PP_ALIGN.LEFT)

    # Tagger
    tag_h = 0.40
    tag_y = inner_y + 0.40
    tag_w = inner_w - 0.60
    tag_x = inner_x + 0.30
    add_node(slide, tag_x, tag_y, tag_w, tag_h,
             "Tagger  (persona · test · capability)",
             kind="agent", font_size=10.5)
    # Small inheritance caption under Tagger
    add_text(slide, tag_x, tag_y + tag_h + 0.01, tag_w, 0.18,
             "+ theme_ref / epic_ref inherited from Phase 1 cascade buckets (no LLM call)",
             font_size=8, italic=True, color=TEXT_GREY,
             align=PP_ALIGN.CENTER)

    # Story Extractor
    se_h = 0.50
    se_y = tag_y + tag_h + 0.20
    add_node(slide, tag_x, se_y, tag_w, se_h,
             "Story Extractor\n(generates 4-shape story · self_check rubric)",
             kind="agent", font_size=10)

    # arrows inside inner
    inner_cx = inner_x + inner_w / 2
    add_arrow(slide, inner_cx, tag_y + tag_h, inner_cx, se_y)

    # arrows outside inner (above)
    add_arrow(slide, fx + nw / 2, sy + nh,  fx + nw / 2, sy2)
    add_arrow(slide, fx + nw / 2, sy2 + nh, fx + nw / 2, inner_y)

    # Single destination below inner: stories → corpus + audit
    dest_y = inner_y + inner_h + 0.20
    dest_h = 0.50
    add_node(slide, fx, dest_y, nw, dest_h,
             "Stories → corpus + audit log",
             kind="end", font_size=11, bold=True)
    add_arrow(slide, fx + nw / 2, inner_y + inner_h,
              fx + nw / 2, dest_y)

    # Side note for review-queue routing
    note_y = dest_y + dest_h + 0.12
    add_text(slide, fx, note_y, nw, 0.30,
             "↘ FAIL or low-confidence stories also flagged for priority review",
             font_size=9.5, italic=True, color=PARK_LINE,
             align=PP_ALIGN.LEFT)

    # RIGHT: example trace (code-box titles label each step; no separate header)
    rx = 7.05
    rw = 5.80

    chunk_yaml = """# Input chunk (already produced by Phase 1 chunking)
chunk_id: SOP_023.sec_4_2.chunk_07
sop_ref: SOP_023
text: |
  After accessioning, the lab assistant
  prepares a Gram stain on the specimen.
  The microbiologist reviews the smear."""

    tagger_yaml = """# Tagger output  (3 LLM lookups + 2 inherited)
tags:
  # Looked up via LLM:
  persona_refs:    [lab_assistant, microbiologist]
  test_ref:        gram_stain
  capability_ref:  diagnostic_execution
  # Inherited from chunk's Phase 1 cascade buckets:
  theme_ref:       culture_setup       # from capability cascade
  epic_ref:        gram_stain_prep     # from theme cascade
  confidence: 0.94"""

    story_yaml = """# Story Extractor output (one of possibly several stories per chunk)
story_id: ST_023_07_a
shape: workflow_stage_split
title: Lab assistant preps Gram stain before reading
sop_refs: [SOP_023]
capability_ref: diagnostic_execution
theme_ref: culture_setup     epic_ref: gram_stain_prep
persona_refs: [lab_assistant]            # narrowed from tagger's 2 personas
test_ref: gram_stain                     # this story is about the prep action
predecessor_capability: specimen_lifecycle
successor_capability:  diagnostic_execution
acceptance_criteria: ["Smear stained per protocol"]
extraction_confidence: 0.91
self_check: pass
hitl_status: auto"""

    ch1, ch2, ch3 = 1.05, 1.45, 2.10
    y_a = 1.80
    y_b = y_a + ch1 + 0.18
    y_c = y_b + ch2 + 0.18
    add_code_box(slide, rx, y_a, rw, ch1, chunk_yaml,
                 font_size=8, title="Chunk in")
    add_code_box(slide, rx, y_b, rw, ch2, tagger_yaml,
                 font_size=8, title="Tagger")
    add_code_box(slide, rx, y_c, rw, ch3, story_yaml,
                 font_size=8, title="Story Extractor")


# ============================================================================
# SLIDE 8 — Story schema + 4 shapes
# ============================================================================

def slide_eight(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Story schema + 4 shapes")
    add_subtitle(
        slide,
        "Shared fields, then shape-specific fields. Validator rubric is shape-aware.",
    )

    # Top: shared schema strip
    schema_yaml = """story_id        | shape ∈ {capability, workflow_stage_split, configuration_instance, cleanup}
title           | sop_refs[]                | capability_ref  | theme_ref | epic_ref
persona_refs[]  | test_ref                  | acceptance_criteria[]      | self_check
extraction_confidence (0.0-1.0)  |  hitl_status ∈ {auto, confirmed, corrected, rejected}
shape-specific  | (capability)             sop_refs ≥ 2 distinct SOPs  (recurring pattern)
                | (workflow_stage_split)   predecessor_capability, successor_capability
                | (configuration_instance) config_keys[]
                | (cleanup)                deprecates[], replaced_by

note: shape=capability is a STORY TYPE (lifted recurring pattern, sop_breadth ≥ 2).
      It references — but is distinct from — the business_capability_v1 catalog."""
    add_code_box(slide, 0.5, 1.80, 12.33, 1.50, schema_yaml,
                 font_size=8.5, title="Story schema (all shapes share the top block)")

    # 2x2 grid of example stories
    sx1, sx2 = 0.5, 6.70
    sy1 = 3.55
    sy2 = sy1 + 1.80
    sw, sh = 6.13, 1.70

    cap_yaml = """story_id: ST_C_017
shape: capability
title: System tracks specimens by barcode end-to-end
sop_refs: [SOP_023, SOP_041, SOP_055]   # ≥ 2 SOPs
capability_ref: specimen_lifecycle
theme_ref: specimen_handling   epic_ref: specimen_intake
persona_refs: [lims]
acceptance_criteria: ["Unique barcode at intake", "Resolves to location"]
extraction_confidence: 0.92
self_check: pass
hitl_status: auto"""

    split_yaml = """story_id: ST_S_044    shape: workflow_stage_split
title: Microbiologist signs off culture before release
sop_refs: [SOP_017]    persona_refs: [microbiologist]
capability_ref: result_delivery
theme_ref: result_review    epic_ref: result_sign_off
predecessor_capability: diagnostic_execution
successor_capability:  result_delivery
acceptance_criteria: ["Sign-off requires organism ID confirmation"]
extraction_confidence: 0.88
self_check: pass    hitl_status: auto"""

    cfg_yaml = """story_id: ST_F_009    shape: configuration_instance
title: Configurable reject criteria for urine specimens
sop_refs: [SOP_039]    persona_refs: [microbiologist]    test_ref: urine_culture
capability_ref: specimen_lifecycle
theme_ref: specimen_handling    epic_ref: specimen_intake
config_keys: [colony_count_threshold, contamination_rule]
acceptance_criteria: ["Thresholds editable per site without code change"]
extraction_confidence: 0.79      # below 0.85
self_check: pass
hitl_status: auto                # in corpus, flagged for priority review"""

    clean_yaml = """story_id: ST_X_103    shape: cleanup
title: Remove deprecated Gram stain manual reader path
sop_refs: [SOP_044, SOP_044_amended]
capability_ref: diagnostic_execution
theme_ref: culture_setup    epic_ref: gram_stain_prep
test_ref: gram_stain
deprecates: [legacy_manual_reader_v0]
replaced_by: [automated_reader_v1]
acceptance_criteria: ["All v0 references removed"]
extraction_confidence: 0.91    self_check: pass    hitl_status: auto"""

    add_code_box(slide, sx1, sy1, sw, sh, cap_yaml,
                 font_size=8, title="capability")
    add_code_box(slide, sx2, sy1, sw, sh, split_yaml,
                 font_size=8, title="workflow_stage_split")
    add_code_box(slide, sx1, sy2, sw, sh, cfg_yaml,
                 font_size=8, title="configuration_instance")
    add_code_box(slide, sx2, sy2, sw, sh, clean_yaml,
                 font_size=8, title="cleanup")


# ============================================================================
# SLIDE 9 — Sampled Validator QC
# ============================================================================

def slide_nine(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Sampled Validator QC")
    add_subtitle(
        slide,
        "Drift detection, not a gate. ~10% of stories audited asynchronously; "
        "Story Extractor's self-check is the hot path.",
    )

    # Left: how it works
    fx = 0.5
    nw = 6.0
    nh = 0.50

    add_section_header(slide, fx, 1.50, nw, "How it works")

    y = 1.85
    add_node(slide, fx, y, nw, nh,
             "Story Extractor emits stories (with self_check)",
             kind="agent", font_size=11)
    y2 = y + nh + 0.22
    add_node(slide, fx, y2, nw, nh,
             "Sampler (~10% rate, configurable per discipline)",
             kind="process", font_size=11)
    y3 = y2 + nh + 0.22
    add_node(slide, fx, y3, nw, nh,
             "Validator runs full rubric (shape-aware)",
             kind="agent", font_size=11)

    # branch on verdict
    y4 = y3 + nh + 0.30
    diamond_w = 1.6
    diamond_h = 0.85
    dx = fx + (nw - diamond_w) / 2
    add_diamond(slide, dx, y4, diamond_w, diamond_h, "verdict?",
                font_size=11)

    pass_y = y4 + diamond_h + 0.25
    pass_w = (nw - 0.20) / 2
    add_node(slide, fx, pass_y, pass_w, nh, "PASS → log audit row",
             kind="end", font_size=10.5)
    add_node(slide, fx + pass_w + 0.20, pass_y, pass_w, nh,
             "FAIL → log to review queue\n(priority; story stays in corpus)",
             kind="park", font_size=10)

    # arrows
    add_arrow(slide, fx + nw / 2, y + nh, fx + nw / 2, y2)
    add_arrow(slide, fx + nw / 2, y2 + nh, fx + nw / 2, y3)
    add_arrow(slide, fx + nw / 2, y3 + nh, fx + nw / 2, y4)
    # both branches emanate from the diamond's bottom vertex
    diamond_bottom_x = dx + diamond_w / 2
    diamond_bottom_y = y4 + diamond_h
    add_arrow(slide, diamond_bottom_x, diamond_bottom_y,
              fx + pass_w / 2, pass_y, label="pass",
              label_offset=(-0.55, -0.05))
    add_arrow(slide, diamond_bottom_x, diamond_bottom_y,
              fx + pass_w + 0.20 + pass_w / 2, pass_y, label="fail",
              label_offset=(0.10, -0.05))

    # Bottom-left rubric note
    note_y = pass_y + nh + 0.30
    add_box(slide, fx, note_y, nw, 0.85,
            "Rubric checks (shape-aware): actionability · testable acceptance "
            "criteria · shape alignment · persona resolved · SOP traceable · "
            "sop_breadth ≥ 2 (capability shape only).",
            fill=SIDE_FILL, line=SIDE_LINE,
            font_size=10.5, align=PP_ALIGN.LEFT)

    # Right: example pass + fail traces (code-box titles label each)
    rx = 6.85
    rw = 6.0

    pass_yaml = """validator_run:
  story_id: ST_023_07_a
  shape: workflow_stage_split
  rubric_version: v1
checks:
  - actionability:    pass    # has concrete action
  - testable_AC:      pass    # criteria are verifiable
  - shape_alignment:  pass    # has predecessor + successor
  - persona_resolved: pass    # lab_assistant in catalog
  - sop_traceable:    pass    # SOP_023.sec_4_2
verdict: PASS
sample_decision: kept"""

    fail_yaml = """validator_run:
  story_id: ST_119_03_b
  shape: capability
checks:
  - actionability:    pass
  - testable_AC:      FAIL    # criterion 2 unverifiable
  - shape_alignment:  pass
  - persona_resolved: pass
  - sop_traceable:    pass
  - sop_breadth_>=2:  pass    # capability-only check
verdict: FAIL
reason: |
  Criterion 2: "system handles specimens
  efficiently" — no measurable threshold.
action: log to review queue (priority: high)
note: story remains in corpus; reviewer
      decides confirm/correct/reject"""

    add_code_box(slide, rx, 1.85, rw, 2.30, pass_yaml,
                 font_size=8.5, title="PASS — story kept, audit logged")
    add_code_box(slide, rx, 4.45, rw, 2.65, fail_yaml,
                 font_size=8.5, title="FAIL — alert raised (story already emitted)")


# ============================================================================
# SLIDE 10 — HITL & confidence routing
# ============================================================================

def slide_ten(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "HITL — non-blocking review queue")
    add_subtitle(
        slide,
        "Every emission flows to corpus AND to the review queue.  Reviewer "
        "can confirm / correct / reject post-hoc — pipeline never waits.",
    )

    # ---- LEFT: routing flow ------------------------------------------------
    fx = 0.5
    full_w = 6.05

    add_section_header(slide, fx, 1.45, full_w, "Sources logged to the review queue")

    # 3 trigger boxes side-by-side
    trig_y = 1.90
    trig_h = 0.95
    trig_w = (full_w - 2 * 0.10) / 3
    triggers = [
        ("Story Extractor",
         "self_check FAIL\nor confidence < 0.85"),
        ("Sampled Validator",
         "rubric verdict\nFAIL"),
        ("Catalog Builder",
         "admit_confidence\n< 0.80 (any of\n5 catalogs)"),
    ]
    trig_centers_x = []
    for i, (head, body) in enumerate(triggers):
        x = fx + i * (trig_w + 0.10)
        # source label header
        add_text(slide, x + 0.04, trig_y, trig_w - 0.08, 0.26,
                 head, font_size=10.5, bold=True, color=ACCENT,
                 align=PP_ALIGN.CENTER)
        # trigger box
        add_box(slide, x, trig_y + 0.28, trig_w, trig_h - 0.28,
                body, fill=PARK_FILL, line=PARK_LINE,
                font_size=10.5, bold=False, align=PP_ALIGN.CENTER)
        trig_centers_x.append(x + trig_w / 2)

    # Review queue box (full width) — every emission logged
    queue_y = trig_y + trig_h + 0.40
    queue = add_node(slide, fx, queue_y, full_w, 0.50,
                     "Review queue — every emission logged; failures + low-conf prioritized",
                     kind="agent", font_size=10.5)

    # arrows from each trigger funneling into queue (top-center)
    queue_top_cx = fx + full_w / 2
    for cx in trig_centers_x:
        add_arrow(slide, cx, trig_y + trig_h + 0.01,
                  queue_top_cx, queue_y - 0.01)

    # Reviewer
    reviewer_y = queue_y + 0.50 + 0.25
    add_node(slide, fx, reviewer_y, full_w, 0.45,
             "Human reviewer", kind="side", font_size=11, bold=True)
    add_arrow(slide, queue_top_cx, queue_y + 0.50,
              queue_top_cx, reviewer_y)

    # Decision diamond
    diamond_w = 1.70
    diamond_h = 0.78
    diamond_x = fx + (full_w - diamond_w) / 2
    diamond_y = reviewer_y + 0.45 + 0.18
    add_diamond(slide, diamond_x, diamond_y, diamond_w, diamond_h,
                "decision", font_size=11)
    add_arrow(slide, queue_top_cx, reviewer_y + 0.45,
              queue_top_cx, diamond_y)

    # 3 outcome boxes
    out_y = diamond_y + diamond_h + 0.30
    out_h = 0.55
    outcomes = [
        ("Confirm",  "no change\n(item stays in corpus)",       END_FILL,  END_LINE),
        ("Correct",  "edit + apply\n→ manual emit-back",        SIDE_FILL, SIDE_LINE),
        ("Reject",   "remove from corpus\n→ manual emit-back",  PARK_FILL, PARK_LINE),
    ]
    out_centers_x = []
    diamond_bot_cx = diamond_x + diamond_w / 2
    diamond_bot_y = diamond_y + diamond_h
    for i, (head, body, fill, line) in enumerate(outcomes):
        x = fx + i * (trig_w + 0.10)
        add_text(slide, x + 0.04, out_y, trig_w - 0.08, 0.24,
                 head, font_size=11, bold=True, color=line,
                 align=PP_ALIGN.CENTER)
        add_box(slide, x, out_y + 0.26, trig_w, out_h - 0.04,
                body, fill=fill, line=line,
                font_size=10, align=PP_ALIGN.CENTER)
        out_centers_x.append(x + trig_w / 2)
        add_arrow(slide, diamond_bot_cx, diamond_bot_y,
                  x + trig_w / 2, out_y + 0.26)

    # bottom note
    note_y = out_y + 0.26 + out_h - 0.04 + 0.10
    add_text(slide, fx, note_y, full_w, 0.50,
             "hitl_status: auto · confirmed · corrected · rejected.  Tagger has no "
             "direct trigger — confidence rolls into Story Extractor.  Milestone 1: "
             "decisions applied to corpus manually (auto emit-back on roadmap).",
             font_size=8.5, italic=True, color=TEXT_GREY,
             align=PP_ALIGN.LEFT)

    # ---- RIGHT: confidence schema + queue example --------------------------
    rx = 6.85
    rw = 5.98

    confidence_yaml = """# At every agent boundary, a confidence is emitted.

catalog item     →  admit_confidence: 0.92    # rubric score (cluster or entity)
tagger output    →  confidence: 0.94          # rolls into story extraction_confidence
story extractor  →  extraction_confidence: 0.89
sampled validator →  verdict: PASS / FAIL  +  reason

# Review-priority thresholds (configurable per discipline)
catalog_admit_min:    0.80    # below → priority review
story_extract_min:    0.85    # below → priority review
validator_drift:      any FAIL → priority review
queue_priority:       weighted(confidence, impact)"""

    add_code_box(slide, rx, 1.85, rw, 2.20, confidence_yaml, font_size=8,
                 title="Confidence — where it surfaces")

    queue_yaml = """queue_id: RQ_2026_05_08_0042
source: story_extractor
ref:
  story_id: ST_119_03_b
  shape: capability
flags: ["low_confidence", "self_check_fail"]
confidence: 0.62
priority: high              # high | medium | low
state: pending_review       # pending_review | reviewed
# Reviewer decision (post-hoc):
reviewer_actions: [confirm, correct, reject]
# After decision, ref.hitl_status is updated;
# applying the decision to the corpus is a
# manual step in milestone 1."""

    add_code_box(slide, rx, 4.60, rw, 2.55, queue_yaml, font_size=8,
                 title="Review queue item — example")


# ============================================================================
# SLIDE 11 — Roadmap + ships
# ============================================================================

def slide_eleven(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "What ships + future roadmap")
    add_subtitle(
        slide,
        "Milestone 1 deliverables on the left. Roadmap items on the right — "
        "explicit, deferred to milestone 2.",
    )

    # Left: v1 ships
    lx, ly, lw = 0.5, 1.55, 6.0

    add_box(slide, lx, ly, lw, 0.50, "Milestone 1 — what ships",
            fill=END_FILL, line=END_LINE,
            font_size=14, bold=True, align=PP_ALIGN.CENTER)

    items_y = ly + 0.65
    items_h = 4.95
    add_box(slide, lx, items_y, lw, items_h, "",
            fill=CARD_FILL, line=CARD_LINE)

    sections = [
        ("Story corpus",
         "4-shape stories with extraction_confidence + hitl_status, "
         "self-checked at emission, traceable to SOP chunks."),
        ("5 catalogs",
         "business_capability · theme · epic · persona · test — each "
         "cold-discovered with per-item confidence; slug-based keys."),
        ("Review queue + audit log",
         "Every emission logged for human review; failures + low-conf "
         "prioritized.  Decisions captured immediately."),
        ("Sampled Validator QC",
         "~10% sample rate, async; drift surfaces in the review queue."),
        ("Onboarding playbook",
         "Hand the system SOPs for a new discipline; outputs follow with "
         "no SME pre-curation."),
    ]
    sy = items_y + 0.18
    for header, body in sections:
        add_text(slide, lx + 0.25, sy, lw - 0.50, 0.28,
                 "•  " + header, font_size=12, bold=True, color=ACCENT)
        add_text(slide, lx + 0.55, sy + 0.28, lw - 0.80, 0.62, body,
                 font_size=10.5, color=TEXT_DARK)
        sy += 0.92

    # Right: future roadmap
    rx, ry, rw = 6.83, 1.55, 6.0

    add_box(slide, rx, ry, rw, 0.50, "Roadmap — milestone 2 and beyond",
            fill=FUTURE_FILL, line=FUTURE_LINE,
            font_size=14, bold=True, align=PP_ALIGN.CENTER, font_color=FUTURE_LINE)

    add_box(slide, rx, ry + 0.65, rw, 4.95, "",
            fill=CARD_FILL, line=CARD_LINE)

    future = [
        ("Phase 3 — Cross-SOP Synthesis",
         "Cluster stories across SOPs; lift recurring patterns to capability "
         "stories at threshold ≥ 2 distinct SOPs."),
        ("Auto emit-back from reviewer decisions",
         "Apply confirm / correct / reject decisions to the corpus automatically; "
         "current milestone applies them manually."),
        ("Reviewer dashboard + queue analytics",
         "Production-grade reviewer UI: prioritization, batch operations, "
         "reviewer-performance analytics, decision provenance."),
        ("Threshold auto-tuning",
         "Replace fixed rubric and confidence thresholds with telemetry-driven "
         "tuning (precision and recall against reviewer decision history)."),
    ]
    ry2 = ry + 0.85
    for header, body in future:
        add_text(slide, rx + 0.25, ry2, rw - 0.50, 0.28,
                 "•  " + header, font_size=12, bold=True, color=FUTURE_LINE)
        add_text(slide, rx + 0.55, ry2 + 0.28, rw - 0.80, 0.85, body,
                 font_size=10.5, color=TEXT_DARK)
        ry2 += 1.13


# ============================================================================
# main
# ============================================================================

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_one(prs)
    slide_two(prs)
    slide_three(prs)
    slide_four(prs)
    slide_five(prs)
    slide_six(prs)
    slide_seven(prs)
    slide_eight(prs)
    slide_nine(prs)
    slide_ten(prs)
    slide_eleven(prs)

    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, start=1):
        add_footer(slide, i, total)

    out = Path(__file__).parent / "client_alignment_deck_v2.pptx"
    prs.save(str(out))
    print(f"wrote: {out}  ({total} slides)")


if __name__ == "__main__":
    main()
