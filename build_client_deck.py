"""
Generate the client alignment deck as native PowerPoint shapes.

Usage:
    python3 build_client_deck.py

Produces: client_alignment_deck.pptx alongside this script.

13 slides:
  1.  What we are building
  2.  What "prior discipline to learn from" means
  3.  The agents (cast of characters)
  4.  The full flow at a glance (3 phases)
  5.  Phase 1 — Setup: flowchart + theme catalog schema example
  6.  Epic Extractor — agent deep-dive (per-SOP + batch novelty)
  7.  Phase 2 — Per-SOP: flowchart (in ForEach container) + story schema
  8.  Phase 3 — Combine: flowchart + capability story example
  9.  The 4 story shapes — one example each (capability / stage-split /
      configuration-instance / cleanup)
  10. Inside two key agents — Validator + Dependency Resolver
  11. Validator at work — sample story walks through the checks
  12. What happens when something fails (3-level fallback + 3 queues)
  13. What ships and the alignment ask
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ---- palette ---------------------------------------------------------------
BOX_FILL = RGBColor(0xE7, 0xEF, 0xF8)
BOX_LINE = RGBColor(0x4A, 0x6E, 0xA0)
ACCENT = RGBColor(0x1F, 0x4E, 0x79)
TEXT_DARK = RGBColor(0x1F, 0x2D, 0x3D)
TEXT_GREY = RGBColor(0x55, 0x60, 0x6E)
ARROW_COLOR = RGBColor(0x4A, 0x6E, 0xA0)
SYS_FILL = RGBColor(0xFA, 0xE9, 0xCB)
SYS_LINE = RGBColor(0xC9, 0x95, 0x2C)
DECISION_FILL = RGBColor(0xFC, 0xE7, 0xC4)
DECISION_LINE = RGBColor(0xC9, 0x95, 0x2C)
START_FILL = RGBColor(0xE0, 0xEC, 0xF8)
START_LINE = RGBColor(0x4A, 0x6E, 0xA0)
END_FILL = RGBColor(0xDC, 0xE9, 0xD7)
END_LINE = RGBColor(0x5C, 0x83, 0x47)
OUT_FILL = RGBColor(0xDC, 0xE9, 0xD7)
OUT_LINE = RGBColor(0x5C, 0x83, 0x47)
HUB_FILL = RGBColor(0xFC, 0xE0, 0xC4)
HUB_LINE = RGBColor(0xB8, 0x6A, 0x29)
PHASE_FILL = RGBColor(0xEE, 0xE7, 0xF8)
PHASE_LINE = RGBColor(0x6E, 0x4F, 0xA0)
CODE_FILL = RGBColor(0xF6, 0xF6, 0xF6)
CODE_LINE = RGBColor(0xC0, 0xC8, 0xD0)
SIDE_BRANCH_FILL = RGBColor(0xEC, 0xF4, 0xE8)
SIDE_BRANCH_LINE = RGBColor(0x5C, 0x83, 0x47)
PARK_FILL = RGBColor(0xF8, 0xE3, 0xE3)
PARK_LINE = RGBColor(0xA8, 0x4A, 0x4A)


# ---- text/shape helpers ---------------------------------------------------

def _set_text(shape, text, font_size, bold, color, align,
              anchor=MSO_ANCHOR.MIDDLE, font_name="Calibri", italic=False):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name


def add_box(slide, left, top, width, height, text, fill=BOX_FILL, line=BOX_LINE,
            font_size=12, bold=False, font_color=TEXT_DARK, align=PP_ALIGN.CENTER,
            shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    box = slide.shapes.add_shape(shape_type, Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line
    box.line.width = Pt(1.0)
    _set_text(box, text, font_size, bold, font_color, align)
    return box


def add_diamond(slide, left, top, width, height, text,
                fill=DECISION_FILL, line=DECISION_LINE,
                font_size=11, bold=True, color=TEXT_DARK):
    d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(left), Inches(top),
                               Inches(width), Inches(height))
    d.fill.solid()
    d.fill.fore_color.rgb = fill
    d.line.color.rgb = line
    d.line.width = Pt(1.25)
    _set_text(d, text, font_size, bold, color, PP_ALIGN.CENTER)
    return d


def add_node(slide, left, top, width, height, text, kind="process",
             font_size=11, bold=False):
    """kind ∈ {start, end, process, agent, side, park}."""
    if kind == "start":
        return add_box(slide, left, top, width, height, text,
                       fill=START_FILL, line=START_LINE,
                       font_size=font_size, bold=True,
                       shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    if kind == "end":
        return add_box(slide, left, top, width, height, text,
                       fill=END_FILL, line=END_LINE,
                       font_size=font_size, bold=True,
                       shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    if kind == "agent":
        return add_box(slide, left, top, width, height, text,
                       fill=SYS_FILL, line=SYS_LINE,
                       font_size=font_size, bold=True)
    if kind == "side":
        return add_box(slide, left, top, width, height, text,
                       fill=SIDE_BRANCH_FILL, line=SIDE_BRANCH_LINE,
                       font_size=font_size, bold=bold)
    if kind == "park":
        return add_box(slide, left, top, width, height, text,
                       fill=PARK_FILL, line=PARK_LINE,
                       font_size=font_size, bold=bold)
    return add_box(slide, left, top, width, height, text,
                   font_size=font_size, bold=bold)


def add_text(slide, left, top, width, height, text, font_size=14,
             bold=False, italic=False, color=TEXT_DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font_name="Calibri"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return tb


def add_bullets(slide, left, top, width, height, items, font_size=11,
                color=TEXT_DARK):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = "•  " + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
        p.space_after = Pt(3)
    return tb


def add_arrow(slide, x1, y1, x2, y2, color=ARROW_COLOR, weight=2.0,
              connector_type=MSO_CONNECTOR.STRAIGHT, label=None,
              label_offset=(0, -0.18), label_color=TEXT_GREY,
              label_size=10, label_bold=False):
    conn = slide.shapes.add_connector(connector_type,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    line_elem = conn.line._get_or_add_ln()
    tail_end = etree.SubElement(line_elem, qn("a:tailEnd"))
    tail_end.set("type", "triangle")
    tail_end.set("w", "med")
    tail_end.set("h", "med")
    if label:
        lx = (x1 + x2) / 2 + label_offset[0]
        ly = (y1 + y2) / 2 + label_offset[1]
        tb = slide.shapes.add_textbox(Inches(lx), Inches(ly),
                                      Inches(0.9), Inches(0.3))
        tf = tb.text_frame
        tf.margin_left = Inches(0.04)
        tf.margin_right = Inches(0.04)
        tf.margin_top = Inches(0.0)
        tf.margin_bottom = Inches(0.0)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = label
        run.font.size = Pt(label_size)
        run.font.bold = label_bold
        run.font.color.rgb = label_color
        run.font.name = "Calibri"
    return conn


def add_title(slide, text, color=ACCENT):
    add_text(slide, 0.5, 0.32, 12.5, 0.6, text,
             font_size=26, bold=True, color=color)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.5), Inches(0.95),
                                 Inches(1.6), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def add_footer(slide, slide_num, total_slides, deck_label="LIMS Story Generation — Client Alignment"):
    """Subtle footer: thin horizontal line + deck label (left) + slide number (right)."""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.5), Inches(7.18),
                                  Inches(12.33), Inches(0.012))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0xD0, 0xD6, 0xDC)
    line.line.fill.background()
    add_text(slide, 0.5, 7.22, 8.0, 0.25, deck_label,
             font_size=8, italic=True, color=RGBColor(0x90, 0x97, 0x9F),
             align=PP_ALIGN.LEFT)
    add_text(slide, 11.83, 7.22, 1.0, 0.25,
             f"{slide_num} / {total_slides}",
             font_size=8, color=RGBColor(0x90, 0x97, 0x9F),
             align=PP_ALIGN.RIGHT)


def add_code_box(slide, left, top, width, height, content, font_size=10,
                 title=None, fill=CODE_FILL, line=CODE_LINE):
    """Code-style block with monospace font for schema examples."""
    if title:
        add_text(slide, left, top - 0.32, width, 0.3, title,
                 font_size=11, bold=True, color=ACCENT)
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line
    box.line.width = Pt(0.75)
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    lines = content.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = ln
        run.font.size = Pt(font_size)
        run.font.color.rgb = TEXT_DARK
        run.font.name = "Consolas"
        p.space_after = Pt(0)
    return box


# ---- slide 1 --------------------------------------------------------------

def slide_one(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "What we are building")

    add_text(slide, 0.5, 1.1, 12.5, 0.6,
             "A system that turns SOPs (Standard Operating Procedures) into a "
             "ready-to-use backlog of dev work — without an expert reviewing each item.",
             font_size=15, color=TEXT_DARK)

    box_w, box_h, y = 3.4, 1.4, 2.0
    x1, x2, x3 = 0.6, 4.95, 9.3
    add_box(slide, x1, y, box_w, box_h,
            "SOPs\n(Standard Operating Procedures)",
            font_size=14, bold=True)
    add_box(slide, x2, y, box_w, box_h, "Our system",
            fill=SYS_FILL, line=SYS_LINE, font_size=18, bold=True)
    add_box(slide, x3, y, box_w, box_h,
            "Dev backlog\n(Jira: Epics + Stories)",
            fill=OUT_FILL, line=OUT_LINE, font_size=14, bold=True)
    add_arrow(slide, x1 + box_w + 0.05, y + box_h / 2,
              x2 - 0.05, y + box_h / 2)
    add_arrow(slide, x2 + box_w + 0.05, y + box_h / 2,
              x3 - 0.05, y + box_h / 2)
    add_text(slide, x2, y + box_h + 0.1, box_w, 0.4,
             "Plus: configuration files + audit log",
             font_size=11, italic=True, color=TEXT_GREY,
             align=PP_ALIGN.CENTER)

    add_text(slide, 0.5, 4.4, 12.5, 0.4,
             "Why this matters", font_size=18, bold=True, color=ACCENT)
    add_bullets(slide, 0.7, 4.9, 12.0, 1.4, [
        "Today: every new SOP → manual expert review → hand-written stories. Slow, costly, inconsistent.",
        "With this system: SOPs in, structured stories out. Quality is built into the rules, not into one expert's review.",
        "Same system, any clinical discipline. Each new discipline is a configuration change, not an engineering project.",
    ], font_size=13)

    # ----- mini visual: what's actually in the dev backlog -----
    tree_y = 6.45
    add_text(slide, 0.5, tree_y, 12.5, 0.20,
             "What's inside the Dev backlog (Jira hierarchy)",
             font_size=11, italic=True, bold=True, color=ACCENT)
    row_y = tree_y + 0.27
    row_h = 0.45

    ep_x, ep_w = 2.4, 1.4
    add_box(slide, ep_x, row_y, ep_w, row_h, "Epic",
            fill=PHASE_FILL, line=PHASE_LINE, font_size=12, bold=True)
    st_w = 1.4
    gap = 0.15
    st_x0 = ep_x + ep_w + 0.50
    for i in range(3):
        x = st_x0 + i * (st_w + gap)
        add_box(slide, x, row_y, st_w, row_h, "Story",
                fill=OUT_FILL, line=OUT_LINE, font_size=12, bold=True)
        add_arrow(slide, ep_x + ep_w, row_y + row_h / 2,
                  x, row_y + row_h / 2, weight=1.0)
    # Task box (faded, marked out of scope)
    task_x = st_x0 + 3 * (st_w + gap) + 0.30
    task_w = 2.8
    add_box(slide, task_x, row_y, task_w, row_h,
            "Task    (out of scope — dev team owns)",
            fill=RGBColor(0xF1, 0xF1, 0xF1),
            line=RGBColor(0xB0, 0xB0, 0xB0),
            font_size=10, bold=False,
            font_color=RGBColor(0x80, 0x80, 0x80))


# ---- slide 2 --------------------------------------------------------------

def slide_two(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "What does \"prior discipline to learn from\" mean?")

    add_text(slide, 0.5, 1.1, 12.5, 0.5,
             "It's a teaching reference — a discipline whose dev work has already been "
             "validated. For our project, that's Cytology.",
             font_size=14, color=TEXT_DARK)

    hub_w, hub_h = 2.6, 1.1
    hub_x = (13.333 - hub_w) / 2
    hub_y = 2.0
    add_box(slide, hub_x, hub_y, hub_w, hub_h,
            "Prior discipline\n(Cytology)",
            fill=HUB_FILL, line=HUB_LINE, font_size=15, bold=True)

    # MANDATORY row (top): two larger boxes
    mand_w, mand_h = 3.6, 1.35
    mand_y = 2.00
    add_box(slide, 0.5, mand_y, mand_w, mand_h,
            "Theme starter set\nVocabulary of workflow categories\n(Pre-Analytic, Analytic, Reporting, …)",
            fill=BOX_FILL, line=BOX_LINE, font_size=12, bold=True)
    add_box(slide, 13.333 - mand_w - 0.5, mand_y, mand_w, mand_h,
            "Epic starter set\nHigh-level dev groupings\n(Specimen Receiving, Reporting, Billing, …)",
            fill=BOX_FILL, line=BOX_LINE, font_size=12, bold=True)
    # mandatory badges
    add_text(slide, 0.5, mand_y - 0.3, 1.4, 0.3, "MANDATORY",
             font_size=9, bold=True, color=ACCENT)
    add_text(slide, 13.333 - mand_w - 0.5, mand_y - 0.3, 1.4, 0.3,
             "MANDATORY", font_size=9, bold=True, color=ACCENT)

    # arrows from hub to two mandatory boxes
    add_arrow(slide, hub_x, hub_y + hub_h * 0.45,
              0.5 + mand_w, mand_y + mand_h * 0.5,
              color=HUB_LINE, weight=2.0)
    add_arrow(slide, hub_x + hub_w, hub_y + hub_h * 0.45,
              13.333 - mand_w - 0.5, mand_y + mand_h * 0.5,
              color=HUB_LINE, weight=2.0)

    # OPTIONAL row (bottom): single box, dimmer
    opt_w, opt_h = 6.0, 1.2
    opt_y = 4.05
    opt_x = (13.333 - opt_w) / 2
    add_box(slide, opt_x, opt_y, opt_w, opt_h,
            "Story exemplars\nFew-shot examples for the Story Extractor\n"
            "Used IF prior-discipline Jira is accessible — system runs without them",
            fill=RGBColor(0xF6, 0xF6, 0xF6), line=RGBColor(0xA0, 0xA0, 0xA0),
            font_size=11)
    add_text(slide, opt_x, opt_y - 0.3, 2.0, 0.3,
             "OPTIONAL (enhancement)", font_size=9, bold=True,
             color=RGBColor(0x80, 0x80, 0x80))
    # arrow from hub down to optional box (dashed-feel via lighter color)
    add_arrow(slide, hub_x + hub_w / 2, hub_y + hub_h,
              opt_x + opt_w / 2, opt_y,
              color=RGBColor(0xA0, 0xA0, 0xA0), weight=1.5)

    # bottom: not inherited
    add_text(slide, 0.5, 5.55, 12.5, 0.35,
             "Not inherited from the prior discipline",
             font_size=14, bold=True, color=ACCENT)
    add_text(slide, 0.7, 5.92, 12.5, 0.45,
             "Tests   •   Roles   •   The new discipline's own SOPs   •   "
             "Tasks (out of scope; dev team owns decomposition)",
             font_size=11, color=TEXT_GREY)
    add_text(slide, 0.5, 6.30, 12.5, 0.45,
             "Sourcing:  theme + epic catalogs are typically snapshotted from a validated prior discipline (Cytology, in our case).  "
             "If no prior is available, a one-time cold-start SOP analysis discovers themes and epics from scratch — that's how Cytology's catalog was built originally.",
             font_size=10, italic=True, color=ACCENT)
    add_text(slide, 0.5, 6.78, 12.5, 0.30,
             "Quality-check thresholds use defaults tuned by first-run telemetry — no holdout calibration against the prior discipline is needed.",
             font_size=9, italic=True, color=TEXT_GREY)


# ---- slide 3: the agents --------------------------------------------------

def slide_three(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "The agents — cast of characters")

    add_text(slide, 0.5, 1.1, 12.5, 0.5,
             "Six agents do the work. Plus a quorum — 5 conditional checks evaluated "
             "when admitting new categories.",
             font_size=13, color=TEXT_DARK)

    # table of agents
    headers = ["#", "Agent", "What it does", "When it runs"]
    rows = [
        ["1", "Theme Discovery",
         "Builds the new discipline's theme catalog (categories) by classifying "
         "sample SOPs against the prior discipline's themes; new themes "
         "emerge only with evidence.",
         "Once per new discipline\n(re-runs only on alarm)"],
        ["2", "Epic Extractor (conditioned)",
         "Drafts high-level epics from each SOP; matches each draft "
         "against the prior discipline's epic catalog. Matches inherit; "
         "non-matches are clustered for novelty review.",
         "Once per SOP\n+ batch novelty pass"],
        ["3", "Story Extractor",
         "Drafts user stories from SOP chunks using schema + shape "
         "definitions + closed-enum catalogs. Optionally retrieves "
         "(test, persona, stage, shape)-matched exemplars from the prior "
         "discipline's Jira if it's available — improves style fidelity.",
         "Per epic, per SOP"],
        ["4", "Validator (quality check)",
         "Quality-checks every emitted story against shape-specific rules "
         "(MUST/SHALL for capabilities; concrete values for configurations; "
         "named artifact for cleanups). Failures auto-park.",
         "Twice per story\n(after extract + after synthesis)"],
        ["5", "Cross-SOP Synthesis",
         "After all SOPs are processed, finds patterns recurring across "
         "≥ 2 SOPs and lifts them into broader \"capability stories\" "
         "with parameters.",
         "Once per batch"],
        ["6", "Dependency Resolver",
         "Resolves dependencies between stories into a topological order — "
         "the final ordered backlog the dev team consumes.",
         "Once per batch"],
    ]

    n_rows = len(rows) + 1
    table_x = 0.4
    table_y = 1.65
    table_w = 12.55
    table_h = 0.4 + (n_rows - 1) * 0.6
    tbl_shape = slide.shapes.add_table(n_rows, len(headers),
                                       Inches(table_x), Inches(table_y),
                                       Inches(table_w), Inches(table_h))
    table = tbl_shape.table

    widths_in = [0.5, 2.4, 6.5, 3.15]
    for i, w in enumerate(widths_in):
        table.columns[i].width = Inches(w)

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        tf = cell.text_frame
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if j > 0 else PP_ALIGN.CENTER
        run = p.add_run()
        run.text = h
        run.font.size = Pt(12)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        run.font.name = "Calibri"
    table.rows[0].height = Inches(0.4)

    for i, row in enumerate(rows):
        table.rows[i + 1].height = Inches(0.6)
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                RGBColor(0xFA, 0xE9, 0xCB) if j == 1 else
                RGBColor(0xF8, 0xFA, 0xFC) if i % 2 == 0 else
                RGBColor(0xFF, 0xFF, 0xFF)
            )
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.08)
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j in (0, 3) else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = val
            run.font.size = Pt(10) if j == 2 else Pt(10)
            run.font.bold = (j == 1)
            run.font.color.rgb = TEXT_DARK
            run.font.name = "Calibri"

    callout_y = table_y + table_h + 0.18
    callout_h = 7.05 - callout_y
    # Quorum callout — generalized to conditional logic (5 parallel checks)
    add_text(slide, 0.4, callout_y, 12.55, 0.3,
             "Quorum  —  conditional logic with 5 parallel checks  (used inside Theme Discovery and Epic Extractor)",
             font_size=12, bold=True, color=ACCENT)
    # 5 small boxes side by side — condition labels
    qp_y = callout_y + 0.32
    qp_h = 0.55
    qp_gap = 0.08
    qp_total_w = 12.55
    qp_w = (qp_total_w - 4 * qp_gap) / 5
    conditions = [
        "1. Coherence\nas a new theme",
        "2. Distinctness\nfrom existing",
        "3. Description\nsharpness",
        "4. Cluster member\nconsistency",
        "5. No catalog\nfragmentation",
    ]
    for i, cond in enumerate(conditions):
        x = 0.4 + i * (qp_w + qp_gap)
        add_box(slide, x, qp_y, qp_w, qp_h, cond,
                fill=SYS_FILL, line=SYS_LINE,
                font_size=9, bold=False)
    # below the 5 boxes: vote-aggregation note
    note_y = qp_y + qp_h + 0.05
    note_h = 7.05 - note_y
    add_text(slide, 0.4, note_y, 12.55, note_h,
             "5 conditions evaluated in parallel. Need 3 of 5 to clear plus consensus on the proposed entry.  "
             "Otherwise the candidate goes to the review pile for the next round.",
             font_size=10, italic=True, color=TEXT_DARK,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ---- slide 4: full flow at a glance ---------------------------------------

def slide_four(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "The full flow — three phases")

    add_text(slide, 0.5, 1.1, 12.5, 0.5,
             "End-to-end view. Each phase is detailed in the next three slides.",
             font_size=13, italic=True, color=TEXT_GREY)

    phase_w, phase_h = 3.95, 3.4
    phase_y = 2.0
    px = [0.4, 4.7, 9.0]
    titles = [
        ("Phase 1\nSetup",
         "Build theme & epic catalogs for the new discipline, conditioned on the prior discipline.",
         "Agents: Theme Discovery, Epic Extractor, Quorum panel.",
         "Runs once per new discipline.\nRepeats only if alarm fires."),
        ("Phase 2\nPer-SOP",
         "Read each SOP, draft stories, quality-check.",
         "Agents: Epic Extractor, Story Extractor, Validator.",
         "Runs for every new SOP."),
        ("Phase 3\nCombine + ship",
         "Look across all SOPs, lift recurring patterns, generate outputs.",
         "Agents: Cross-SOP Synthesis, Validator, Dependency Resolver.",
         "Runs after a batch of SOPs is processed."),
    ]
    for x, (head, mid, agents, foot) in zip(px, titles):
        add_box(slide, x, phase_y, phase_w, 0.85, head,
                fill=PHASE_FILL, line=PHASE_LINE,
                font_size=15, bold=True, font_color=ACCENT)
        add_text(slide, x + 0.15, phase_y + 1.0, phase_w - 0.3, 1.1,
                 mid, font_size=12, color=TEXT_DARK,
                 align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.15, phase_y + 2.1, phase_w - 0.3, 0.65,
                 agents, font_size=10, italic=True, color=TEXT_GREY,
                 align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.15, phase_y + 2.7, phase_w - 0.3, 0.6,
                 foot, font_size=10, italic=True, color=TEXT_GREY,
                 align=PP_ALIGN.CENTER)

    arrow_y = phase_y + 0.4
    add_arrow(slide, px[0] + phase_w + 0.05, arrow_y, px[1] - 0.05, arrow_y, weight=2.5)
    add_arrow(slide, px[1] + phase_w + 0.05, arrow_y, px[2] - 0.05, arrow_y, weight=2.5)

    cont_y = 5.85
    add_box(slide, 0.4, cont_y, 12.55, 0.6,
            "Continuous (alongside all phases): "
            "drift report   •   alarm-driven re-runs when too many items go to the review pile",
            fill=RGBColor(0xF1, 0xF5, 0xFA), line=BOX_LINE,
            font_size=11, font_color=TEXT_DARK,
            align=PP_ALIGN.CENTER, shape_type=MSO_SHAPE.RECTANGLE)


# ---- slide 5: phase 1 — flowchart + theme catalog example ----------------

def slide_five(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Phase 1 — Setup: building the theme catalog")

    add_text(slide, 0.5, 1.05, 12.5, 0.4,
             "Run once for a new discipline. Result: a theme catalog mostly inherited "
             "from the prior discipline + 0–N new themes (each with evidence).",
             font_size=12, color=TEXT_DARK)

    # ----- left half: flowchart (compact for footer fit) -----
    fx = 0.5
    nw = 3.0
    nh = 0.42
    cx = fx + nw / 2
    dh = 0.55
    gap = 0.10
    post_diamond = 0.22

    y1 = 1.55
    add_node(slide, fx, y1, nw, nh, "Sample SOPs (~20)\n+ prior discipline's themes",
             kind="start", font_size=9, bold=True)
    # NEW: Section SOP into logical chunks (separate from tagging)
    y_chunk = y1 + nh + gap
    add_node(slide, fx, y_chunk, nw, nh, "Section each SOP\ninto logical chunks",
             kind="process", font_size=9)
    y2 = y_chunk + nh + gap
    add_node(slide, fx, y2, nw, nh, "Tag chunks:\ntheme, test, persona, stage",
             kind="process", font_size=9)
    y3 = y2 + nh + gap
    add_node(slide, fx, y3, nw, nh, "Compare each chunk to\nprior discipline's themes",
             kind="process", font_size=9)
    y4 = y3 + nh + gap
    add_diamond(slide, fx + 0.4, y4, nw - 0.8, dh,
                "Score ≥ match\nthreshold?", font_size=9)
    inh_x = fx + nw + 0.4
    inh_y = y4 + (dh - nh) / 2
    add_node(slide, inh_x, inh_y, 2.5, nh,
             "Inherit theme\n(into new catalog)",
             kind="side", font_size=9, bold=True)
    y6 = y4 + dh + post_diamond
    add_node(slide, fx, y6, nw, nh, "Cluster residual into candidates",
             kind="process", font_size=9)
    # 5 parallel conditional checks (relabeled from "reviewer agents")
    y7 = y6 + nh + gap
    qph = 0.42
    qpw_total = nw
    qpw_each = (qpw_total - 4 * 0.04) / 5
    for i in range(5):
        qx = fx + i * (qpw_each + 0.04)
        add_box(slide, qx, y7, qpw_each, qph, str(i + 1),
                fill=SYS_FILL, line=SYS_LINE,
                font_size=11, bold=True)
    # label below the row of 5 — generalized to conditional logic
    add_text(slide, fx, y7 + qph + 0.01, nw, 0.18,
             "5 conditional checks evaluated in parallel",
             font_size=8, italic=True, color=TEXT_GREY,
             align=PP_ALIGN.CENTER)
    y8 = y7 + qph + 0.25
    add_diamond(slide, fx + 0.4, y8, nw - 0.8, dh,
                "3 of 5 conditions\ncleared?", font_size=9)
    park_x = fx + nw + 0.4
    park_y = y8 + (dh - nh) / 2
    add_node(slide, park_x, park_y, 2.5, nh,
             "Park to review pile\n(retry next cycle)",
             kind="park", font_size=9, bold=True)
    y9 = y8 + dh + post_diamond
    add_node(slide, fx, y9, nw, nh, "New discipline's theme catalog",
             kind="end", font_size=10, bold=True)

    # arrows in the flowchart
    add_arrow(slide, cx, y1 + nh, cx, y_chunk)
    add_arrow(slide, cx, y_chunk + nh, cx, y2)
    add_arrow(slide, cx, y2 + nh, cx, y3)
    add_arrow(slide, cx, y3 + nh, cx, y4)
    # diamond branches
    add_arrow(slide, fx + nw - 0.4, y4 + dh / 2, inh_x, inh_y + nh / 2,
              label="yes\n(inherit)", label_offset=(0.05, -0.32),
              label_color=END_LINE, label_bold=True)
    add_arrow(slide, cx, y4 + dh, cx, y6,
              label="no\n(residual)", label_offset=(0.08, -0.18),
              label_color=PARK_LINE, label_bold=True)
    # cluster → top-center of the 5-agent row
    add_arrow(slide, cx, y6 + nh, cx, y7)
    # bottom-center of the 5-agent row → quorum decision diamond
    add_arrow(slide, cx, y7 + qph + 0.22, cx, y8)
    # quorum branches
    add_arrow(slide, fx + nw - 0.4, y8 + dh / 2, park_x, park_y + nh / 2,
              label="no", label_offset=(0.05, -0.18),
              label_color=PARK_LINE, label_bold=True)
    add_arrow(slide, cx, y8 + dh, cx, y9,
              label="yes (admit)", label_offset=(0.05, -0.18),
              label_color=END_LINE, label_bold=True)
    # inherit branch converges into final node — elbow: down from inherit,
    # then left into the right edge of the final node
    add_arrow(slide, inh_x + 1.25, inh_y + nh,        # start: bottom-center of inherit box
              fx + nw, y9 + nh / 2,                    # end: right-edge-middle of final node
              color=SIDE_BRANCH_LINE, weight=1.25,
              connector_type=MSO_CONNECTOR.ELBOW)

    # ----- right half: theme catalog YAML example -----
    code_x = 7.7
    code_y = 1.6
    code_w = 5.4
    code_h = 4.0
    yaml_text = """theme_catalog:
  catalog_id:     micro_v1
  parent_catalog: cyto_v1     # prior discipline

  inherited_themes:           # carried over by Pass 1
    - G1, G2, G3, G4, G5, G6, G7, G8

  novel_themes:               # admitted by quorum
    - id: MI1
      name: Susceptibility Testing
      evidence: 42 chunks
      votes: [admit×5]    decision: admit
    - id: MI2
      name: Biosafety Containment
      evidence: 28 chunks
      votes: [admit×4, fold→G6]    decision: admit

  folded_in:                  # M=3 majority on fold target
    - id: MI3
      name: Critical Result Notification
      evidence: 15 chunks
      votes: [fold→G4×3, admit×2]    decision: fold→G4

  unclassified_bucket:        # left for next round
    count: 38                 # = 6.3% of 600 (above 5% alarm)"""
    add_code_box(slide, code_x, code_y, code_w, code_h, yaml_text,
                 font_size=8.5, title="Example output: theme catalog (simplified)")

    # ----- Pass 1 distribution visual (replaces prior explainer text) -----
    dist_y = code_y + code_h + 0.10
    add_text(slide, code_x, dist_y, code_w, 0.22,
             "What Pass 1 produced  (Microbiology classified vs Cytology's themes)",
             font_size=10, bold=True, color=ACCENT)
    bucket_y = dist_y + 0.27
    bucket_h = 0.42
    inherited = [
        ("G1", 145), ("G2", 88), ("G3", 60), ("G4", 52),
        ("G5", 44), ("G6", 30), ("G7", 38), ("G8", 20),
    ]
    residual = ("Residual", 123)
    n = len(inherited)
    avail_w = code_w
    div_gap = 0.12
    inh_total_w = avail_w * 0.62
    res_w = avail_w * 0.34
    inh_box_w = (inh_total_w - (n - 1) * 0.04) / n
    for i, (name, cnt) in enumerate(inherited):
        x = code_x + i * (inh_box_w + 0.04)
        add_box(slide, x, bucket_y, inh_box_w, bucket_h,
                f"{name}\n{cnt}",
                fill=BOX_FILL, line=BOX_LINE,
                font_size=8, bold=True)
    res_x = code_x + inh_total_w + div_gap
    add_box(slide, res_x, bucket_y, res_w - 0.04, bucket_h,
            f"{residual[0]}\n{residual[1]} chunks",
            fill=RGBColor(0xE8, 0xE2, 0xD4), line=RGBColor(0x90, 0x88, 0x78),
            font_size=9, bold=True)
    # tiny caption hint just under the bucket row
    add_text(slide, code_x, bucket_y + bucket_h + 0.02, inh_total_w, 0.18,
             "inherited from Cytology (477 chunks total)",
             font_size=8, italic=True, color=TEXT_GREY,
             align=PP_ALIGN.CENTER)
    add_text(slide, res_x, bucket_y + bucket_h + 0.02, res_w - 0.04, 0.18,
             "feeds Pass 2 cluster",
             font_size=8, italic=True, color=RGBColor(0x70, 0x68, 0x58),
             align=PP_ALIGN.CENTER)
    # Multi-label clarification footnote
    note_y = bucket_y + bucket_h + 0.24
    add_text(slide, code_x, note_y, code_w, 0.40,
             "Note:  Pass 1 here is single-label — each chunk goes to "
             "its best-matching theme above τ_match=0.65 (or to residual). "
             "Runtime story / chunk theme tags are multi-label per D3 "
             "(themes are soft tags, not a forced taxonomy).",
             font_size=8, italic=True, color=TEXT_GREY)


# ---- slide 6: epic extractor — agent deep-dive ---------------------------

def slide_epic_extractor(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Epic Extractor — building the epic catalog")

    add_text(slide, 0.5, 1.05, 12.5, 0.4,
             "Same conditioned-discovery mechanism as Theme Discovery (slide 5), "
             "applied at the epic level. Runs during Phase 2 per-SOP and at the "
             "Phase 3 batch boundary — not a separate phase.",
             font_size=12, color=TEXT_DARK)

    # ----- left half: flowchart -----
    fx = 0.5
    nw = 3.0
    nh = 0.42
    cx = fx + nw / 2
    dh = 0.55
    gap = 0.10
    post_diamond = 0.22

    # per-SOP (top) section label
    add_text(slide, fx, 1.50, nw, 0.20,
             "Runs PER SOP (in parallel)",
             font_size=8, italic=True, bold=True, color=TEXT_GREY,
             align=PP_ALIGN.CENTER)

    y1 = 1.72
    add_node(slide, fx, y1, nw, nh,
             "Epic Extractor:\ndraft epics from this SOP",
             kind="agent", font_size=9, bold=True)
    y2 = y1 + nh + gap
    add_node(slide, fx, y2, nw, nh,
             "Match each draft against\nprior epic catalog (cyto_epic_v1)",
             kind="process", font_size=9)
    y3 = y2 + nh + gap
    add_diamond(slide, fx + 0.4, y3, nw - 0.8, dh,
                "Score ≥ τ_match?", font_size=9)
    # YES branch — inherit (with epic_analog populated by construction)
    inh_x = fx + nw + 0.4
    inh_y = y3 + (dh - nh) / 2
    add_node(slide, inh_x, inh_y, 2.5, nh,
             "Inherit\n(epic_analog populated\nby construction)",
             kind="side", font_size=8.5, bold=True)

    # NO branch — residual draft, continues down
    y4 = y3 + dh + post_diamond
    add_node(slide, fx, y4, nw, nh,
             "Residual draft\n(per-SOP output → batch pool)",
             kind="process", font_size=9)

    # batch boundary — separator label
    sep_y = y4 + nh + 0.06
    sep = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                     Inches(fx), Inches(sep_y),
                                     Inches(fx + nw), Inches(sep_y))
    sep.line.color.rgb = RGBColor(0xC0, 0xC8, 0xD0)
    sep.line.width = Pt(0.75)
    add_text(slide, fx, sep_y + 0.01, nw, 0.18,
             "After all SOPs processed (BATCH BOUNDARY)",
             font_size=8, italic=True, bold=True, color=TEXT_GREY,
             align=PP_ALIGN.CENTER)

    y5 = sep_y + 0.22
    add_node(slide, fx, y5, nw, nh,
             "Cluster residual drafts\nacross all SOPs",
             kind="process", font_size=9)

    # 5 parallel conditional checks
    y6 = y5 + nh + gap
    qph = 0.40
    qpw_each = (nw - 4 * 0.04) / 5
    for i in range(5):
        qx = fx + i * (qpw_each + 0.04)
        add_box(slide, qx, y6, qpw_each, qph, str(i + 1),
                fill=SYS_FILL, line=SYS_LINE,
                font_size=11, bold=True)
    add_text(slide, fx, y6 + qph + 0.01, nw, 0.18,
             "Same 5 conditional checks as slide 5",
             font_size=8, italic=True, color=TEXT_GREY,
             align=PP_ALIGN.CENTER)

    y7 = y6 + qph + 0.25
    add_diamond(slide, fx + 0.4, y7, nw - 0.8, dh,
                "3 of 5 conditions\ncleared?", font_size=9)
    # NO branch — park to E0
    park_x = fx + nw + 0.4
    park_y = y7 + (dh - nh) / 2
    add_node(slide, park_x, park_y, 2.5, nh,
             "Park to E0\n(unclassified bucket)",
             kind="park", font_size=9, bold=True)

    # final node
    y8 = y7 + dh + post_diamond
    add_node(slide, fx, y8, nw, nh,
             "Epic catalog\n(inherited + novel)",
             kind="end", font_size=10, bold=True)

    # arrows
    add_arrow(slide, cx, y1 + nh, cx, y2)
    add_arrow(slide, cx, y2 + nh, cx, y3)
    # decision 1 branches
    add_arrow(slide, fx + nw - 0.4, y3 + dh / 2, inh_x, inh_y + nh / 2,
              label="yes\n(inherit)", label_offset=(0.05, -0.32),
              label_color=END_LINE, label_bold=True)
    add_arrow(slide, cx, y3 + dh, cx, y4,
              label="no\n(residual)", label_offset=(0.08, -0.18),
              label_color=PARK_LINE, label_bold=True)
    # residual → cluster (crosses the dashed separator visually)
    add_arrow(slide, cx, y4 + nh, cx, y5)
    # cluster → 5-check row
    add_arrow(slide, cx, y5 + nh, cx, y6)
    # 5-check row → quorum decision
    add_arrow(slide, cx, y6 + qph + 0.22, cx, y7)
    # decision 2 branches
    add_arrow(slide, fx + nw - 0.4, y7 + dh / 2, park_x, park_y + nh / 2,
              label="no", label_offset=(0.05, -0.18),
              label_color=PARK_LINE, label_bold=True)
    add_arrow(slide, cx, y7 + dh, cx, y8,
              label="yes (admit)", label_offset=(0.05, -0.18),
              label_color=END_LINE, label_bold=True)
    # inherit branch elbows down + left into the final node's right edge
    add_arrow(slide, inh_x + 1.25, inh_y + nh,
              fx + nw, y8 + nh / 2,
              color=SIDE_BRANCH_LINE, weight=1.25,
              connector_type=MSO_CONNECTOR.ELBOW)

    # ----- right half: epic catalog YAML example -----
    code_x = 7.7
    code_y = 1.6
    code_w = 5.4
    code_h = 4.5
    yaml_text = """epic_catalog:
  catalog_id:     micro_epic_v1
  parent_catalog: cyto_epic_v1     # prior discipline

  inherited_epics:                 # score ≥ τ_match
    - id:    EPIC-MICRO-001
      title: Specimen Receiving
      epic_analog:
        catalog_id: cyto_epic_v1
        epic_id:    EPIC-CYTO-014
        equivalence: identical     # populated by construction
      auto_inherited: true
    - id:    EPIC-MICRO-002
      title: Reporting
      epic_analog: { epic_id: EPIC-CYTO-021, ... }

  novel_epics:                     # admitted via quorum
    - id:    EPIC-MICRO-007
      title: Antibiotic Susceptibility Testing (AST)
      cluster_evidence:
        n_drafts:  4
        from_sops: [SOP-007, SOP-019, SOP-022, SOP-031]
      votes: [admit×5]   decision: admit
      epic_analog: null              # no Cyto correspondent

  unclassified_drafts:             # E0 bucket
    count: 3"""
    add_code_box(slide, code_x, code_y, code_w, code_h, yaml_text,
                 font_size=8.5, title="Example output: epic catalog (simplified)")

    # ----- inheritance-lift mini-stat under YAML -----
    stat_y = code_y + code_h + 0.15
    add_text(slide, code_x, stat_y, code_w, 0.22,
             "What this looks like in practice",
             font_size=10, bold=True, color=ACCENT)
    stat_row_y = stat_y + 0.27
    stat_h = 0.42
    # 3-tile mini summary
    tiles = [
        ("inherited", "~70%", "auto, no quorum",
         BOX_FILL, BOX_LINE, ACCENT),
        ("novel via quorum", "~25%", "epic_analog = null",
         END_FILL, END_LINE, RGBColor(0x3E, 0x6A, 0x35)),
        ("E0 (parked)", "~5%", "discard quorum (D32)",
         PARK_FILL, PARK_LINE, PARK_LINE),
    ]
    tile_gap = 0.08
    tile_w = (code_w - 2 * tile_gap) / 3
    for i, (label, pct, sub, fill, line, color) in enumerate(tiles):
        x = code_x + i * (tile_w + tile_gap)
        # header
        add_box(slide, x, stat_row_y, tile_w, stat_h * 0.55, label,
                fill=fill, line=line, font_size=9, bold=True,
                font_color=color)
        # body — pct (large) + sub (small)
        add_text(slide, x, stat_row_y + stat_h * 0.55 + 0.02,
                 tile_w, 0.20, pct,
                 font_size=14, bold=True, color=color,
                 align=PP_ALIGN.CENTER)
        add_text(slide, x, stat_row_y + stat_h * 0.55 + 0.24,
                 tile_w, 0.18, sub,
                 font_size=8, italic=True, color=TEXT_GREY,
                 align=PP_ALIGN.CENTER)


# ---- slide 7: phase 2 — flowchart + story example ------------------------

def slide_six(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Phase 2 — Per-SOP processing")

    add_text(slide, 0.5, 1.05, 12.5, 0.4,
             "Runs once per SOP. Result: validated stories under the right epics. "
             "The dashed box below wraps everything that runs per SOP — N instances run in parallel.",
             font_size=12, color=TEXT_DARK)

    # ----- left half: flowchart -----
    fx = 0.5
    nw = 3.0
    nh = 0.55
    cx = fx + nw / 2
    dh = 0.75

    # ForEach-SOP container — wraps the per-SOP pipeline. Visual hint
    # that the entire pipeline runs once per SOP, with N instances in parallel.
    container = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       Inches(0.30), Inches(1.45),
                                       Inches(6.20), Inches(5.70))
    container.fill.background()
    container.line.color.rgb = RGBColor(0x90, 0x99, 0xA3)
    container.line.width = Pt(1.0)
    try:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        container.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    except Exception:
        pass
    # ForEach label hugging the container's top border
    label_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(0.45), Inches(1.36),
                                      Inches(3.5), Inches(0.22))
    label_bg.fill.solid()
    label_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    label_bg.line.fill.background()
    add_text(slide, 0.45, 1.36, 3.5, 0.22,
             "ForEach SOP   —   N instances run in parallel",
             font_size=10, italic=True, bold=True,
             color=RGBColor(0x60, 0x68, 0x70),
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # compressed spacing to fit 6 process steps + diamond + final
    nh_local = 0.50
    gap = 0.15
    dh_local = 0.60
    post_diamond = 0.35

    y1 = 1.55
    add_node(slide, fx, y1, nw, nh_local, "SOP arrives", kind="start",
             font_size=10, bold=True)

    # NEW: Section SOP into logical chunks (separate from tagging)
    y_chunk = y1 + nh_local + gap
    add_node(slide, fx, y_chunk, nw, nh_local,
             "Section SOP into logical chunks",
             kind="process", font_size=10)

    y2 = y_chunk + nh_local + gap
    add_node(slide, fx, y2, nw, nh_local,
             "Tag chunks: theme, test, persona, stage",
             kind="process", font_size=10)

    y3 = y2 + nh_local + gap
    add_node(slide, fx, y3, nw, nh_local,
             "Epic Extractor: draft + match",
             kind="agent", font_size=10, bold=True)

    y4 = y3 + nh_local + gap
    add_node(slide, fx, y4, nw, nh_local,
             "Story Extractor: draft\n(uses exemplars if available)",
             kind="agent", font_size=10, bold=True)

    y5 = y4 + nh_local + gap
    add_node(slide, fx, y5, nw, nh_local,
             "Validator: quality-check (shape rules)",
             kind="agent", font_size=10, bold=True)

    y6 = y5 + nh_local + gap
    add_diamond(slide, fx + 0.4, y6, nw - 0.8, dh_local,
                "Pass quality?", font_size=10)
    # fail branch (right)
    fail_x = fx + nw + 0.4
    fail_y = y6 + 0.05
    add_node(slide, fail_x, fail_y, 2.4, nh_local - 0.05,
             "Revise (≤2 tries)\nthen review pile",
             kind="park", font_size=10, bold=True)

    y7 = y6 + dh_local + post_diamond
    add_node(slide, fx, y7, nw, nh_local, "Story kept (validated)",
             kind="end", font_size=11, bold=True)

    # arrows
    add_arrow(slide, cx, y1 + nh_local, cx, y_chunk)
    add_arrow(slide, cx, y_chunk + nh_local, cx, y2)
    add_arrow(slide, cx, y2 + nh_local, cx, y3)
    add_arrow(slide, cx, y3 + nh_local, cx, y4)
    add_arrow(slide, cx, y4 + nh_local, cx, y5)
    add_arrow(slide, cx, y5 + nh_local, cx, y6)
    add_arrow(slide, fx + nw - 0.4, y6 + dh_local / 2, fail_x, fail_y + (nh_local - 0.05) / 2,
              label="no", label_color=PARK_LINE, label_bold=True,
              label_offset=(0.05, -0.18))
    add_arrow(slide, cx, y6 + dh_local, cx, y7,
              label="yes", label_color=END_LINE, label_bold=True,
              label_offset=(0.05, -0.18))

    # ----- right half: story schema example -----
    code_x = 7.7
    code_y = 1.55
    code_w = 5.4
    code_h = 3.9
    yaml_text = """story:
  id: STORY-MICRO-0042
  shape: capability       # one of:
                          #   capability,
                          #   workflow-stage-split,
                          #   configuration-instance,
                          #   cleanup
  title: "Validate received specimens
          against open orders"
  acceptance_criteria:
    - when: "specimen received with
             accession number"
      then: "system MUST match against
             open order"
  tests: [gram_stain, blood_culture_incubation]
  persona: lims           # closed enum
  stage: accessioning_verification
  source_citation: SOP-MICRO-014, lines 23-31
  epic_id: EPIC-MICRO-001 (Specimen Receiving)
  quality: passed         # or: parked"""
    add_code_box(slide, code_x, code_y, code_w, code_h, yaml_text,
                 font_size=8.5, title="Example output: a single story (simplified)")

    # ----- 4 shape rules (visual) — replaces prior explainer text -----
    rules_y = code_y + code_h + 0.10
    add_text(slide, code_x, rules_y, code_w, 0.25,
             "The 4 quality rules — one per story shape",
             font_size=11, bold=True, color=ACCENT)
    shape_items = [
        ("Capability",      "MUST/SHALL\n+ parameters"),
        ("Stage-split",     "Stage in title\n+ siblings"),
        ("Config-instance", "Concrete typed values\nno MUST/SHALL"),
        ("Cleanup",         "Named artifact\n+ before/after"),
    ]
    n = len(shape_items)
    gap = 0.06
    cards_y = rules_y + 0.30
    card_w = (code_w - (n - 1) * gap) / n
    head_h = 0.25
    body_h = 0.55
    for i, (head, body) in enumerate(shape_items):
        x = code_x + i * (card_w + gap)
        add_box(slide, x, cards_y, card_w, head_h, head,
                fill=SYS_FILL, line=SYS_LINE,
                font_size=10, bold=True, font_color=ACCENT)
        add_box(slide, x, cards_y + head_h, card_w, body_h, body,
                fill=RGBColor(0xFF, 0xFF, 0xFF), line=SYS_LINE,
                font_size=8)


# ---- slide 8: phase 3 — flowchart + capability story example ------------

def slide_seven(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Phase 3 — Combine patterns + continuous monitoring")

    add_text(slide, 0.5, 1.05, 12.5, 0.4,
             "Runs after a batch of SOPs. Recurring patterns lift to "
             "broader \"capability stories\" with parameters.",
             font_size=12, color=TEXT_DARK)

    # ----- left half: flowchart (compact for footer fit) -----
    fx = 0.5
    nw = 3.0
    nh = 0.45
    cx = fx + nw / 2
    dh = 0.6

    y1 = 1.55
    add_node(slide, fx, y1, nw, nh,
             "All kept stories\nfrom all SOPs",
             kind="start", font_size=9, bold=True)

    y2 = y1 + nh + 0.15
    add_node(slide, fx, y2, nw, nh,
             "Cluster on (test, persona, stage,\nshape, behavior)",
             kind="process", font_size=9)

    y3 = y2 + nh + 0.15
    add_diamond(slide, fx + 0.3, y3, nw - 0.6, dh,
                "Pattern in\n≥ 2 SOPs?", font_size=9)
    keep_x = fx + nw + 0.4
    keep_y = y3 + (dh - nh) / 2
    add_node(slide, keep_x, keep_y, 2.5, nh,
             "Keep concrete\nstory only",
             kind="side", font_size=9, bold=True)

    y4 = y3 + dh + 0.3
    add_node(slide, fx, y4, nw, nh,
             "Lift to capability story\n(parameterized)",
             kind="agent", font_size=9, bold=True)

    y5 = y4 + nh + 0.15
    add_node(slide, fx, y5, nw, nh,
             "Validator (gate 2)",
             kind="agent", font_size=9, bold=True)

    y6 = y5 + nh + 0.15
    add_node(slide, fx, y6, nw, nh,
             "Dependency Resolver\n(topological order)",
             kind="agent", font_size=9, bold=True)

    y7 = y6 + nh + 0.15
    add_node(slide, fx, y7, nw, nh,
             "Final ordered backlog",
             kind="end", font_size=10, bold=True)

    add_arrow(slide, cx, y1 + nh, cx, y2)
    add_arrow(slide, cx, y2 + nh, cx, y3)
    add_arrow(slide, fx + nw - 0.3, y3 + dh / 2, keep_x, keep_y + nh / 2,
              label="no", label_color=SIDE_BRANCH_LINE, label_bold=True,
              label_offset=(0.05, -0.18))
    add_arrow(slide, cx, y3 + dh, cx, y4,
              label="yes", label_color=END_LINE, label_bold=True,
              label_offset=(0.05, -0.18))
    add_arrow(slide, cx, y4 + nh, cx, y5)
    add_arrow(slide, cx, y5 + nh, cx, y6)
    add_arrow(slide, cx, y6 + nh, cx, y7)
    # keep-concrete branch converges into final node — elbow: down from
    # keep-concrete, then left into the right edge of the final backlog node
    add_arrow(slide, keep_x + 1.25, keep_y + nh,    # start: bottom-center of keep-concrete box
              fx + nw, y7 + nh / 2,                  # end: right-edge-middle of final node
              color=SIDE_BRANCH_LINE, weight=1.25,
              connector_type=MSO_CONNECTOR.ELBOW)

    # ----- right half: capability story example -----
    code_x = 7.7
    code_y = 1.6
    code_w = 5.4
    code_h = 4.0
    yaml_text = """capability_story:
  id: STORY-MICRO-CAP-0007
  shape: capability
  title: "Configurable critical-result
          notification across tests"

  parameters:                # what varies
    - test:      enum
    - threshold: number+units
    - persona_owner: microbiologist_supervisor

  acceptance_criteria:
    - when: "{test} flagged critical"
      then: "notify {persona_owner}
             within 1 hour"

  child_stories:             # the concrete
    - STORY-MICRO-0023       # cases
    - STORY-MICRO-0041       # this lifted
    - STORY-MICRO-0089       # from
  source_sops: [SOP-014, SOP-019, SOP-022]"""
    add_code_box(slide, code_x, code_y, code_w, code_h, yaml_text,
                 font_size=8.5, title="Example: a lifted capability story")

    # ----- right-column visual: what "pattern in ≥ 2 SOPs" produces -----
    # Sits below the capability-story YAML (right column only — left
    # column has the flowchart all the way down).
    viz_y = 5.7
    add_text(slide, code_x, viz_y, code_w, 0.25,
             "What \"pattern in ≥ 2 SOPs\" produces (worked example)",
             font_size=11, bold=True, color=ACCENT)

    # 2 source stories (left, stacked)
    src_x = code_x
    src_w = 1.30
    src_h = 0.35
    src_y_top = viz_y + 0.32
    src_y_bot = src_y_top + src_h + 0.05
    add_box(slide, src_x, src_y_top, src_w, src_h,
            "Story · SOP-007", fill=BOX_FILL, line=BOX_LINE,
            font_size=9, bold=True)
    add_box(slide, src_x, src_y_bot, src_w, src_h,
            "Story · SOP-014", fill=BOX_FILL, line=BOX_LINE,
            font_size=9, bold=True)

    # Cluster (middle)
    clu_x = src_x + src_w + 0.30
    clu_w = 1.55
    clu_y = src_y_top + 0.08
    clu_h = src_y_bot + src_h - clu_y - 0.08
    add_box(slide, clu_x, clu_y, clu_w, clu_h,
            "Cluster:\n2 distinct SOPs\n≥ 2 → lift",
            fill=SYS_FILL, line=SYS_LINE, font_size=9, bold=True)

    # Capability (right)
    cap_x = clu_x + clu_w + 0.30
    cap_w = code_x + code_w - cap_x
    cap_y = clu_y
    cap_h = clu_h
    add_box(slide, cap_x, cap_y, cap_w, cap_h,
            "Capability story\n+ parameters",
            fill=END_FILL, line=END_LINE, font_size=10, bold=True)

    # arrows
    add_arrow(slide, src_x + src_w + 0.02, src_y_top + src_h / 2,
              clu_x - 0.02, clu_y + clu_h * 0.35, weight=1.2)
    add_arrow(slide, src_x + src_w + 0.02, src_y_bot + src_h / 2,
              clu_x - 0.02, clu_y + clu_h * 0.65, weight=1.2)
    add_arrow(slide, clu_x + clu_w + 0.02, clu_y + clu_h / 2,
              cap_x - 0.02, cap_y + cap_h / 2, weight=1.5)


# ---- slide 8: what ships + alignment --------------------------------------

def slide_shape_examples(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "The 4 story shapes — one example each")

    add_text(slide, 0.5, 1.05, 12.5, 0.40,
             "Different shapes have different quality rules (slide 7). "
             "These examples show what makes each shape that shape.   "
             "(AC = Acceptance Criterion.)",
             font_size=12, italic=True, color=TEXT_GREY)

    grid_x = 0.4
    grid_y = 1.55
    h_gap = 0.18
    v_gap = 0.18
    card_w = (13.333 - 2 * grid_x - h_gap) / 2
    card_h = (7.05 - grid_y - v_gap) / 2
    head_h = 0.36

    examples = [
        ("CAPABILITY",
         "MUST/SHALL  +  configurable parameters",
         """title:  "Validate received specimens
        against open orders"
shape:   capability
persona: lims
stage:   accessioning_verification
tests:   [gram_stain, blood_culture_incubation]

AC1:  When specimen received with
      accession, system MUST match
      against open order
AC2:  When no match found within 24h,
      system MUST flag for review"""),
        ("WORKFLOW-STAGE-SPLIT",
         "Stage in title  +  sibling stories enumerated",
         """title:  "PHI (patient-health-info) update —
         after results are reported"
shape:   workflow-stage-split
persona:  microbiologist_supervisor
stage:   reporting_case_closure

AC1:  When PHI corrected on finalised
      report, system MUST regenerate
      with audit trail entry
AC2:  When the EHR (electronic health
      record) has the report, system MUST
      send corrected-result message

cross_links: STORY-0076 (before work),
             STORY-0077 (after work started)"""),
        ("CONFIGURATION-INSTANCE",
         "Concrete typed values  ·  no MUST/SHALL pretense",
         """title:  "Blood culture incubation
        — bottle parameters"
shape:   configuration-instance
tests:   [blood_culture_incubation]
persona: null     stage: null

AC1:  When bottle received, configure
      incubation
      expected_value:
        temperature_c:    35
        duration_h:       120
        agitation_rpm:    220
        alert_threshold_h: 24"""),
        ("CLEANUP",
         "Named artifact  +  before/after observable",
         """title:  "Remove obsolete
        'Pending Pathologist Review'
        status from results screen"
shape:   cleanup
persona: null     stage: null

AC1:  When viewing the results screen,
      the option MUST NOT appear
AC2:  When existing case has the status,
      MUST migrate to 'Pending Review'
      with audit trail

edge_cases:  existing case w/ status,
             in-flight orders w/ status"""),
    ]

    for i, (shape_name, feature, body_text) in enumerate(examples):
        col = i % 2
        row = i // 2
        x = grid_x + col * (card_w + h_gap)
        y = grid_y + row * (card_h + v_gap)
        # Header strip
        add_box(slide, x, y, card_w, head_h, shape_name,
                fill=SYS_FILL, line=SYS_LINE,
                font_size=13, bold=True, font_color=ACCENT)
        # Sub-feature
        add_text(slide, x + 0.05, y + head_h + 0.02, card_w - 0.10, 0.22,
                 feature,
                 font_size=10, italic=True, color=TEXT_GREY,
                 align=PP_ALIGN.CENTER)
        # Body code-style box
        body_y = y + head_h + 0.30
        body_h = card_h - head_h - 0.30
        add_code_box(slide, x, body_y, card_w, body_h, body_text,
                     font_size=8.5)


def slide_inside_agents(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Inside two key agents — Validator & Dependency Resolver")

    add_text(slide, 0.5, 1.05, 12.5, 0.4,
             "What these two agents actually do, beyond their one-line description on slide 3.",
             font_size=12, italic=True, color=TEXT_GREY)

    # ===== LEFT panel: Validator =====
    left_x = 0.4
    left_w = 6.30
    add_box(slide, left_x, 1.55, left_w, 0.5, "Validator",
            fill=SYS_FILL, line=SYS_LINE, font_size=15, bold=True,
            font_color=ACCENT)
    add_text(slide, left_x, 2.10, left_w, 0.3,
             "Runs at gate 1 (after Story Extractor) and gate 2 (after Synthesis).",
             font_size=10, italic=True, color=TEXT_GREY,
             align=PP_ALIGN.CENTER)

    fx = left_x + 0.15
    nw = left_w - 0.30
    nh = 0.42
    cx = fx + nw / 2

    y = 2.50
    steps = [
        ("start",   "Receives a draft story",                                 True),
        ("process", "1. Closed-enum check  (tests, persona, stage in catalog?)", False),
        ("process", "2. Shape verification  (declared shape matches content?)", False),
        ("process", "3. Apply shape-specific rubric  (4 rules — see slide 7)", False),
    ]
    for kind, text, bold in steps:
        add_node(slide, fx, y, nw, nh, text, kind=kind, font_size=10, bold=bold)
        if y > 2.50:
            add_arrow(slide, cx, y - 0.13, cx, y, weight=1.3)
        y += nh + 0.13

    # decision diamond — narrowed so it doesn't overlap the right branch
    dh = 0.55
    diamond_x = fx + 1.95
    diamond_w = nw - 4.10
    dcx = diamond_x + diamond_w / 2  # diamond's own midline (centers arrows on its vertices)
    # arrow from last process box → diamond top vertex (using dcx, not panel cx)
    add_arrow(slide, dcx, y - 0.13, dcx, y, weight=1.3)
    add_diamond(slide, diamond_x, y, diamond_w, dh,
                "Pass all checks?", font_size=10)

    # right branch: revise/park — moved further right with clear gap
    rev_x = diamond_x + diamond_w + 0.30
    rev_y = y + 0.05
    rev_w = (fx + nw) - rev_x
    add_node(slide, rev_x, rev_y, rev_w, nh - 0.05,
             "revise (≤2)\nthen auto-park",
             kind="park", font_size=8, bold=True)
    # arrow from diamond right vertex → left edge of revise box
    add_arrow(slide, diamond_x + diamond_w, y + dh / 2,
              rev_x, rev_y + (nh - 0.05) / 2,
              label="no", label_color=PARK_LINE, label_bold=True,
              label_offset=(0.0, -0.22))

    # arrow down from diamond bottom vertex (using dcx, not panel cx)
    y2 = y + dh + 0.30
    add_arrow(slide, dcx, y + dh, dcx, y2,
              label="yes", label_color=END_LINE, label_bold=True,
              label_offset=(0.05, -0.18))
    # final node
    add_node(slide, fx, y2, nw, nh, "Story kept (validated) → Output A",
             kind="end", font_size=10, bold=True)

    # ===== RIGHT panel: Dependency Resolver =====
    right_x = 6.95
    right_w = 6.00
    add_box(slide, right_x, 1.55, right_w, 0.5, "Dependency Resolver",
            fill=SYS_FILL, line=SYS_LINE, font_size=15, bold=True,
            font_color=ACCENT)
    add_text(slide, right_x, 2.10, right_w, 0.3,
             "Takes all kept stories; outputs a topological order for sprint planning.",
             font_size=10, italic=True, color=TEXT_GREY,
             align=PP_ALIGN.CENTER)
    add_text(slide, right_x, 2.42, right_w, 0.3,
             "Each story has dependencies[] and cross_links[] — the resolver builds a dependency graph (DAG) and sorts it.",
             font_size=9, italic=True, color=TEXT_GREY,
             align=PP_ALIGN.CENTER)

    # DAG: A → B and A → C; B and C → D; D → E
    s_w, s_h = 1.05, 0.40
    panel_cx = right_x + right_w / 2

    # A at top center
    a_x = panel_cx - s_w / 2
    a_y = 2.85
    add_box(slide, a_x, a_y, s_w, s_h, "Story A",
            fill=BOX_FILL, line=BOX_LINE, font_size=10, bold=True)

    # B (left) and C (right) in middle row
    bc_y = a_y + s_h + 0.40
    b_x = right_x + 0.5
    c_x = right_x + right_w - s_w - 0.5
    add_box(slide, b_x, bc_y, s_w, s_h, "Story B",
            fill=BOX_FILL, line=BOX_LINE, font_size=10, bold=True)
    add_box(slide, c_x, bc_y, s_w, s_h, "Story C",
            fill=BOX_FILL, line=BOX_LINE, font_size=10, bold=True)

    # D below center
    d_x = panel_cx - s_w / 2
    d_y = bc_y + s_h + 0.40
    add_box(slide, d_x, d_y, s_w, s_h, "Story D",
            fill=BOX_FILL, line=BOX_LINE, font_size=10, bold=True)

    # E below D
    e_x = d_x
    e_y = d_y + s_h + 0.40
    add_box(slide, e_x, e_y, s_w, s_h, "Story E",
            fill=BOX_FILL, line=BOX_LINE, font_size=10, bold=True)

    # Dependency arrows (A→B, A→C, B→D, C→D, D→E)
    add_arrow(slide, a_x + s_w/2, a_y + s_h, b_x + s_w/2, bc_y, weight=1.2)
    add_arrow(slide, a_x + s_w/2, a_y + s_h, c_x + s_w/2, bc_y, weight=1.2)
    add_arrow(slide, b_x + s_w/2, bc_y + s_h, d_x + s_w/2, d_y, weight=1.2)
    add_arrow(slide, c_x + s_w/2, bc_y + s_h, d_x + s_w/2, d_y, weight=1.2)
    add_arrow(slide, d_x + s_w/2, d_y + s_h, e_x + s_w/2, e_y, weight=1.2)

    # Output strip below the DAG
    out_y = e_y + s_h + 0.20
    add_box(slide, right_x, out_y, right_w, 0.45,
            "Output:   1. A  →  2. B  →  3. C  →  4. D  →  5. E",
            fill=END_FILL, line=END_LINE, font_size=11, bold=True,
            font_color=ACCENT)


def slide_validator_walkthrough(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Validator at work — a sample story walks through the checks")

    add_text(slide, 0.5, 1.05, 12.5, 0.40,
             "Each check is explicit. Failure on attempt 1; pass on revise (attempt 2).",
             font_size=12, italic=True, color=TEXT_GREY)

    # 2 cards side by side: ATTEMPT 1 (left, fail) and ATTEMPT 2 (right, pass)
    grid_y = 1.55
    h_gap = 0.20
    card_w = (13.333 - 0.8 - h_gap) / 2

    # ===== LEFT: Attempt 1 (fail) =====
    a1_x = 0.4
    add_box(slide, a1_x, grid_y, card_w, 0.40,
            "ATTEMPT 1   ✗   draft",
            fill=PARK_FILL, line=PARK_LINE,
            font_size=14, bold=True, font_color=PARK_LINE)

    a1_yaml = """story (draft):
  shape:   capability
  title:   "Validate accession appropriately"
  persona: lims
  tests:   [gram_stain]
  stage:   accessioning_verification
  AC: "When specimen received, system MUST
       validate appropriately"
  source:  SOP-MICRO-014, lines 23-31"""
    a1_code_y = grid_y + 0.45
    add_code_box(slide, a1_x, a1_code_y, card_w, 1.75, a1_yaml,
                 font_size=8.5)

    checks_y = a1_code_y + 1.85
    add_text(slide, a1_x + 0.05, checks_y, card_w - 0.10, 0.25,
             "Validator checks:",
             font_size=11, bold=True, color=ACCENT)

    a1_checks = [
        ("✓", "1. Closed-enum check  —  tests, persona, stage all in catalog", END_LINE, False),
        ("✓", "2. Shape verification  —  declared 'capability' matches content", END_LINE, False),
        ("✗", "3. Shape rubric  —  capability:no_ambiguous_quantifiers fails on 'appropriately'", PARK_LINE, True),
    ]
    cy = checks_y + 0.32
    for mark, text, color, bold in a1_checks:
        add_text(slide, a1_x + 0.10, cy, card_w - 0.20, 0.30,
                 f"{mark}    {text}",
                 font_size=10, color=color, bold=bold)
        cy += 0.32

    verdict_y = cy + 0.10
    add_box(slide, a1_x, verdict_y, card_w, 0.40,
            "→ Revise  (attempt 1 of 2)  —  Story Extractor reworks the AC",
            fill=PARK_FILL, line=PARK_LINE,
            font_size=11, bold=True, font_color=PARK_LINE)

    # ===== RIGHT: Attempt 2 (pass) =====
    a2_x = a1_x + card_w + h_gap
    add_box(slide, a2_x, grid_y, card_w, 0.40,
            "ATTEMPT 2   ✓   revised",
            fill=END_FILL, line=END_LINE,
            font_size=14, bold=True, font_color=END_LINE)

    a2_yaml = """story (revised):
  shape:   capability
  title:   "Validate accession against open order"
  persona: lims
  tests:   [gram_stain]
  stage:   accessioning_verification
  AC: "When specimen received with accession,
       system MUST match against open order;
       if no match within 24h, MUST flag for
       accessioning review"
  source:  SOP-MICRO-014, lines 23-31"""
    add_code_box(slide, a2_x, a1_code_y, card_w, 1.75, a2_yaml,
                 font_size=8.5)

    add_text(slide, a2_x + 0.05, checks_y, card_w - 0.10, 0.25,
             "Validator checks:",
             font_size=11, bold=True, color=ACCENT)
    a2_checks = [
        ("✓", "1. Closed-enum check  —  passes (unchanged)", END_LINE, False),
        ("✓", "2. Shape verification  —  passes", END_LINE, False),
        ("✓", "3. Shape rubric  —  concrete, observable, parameterized (24h)", END_LINE, True),
    ]
    cy = checks_y + 0.32
    for mark, text, color, bold in a2_checks:
        add_text(slide, a2_x + 0.10, cy, card_w - 0.20, 0.30,
                 f"{mark}    {text}",
                 font_size=10, color=color, bold=bold)
        cy += 0.32

    add_box(slide, a2_x, verdict_y, card_w, 0.40,
            "→ Story kept (validated)  →  Output A",
            fill=END_FILL, line=END_LINE,
            font_size=11, bold=True, font_color=END_LINE)

    # Bottom note (full-width)
    note_y = verdict_y + 0.55
    add_text(slide, 0.5, note_y, 12.5, 0.50,
             "Note:  closed-enum violations (e.g. invalid persona name) hard-reject without entering the revise loop. "
             "Only shape/AC issues consume revision attempts. Failures after 2 revisions auto-park to the review pile (see slide 12).",
             font_size=10, italic=True, color=TEXT_GREY,
             align=PP_ALIGN.CENTER)


def slide_failures(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "What happens when something fails")

    add_text(slide, 0.5, 1.05, 12.5, 0.5,
             "Failures at three levels — each has a deterministic fallback. "
             "No human intervention is required at runtime.",
             font_size=12, color=TEXT_DARK)

    # 3 panels horizontally
    panel_y = 1.7
    panel_h = 3.7
    panel_w = 4.05
    gap = 0.18
    panels_total = 3 * panel_w + 2 * gap
    px0 = (13.333 - panels_total) / 2

    panels = [
        {
            "header": "Per-call",
            "subhead": "(LLM error or timeout)",
            "fill": SYS_FILL,
            "line": SYS_LINE,
            "rows": [
                ("Trigger", "An agent's LLM call returns an error or times out."),
                ("Fallback", "Retry with exponential backoff (3 attempts)."),
                ("If persistent", "Defer the item to the next batch; log the failure to the audit trail."),
                ("Pipeline impact", "Other items continue. No block."),
            ],
        },
        {
            "header": "Per-story",
            "subhead": "(quality check fails)",
            "fill": SIDE_BRANCH_FILL,
            "line": SIDE_BRANCH_LINE,
            "rows": [
                ("Trigger", "Validator rejects a story."),
                ("Fallback", "Up to 2 revision attempts; the Validator drives revisions with the failed-checks list."),
                ("If persistent", "Auto-park to the review pile. Story does not reach Output A in strict mode."),
                ("Pipeline impact", "Subsequent stories continue. No block."),
            ],
        },
        {
            "header": "Per-batch",
            "subhead": "(drift / catalog issues)",
            "fill": PARK_FILL,
            "line": PARK_LINE,
            "rows": [
                ("Trigger", "G0 / E0 bucket > 5% rolling, or auto-park rate > 15% sustained."),
                ("Fallback", "Auto-rerun discovery with τ_match −0.05 (more permissive); re-cluster G0; quorum re-evaluates."),
                ("If τ_match floor (0.50) hit", "Drift report flags the discipline for manual catalog re-curation (one-time architect action)."),
                ("Pipeline impact", "Pipeline continues. The drift report is informational."),
            ],
        },
    ]

    for i, p in enumerate(panels):
        x = px0 + i * (panel_w + gap)
        # header strip
        add_box(slide, x, panel_y, panel_w, 0.5, p["header"],
                fill=p["fill"], line=p["line"],
                font_size=14, bold=True, font_color=ACCENT)
        add_text(slide, x, panel_y + 0.55, panel_w, 0.25, p["subhead"],
                 font_size=10, italic=True, color=TEXT_GREY,
                 align=PP_ALIGN.CENTER)

        # body box (white background, panel-color border)
        body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      Inches(x), Inches(panel_y + 0.85),
                                      Inches(panel_w), Inches(panel_h - 0.85))
        body.fill.solid()
        body.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        body.line.color.rgb = p["line"]
        body.line.width = Pt(0.75)
        tf = body.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.12)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.10)
        tf.margin_bottom = Inches(0.10)
        tf.vertical_anchor = MSO_ANCHOR.TOP
        for j, (label, text) in enumerate(p["rows"]):
            if j == 0:
                pp = tf.paragraphs[0]
            else:
                pp = tf.add_paragraph()
            pp.alignment = PP_ALIGN.LEFT
            r1 = pp.add_run()
            r1.text = label + ": "
            r1.font.size = Pt(10)
            r1.font.bold = True
            r1.font.color.rgb = ACCENT
            r1.font.name = "Calibri"
            r2 = pp.add_run()
            r2.text = text
            r2.font.size = Pt(10)
            r2.font.color.rgb = TEXT_DARK
            r2.font.name = "Calibri"
            pp.space_after = Pt(6)

    # Bottom banner — three queues
    banner_y = panel_y + panel_h + 0.20
    add_text(slide, 0.5, banner_y, 12.5, 0.35,
             "Three queues collect different kinds of failures",
             font_size=14, bold=True, color=ACCENT)
    queue_y = banner_y + 0.45
    queue_h = 7.05 - queue_y
    queue_w = 4.05
    qx0 = (13.333 - 3 * queue_w - 2 * gap) / 2
    queues = [
        ("Review pile", "Quality / rubric failures", "Stories that failed Validator after 2 revisions."),
        ("G0 / E0 buckets", "Classification failures", "Chunks (G0) and draft epics (E0) that didn't classify or cluster."),
        ("Audit trail", "Call-level errors", "LLM-call failures after retry exhaustion. For diagnostics."),
    ]
    for i, (name, sub, desc) in enumerate(queues):
        x = qx0 + i * (queue_w + gap)
        bb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(x), Inches(queue_y),
                                    Inches(queue_w), Inches(queue_h))
        bb.fill.solid()
        bb.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xFA)
        bb.line.color.rgb = BOX_LINE
        bb.line.width = Pt(0.75)
        tf = bb.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.12)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.08)
        tf.margin_bottom = Inches(0.08)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.LEFT
        r1 = p1.add_run()
        r1.text = name
        r1.font.size = Pt(12)
        r1.font.bold = True
        r1.font.color.rgb = ACCENT
        r1.font.name = "Calibri"
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = sub
        r2.font.size = Pt(10)
        r2.font.italic = True
        r2.font.color.rgb = TEXT_GREY
        r2.font.name = "Calibri"
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.LEFT
        r3 = p3.add_run()
        r3.text = desc
        r3.font.size = Pt(10)
        r3.font.color.rgb = TEXT_DARK
        r3.font.name = "Calibri"


def slide_alignment_ask(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "What ships and the alignment ask")

    deliv_y = 1.3
    deliv_w = 4.0
    deliv_h = 1.5
    gap = 0.3
    total = 3 * deliv_w + 2 * gap
    start_x = (13.333 - total) / 2

    deliveries = [
        ("Jira backlog",
         "Epics + Stories\nReady for the dev team\nIncludes traceability to source SOPs"),
        ("Configuration files",
         "Per-test settings\n(e.g. blood culture parameters)\nDeterministic structure"),
        ("Reference + audit",
         "Theme catalog, epic catalog\nCross-discipline mapping\nReview pile + decision log"),
    ]
    for i, (head, body) in enumerate(deliveries):
        x = start_x + i * (deliv_w + gap)
        add_box(slide, x, deliv_y, deliv_w, 0.55, head,
                fill=OUT_FILL, line=OUT_LINE,
                font_size=15, bold=True)
        bb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(x), Inches(deliv_y + 0.55),
                                    Inches(deliv_w), Inches(deliv_h - 0.55))
        bb.fill.solid()
        bb.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        bb.line.color.rgb = OUT_LINE
        bb.line.width = Pt(1.0)
        _set_text(bb, body, font_size=11, bold=False, color=TEXT_DARK,
                  align=PP_ALIGN.CENTER)

    # Inputs section — categorized by sourcing
    add_text(slide, 0.5, 2.95, 12.5, 0.35,
             "Inputs to the system  —  what's already captured vs what's still needed",
             font_size=14, bold=True, color=ACCENT)

    # 3 columns
    col_y = 3.35
    col_h = 2.10
    col_gap = 0.18
    col_w = (13.333 - 2 * 0.4 - 2 * col_gap) / 3
    col_x0 = 0.4

    # Column A — Already captured from project documentation
    a_x = col_x0
    add_box(slide, a_x, col_y, col_w, 0.40,
            "Already captured  ✓",
            fill=END_FILL, line=END_LINE,
            font_size=12, bold=True, font_color=END_LINE)
    add_text(slide, a_x + 0.03, col_y + 0.42, col_w - 0.06, 0.22,
             "from your project documentation",
             font_size=9, italic=True, color=TEXT_GREY,
             align=PP_ALIGN.CENTER)
    add_bullets(slide, a_x + 0.10, col_y + 0.72, col_w - 0.20, col_h - 0.72, [
        "Persona inventory  —  documented in your Roles & Process Map (we'll structure it into micro_persona_v1.yaml)",
        "Universal 6-stage enum  —  derived from the Process Map (one-time, universal)",
        "ANALOGY persona_links  —  from the cross-discipline role inventory",
    ], font_size=9, color=TEXT_DARK)

    # Column B — From the seminal prior (Cytology)
    b_x = col_x0 + col_w + col_gap
    add_box(slide, b_x, col_y, col_w, 0.40,
            "From the seminal prior  ✓",
            fill=BOX_FILL, line=BOX_LINE,
            font_size=12, bold=True, font_color=ACCENT)
    add_text(slide, b_x + 0.03, col_y + 0.42, col_w - 0.06, 0.22,
             "snapshotted from Cytology's Connect build",
             font_size=9, italic=True, color=TEXT_GREY,
             align=PP_ALIGN.CENTER)
    add_bullets(slide, b_x + 0.10, col_y + 0.72, col_w - 0.20, col_h - 0.72, [
        "Theme catalog  (cyto_v1)  —  used as the warm-start prior for Theme Discovery",
        "Epic catalog  (cyto_epic_v1)  —  used by the conditioned Epic Extractor",
    ], font_size=9, color=TEXT_DARK)

    # Column C — Still needed from the client
    c_x = col_x0 + 2 * (col_w + col_gap)
    add_box(slide, c_x, col_y, col_w, 0.40,
            "Still needed from the client",
            fill=SYS_FILL, line=SYS_LINE,
            font_size=12, bold=True, font_color=ACCENT)
    add_text(slide, c_x + 0.03, col_y + 0.42, col_w - 0.06, 0.22,
             "the only blocking inputs",
             font_size=9, italic=True, color=TEXT_GREY,
             align=PP_ALIGN.CENTER)
    add_bullets(slide, c_x + 0.10, col_y + 0.72, col_w - 0.20, col_h - 0.72, [
        "Test catalog  (micro_test_v1.yaml)  —  list of in-scope Microbiology tests (Gram stain, blood culture, MALDI mass-spec ID, AST antimicrobial susceptibility testing, …)",
        "Sample SOPs  (~20 representative SOPs) to seed Phase 1",
    ], font_size=9, color=TEXT_DARK)

    # Optional row
    opt_y = col_y + col_h + 0.10
    add_text(slide, 0.5, opt_y, 12.5, 0.25,
             "Optional  —  quality enhancement",
             font_size=11, bold=True, italic=True, color=TEXT_GREY)
    add_text(slide, 0.7, opt_y + 0.27, 12.0, 0.25,
             "•   Sanitized Cytology Jira export  —  improves story style fidelity at extraction time. The system runs without it.",
             font_size=10, italic=True, color=TEXT_GREY)

    # Alignment ask
    add_box(slide, 0.5, opt_y + 0.65, 12.55, 0.55,
            "Alignment ask:  the test catalog + sample SOPs are the gating inputs.  Everything else is in hand.  Confirm?",
            fill=SYS_FILL, line=SYS_LINE,
            font_size=13, bold=True, font_color=ACCENT,
            align=PP_ALIGN.CENTER)


# ---- main -----------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_one(prs)
    slide_two(prs)
    slide_three(prs)
    slide_four(prs)
    slide_five(prs)
    slide_epic_extractor(prs)
    slide_six(prs)
    slide_seven(prs)
    slide_shape_examples(prs)
    slide_inside_agents(prs)
    slide_validator_walkthrough(prs)
    slide_failures(prs)
    slide_alignment_ask(prs)

    # add footers with slide numbers
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, start=1):
        add_footer(slide, i, total)

    out = Path(__file__).parent / "client_alignment_deck.pptx"
    prs.save(str(out))
    print(f"wrote: {out}  ({total} slides)")


if __name__ == "__main__":
    main()
