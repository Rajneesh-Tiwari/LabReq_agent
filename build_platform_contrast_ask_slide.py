"""
Single exploratory slide: ask to the Connect platform team for the source
artifacts we can use to assemble a picture of Connect's platform capabilities
and common functionality — so we can tag each discovered item with
platform_status (existing / extension / new / divergent).

Honest framing: they won't have our 4-level shape ready-made. We assemble it
from what they have, and they either provide the mapping or QC ours.

Standalone — not yet wired into client_alignment_deck_v2.pptx.
Reuses helpers + palette from build_client_deck_v2.py so the look matches.

Output: platform_contrast_ask.pptx (one slide, 16:9, Calibri).
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from build_client_deck_v2 import (
    add_box, add_text, add_title, add_subtitle, add_footer,
    BOX_FILL, BOX_LINE,
    SYS_FILL, SYS_LINE,
    OUT_FILL, OUT_LINE,
    PHASE_FILL, PHASE_LINE,
    CARD_FILL, CARD_LINE,
    ACCENT, TEXT_DARK, TEXT_GREY,
)


# ---- local helpers ---------------------------------------------------------

def panel_bg(slide, x, y, w, h, fill, line):
    return add_box(slide, x, y, w, h, "", fill=fill, line=line)


def panel_header(slide, x, y, w, text, color):
    add_text(slide, x + 0.20, y + 0.10, w - 0.40, 0.34, text,
             font_size=13, bold=True, color=color)


def rich_bullet(slide, x, y, w, h, lead, tail, body_size=10.5,
                lead_color=TEXT_DARK, tail_color=TEXT_DARK,
                bullet_glyph="•  ", lead_bold=True):
    """A single bullet with a bold-lead phrase and regular-tail explanation."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.0)
    tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = bullet_glyph + lead
    r1.font.size = Pt(body_size)
    r1.font.bold = lead_bold
    r1.font.color.rgb = lead_color
    r1.font.name = "Calibri"
    if tail:
        r2 = p.add_run()
        r2.text = " — " + tail
        r2.font.size = Pt(body_size)
        r2.font.color.rgb = tail_color
        r2.font.name = "Calibri"
    return tb


def plain_text(slide, x, y, w, h, text, size=10.5, bold=False,
               italic=False, color=TEXT_DARK, align=PP_ALIGN.LEFT):
    add_text(slide, x, y, w, h, text,
             font_size=size, bold=bold, italic=italic, color=color, align=align)


# ---- the slide -------------------------------------------------------------

def slide_platform_contrast_ask(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_title(slide, "Contrasting discoveries against the Connect platform")
    add_subtitle(
        slide,
        "We need a picture of Connect's platform capabilities and common "
        "functionality, so we can tag each discovered item — capability, "
        "theme, epic, story — with whether it already exists on Connect, "
        "extends an existing feature, or is net-new. We don't expect a "
        "ready-made catalog; we'll assemble one from artifacts you already "
        "have. The mapping itself is then either provided or QC'd by your team.",
        y=1.02, size=11.5,
    )

    # ----- LEFT PANEL — Sources they probably already have (BLUE) -----
    lx, ly, lw, lh = 0.50, 1.78, 5.30, 4.32
    panel_bg(slide, lx, ly, lw, lh, BOX_FILL, BOX_LINE)
    panel_header(slide, lx, ly, lw,
                 "Sources of platform capabilities — what you likely have",
                 BOX_LINE)

    sources = [
        ("Confluence / architecture docs",
         "module / subsystem pages describing what each part of Connect does.",
         0.50),
        ("Admin / configuration reference  ★",
         "every config knob, organised by module. Highest leverage — this is "
         "how we tell 'extension' from 'existing'.",
         0.62),
        ("Release notes per Connect release",
         "features added, changed, deprecated. Time-ordered. Useful for "
         "status, introduced_in, deprecated_in.",
         0.55),
        ("API / integration documentation",
         "external-system contracts, message schemas.",
         0.42),
        ("Training / onboarding decks, sales collateral",
         "often the cleanest 'what does Connect do' summary, plus feature "
         "shortlists. Useful as a starting outline.",
         0.55),
        ("Per-discipline Jira backlogs (Cyto, etc.)",
         "describes what was built for one discipline — not platform-wide "
         "capability. Use as cross-check only, not primary source.",
         0.62),
    ]

    by = ly + 0.50
    bw = lw - 0.40
    for lead, tail, h in sources:
        rich_bullet(slide, lx + 0.20, by, bw, h + 0.05, lead, tail,
                    body_size=10, lead_color=ACCENT)
        by += h

    # ----- RIGHT TOP PANEL — Our categorisation + mapping (YELLOW) -----
    rx, ry, rw, rh_top = 5.93, 1.78, 6.90, 2.65
    panel_bg(slide, rx, ry, rw, rh_top, SYS_FILL, SYS_LINE)
    panel_header(slide, rx, ry, rw,
                 "Our 4-level categorisation — how we'd map your sources",
                 SYS_LINE)
    plain_text(slide, rx + 0.20, ry + 0.46, rw - 0.40, 0.30,
               "Each row below is a 'platform capability' or 'common "
               "functionality' we'll match against every discovered item.",
               size=10, italic=True, color=TEXT_GREY)

    # mapping rows
    map_rows = [
        ("business capability",
         "module pages, sales collateral, training decks"),
        ("theme",
         "sub-module pages, data-model concepts, training decks"),
        ("epic",
         "behavioural sections in docs, release notes, API docs"),
        ("configuration  ★",
         "admin / configuration reference"),
    ]
    my = ry + 0.86
    label_w, sep, arrow_w = 2.20, 0.10, 0.40
    src_x = rx + 0.20 + label_w + sep + arrow_w + sep
    src_w = rw - (src_x - rx) - 0.20
    for level, sources_str in map_rows:
        # level label (bold, accent colour)
        plain_text(slide, rx + 0.20, my, label_w, 0.30, level,
                   size=11, bold=True, color=ACCENT)
        # arrow
        plain_text(slide, rx + 0.20 + label_w + sep, my, arrow_w, 0.30,
                   "←", size=12, bold=True, color=SYS_LINE,
                   align=PP_ALIGN.CENTER)
        # sources
        plain_text(slide, src_x, my, src_w, 0.30, sources_str,
                   size=10.5, color=TEXT_DARK)
        my += 0.38

    # ----- RIGHT BOTTOM PANEL — Two paths to verify mapping (PURPLE) -----
    rh_bot = 1.55
    rby = ry + rh_top + 0.10  # 1.78 + 2.65 + 0.10 = 4.53
    panel_bg(slide, rx, rby, rw, rh_bot, PHASE_FILL, PHASE_LINE)
    panel_header(slide, rx, rby, rw,
                 "How we ratify the mapping — pick one path",
                 PHASE_LINE)

    rich_bullet(
        slide, rx + 0.20, rby + 0.48, rw - 0.40, 0.40,
        "① You provide the mapping",
        "your team annotates artifacts by level. Cleanest if you already "
        "maintain a structured feature catalog. Rare in practice.",
        body_size=10.5, lead_color=PHASE_LINE,
    )
    rich_bullet(
        slide, rx + 0.20, rby + 0.92, rw - 0.40, 0.40,
        "② You QC our draft  (expected default)",
        "we draft the mapping from your sources; you SME-review; 1–2 review "
        "cycles to ratify. Mapping accuracy signed off by your team.",
        body_size=10.5, lead_color=PHASE_LINE,
    )

    # ----- BOTTOM STRIP — What we deliver back (GREEN) -----
    sy, sh = 6.20, 0.92
    panel_bg(slide, 0.50, sy, 12.33, sh, OUT_FILL, OUT_LINE)
    plain_text(slide, 0.70, sy + 0.10, 12.0, 0.30,
               "What we deliver back, per discipline",
               size=12.5, bold=True, color=OUT_LINE)
    plain_text(
        slide, 0.70, sy + 0.40, 12.0, 0.48,
        "Platform Delta Report — every discovered capability / theme / epic "
        "/ story tagged with platform_status ∈ {existing · extension · new · "
        "divergent}, with the matched platform item where applicable. "
        "Refreshed each Connect platform release.",
        size=10.5, color=TEXT_DARK,
    )


# ---- main ------------------------------------------------------------------

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_platform_contrast_ask(prs)
    add_footer(prs.slides[0], 1, 1)

    out = Path(__file__).parent / "platform_contrast_ask.pptx"
    prs.save(str(out))
    print(f"wrote: {out}  (1 slide)")

    # bounds check — every content shape must end above y=7.18.
    # footer shapes themselves sit at 7.18 / 7.22; those are exempt.
    bad = []
    for shp in prs.slides[0].shapes:
        top = shp.top / 914400
        bot = top + shp.height / 914400
        if bot > 7.18 + 0.001 and not (abs(top - 7.18) < 0.01 or
                                        abs(top - 7.22) < 0.01):
            bad.append((shp.shape_type, top, bot))
    if bad:
        print("  WARN bounds exceeded:")
        for s, t, b in bad:
            print(f"    type={s} top={t:.3f} bottom={b:.3f}")
    else:
        print("  bounds OK (all content above y=7.18)")


if __name__ == "__main__":
    main()
