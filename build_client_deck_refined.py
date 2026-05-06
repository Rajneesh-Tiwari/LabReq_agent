"""
Generate the REFINED client alignment deck — editorial aesthetic with one
bold accent. Cambria + Calibri + Consolas. Warm cream background.

Usage:
    python3 build_client_deck_refined.py

Produces: client_alignment_deck_refined.pptx alongside this script.

The original `build_client_deck.py` is preserved separately. This file is
a from-scratch redesign:
  - Warm cream background (#FAF7F0) instead of plain white
  - Single bold vermillion accent (#C0441F) instead of rainbow pastels
  - Cambria display + Calibri body + Consolas mono (no generic Inter/Arial)
  - Asymmetric titles, numbered eyebrows, hairline dividers
  - Editorial layouts with generous whitespace

12 slides, same content as the original deck.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ============================================================================
# DESIGN TOKENS
# ============================================================================

# Colors
BG_CREAM      = RGBColor(0xFA, 0xF7, 0xF0)
BG_ALT        = RGBColor(0xF1, 0xEC, 0xE3)
INK           = RGBColor(0x1A, 0x1A, 0x2E)
INK_SOFT      = RGBColor(0x5A, 0x68, 0x78)
ACCENT        = RGBColor(0xC0, 0x44, 0x1F)
ACCENT_SOFT   = RGBColor(0xE8, 0xC8, 0xA0)
HAIRLINE      = RGBColor(0xD4, 0xCC, 0xBE)
SUCCESS       = RGBColor(0x2D, 0x5C, 0x3D)
FAILURE       = RGBColor(0x8B, 0x2E, 0x2A)
QUIET         = RGBColor(0xB0, 0xA8, 0x9B)
HIGHLIGHT_BG  = RGBColor(0xF7, 0xE8, 0xC8)
CODE_BG       = RGBColor(0xF1, 0xEC, 0xDD)
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)

# Fonts
F_DISPLAY     = "Cambria"
F_BODY        = "Calibri Light"
F_BODY_REG    = "Calibri"
F_MONO        = "Consolas"

# Layout constants
SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN_L = 0.55
MARGIN_R = 0.55
FOOTER_Y = 7.10        # everything must end above this


# ============================================================================
# TEXT / SHAPE HELPERS
# ============================================================================

def _set_text(shape, text, *, font_size, font, bold=False, italic=False,
              color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
              line_spacing=None, char_spacing=None):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    if char_spacing is not None:
        # spc is hundredths of a point (negative tightens, positive loosens)
        run.font._rPr.set("spc", str(int(char_spacing * 100)))


def add_text(slide, left, top, width, height, text, *, font_size=12,
             font=F_BODY, bold=False, italic=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=None,
             char_spacing=None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    _set_text(tb, text, font_size=font_size, font=font, bold=bold,
              italic=italic, color=color, align=align, anchor=anchor,
              line_spacing=line_spacing, char_spacing=char_spacing)
    return tb


def add_multitext(slide, left, top, width, height, runs, *, align=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP, line_spacing=None):
    """runs: list of dicts {text, size, font, bold, italic, color, char_spacing}."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    for r in runs:
        run = p.add_run()
        run.text = r["text"]
        run.font.size = Pt(r.get("size", 12))
        run.font.bold = r.get("bold", False)
        run.font.italic = r.get("italic", False)
        run.font.color.rgb = r.get("color", INK)
        run.font.name = r.get("font", F_BODY)
        if "char_spacing" in r:
            run.font._rPr.set("spc", str(int(r["char_spacing"] * 100)))
    return tb


def add_hairline(slide, left, top, width, color=HAIRLINE, weight=0.75):
    """A thin horizontal line — replaces heavy borders/boxes."""
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(left), Inches(top),
                                      Inches(left + width), Inches(top))
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_vrule(slide, left, top, height, color=HAIRLINE, weight=0.75):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                      Inches(left), Inches(top),
                                      Inches(left), Inches(top + height))
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_arrow(slide, x1, y1, x2, y2, *, color=INK_SOFT, weight=1.0,
              connector_type=MSO_CONNECTOR.STRAIGHT):
    conn = slide.shapes.add_connector(connector_type,
                                      Inches(x1), Inches(y1),
                                      Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    line_elem = conn.line._get_or_add_ln()
    tail_end = etree.SubElement(line_elem, qn("a:tailEnd"))
    tail_end.set("type", "triangle")
    tail_end.set("w", "sm")
    tail_end.set("h", "sm")
    return conn


def add_filled_rect(slide, left, top, width, height, *,
                    fill=BG_ALT, line=None, line_weight=0):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(left), Inches(top),
                                  Inches(width), Inches(height))
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Pt(line_weight)
    return rect


def add_rounded(slide, left, top, width, height, *,
                fill=BG_ALT, line=None, line_weight=0):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    if line is None:
        box.line.fill.background()
    else:
        box.line.color.rgb = line
        box.line.width = Pt(line_weight)
    return box


def add_diamond(slide, left, top, width, height, *,
                fill=HIGHLIGHT_BG, line=ACCENT, line_weight=1.0):
    d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(left), Inches(top),
                               Inches(width), Inches(height))
    d.fill.solid()
    d.fill.fore_color.rgb = fill
    d.line.color.rgb = line
    d.line.width = Pt(line_weight)
    return d


def set_slide_bg(slide, color=BG_CREAM):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(0), Inches(0),
                                Inches(SLIDE_W), Inches(SLIDE_H))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    # send to back by being added first
    return bg


def add_eyebrow(slide, text, *, top=0.45, color=ACCENT, font_size=10):
    """Small accented uppercase eyebrow at top-left, with a leading mark."""
    add_multitext(slide, MARGIN_L, top, 11.0, 0.25, [
        {"text": "◗ ", "size": font_size, "color": color,
         "font": F_BODY_REG, "bold": True},
        {"text": text.upper(), "size": font_size, "color": color,
         "font": F_BODY_REG, "bold": True, "char_spacing": 1.5},
    ])


def add_display_title(slide, text, *, top=0.85, width=11.0,
                      font_size=32, italic=True, color=INK,
                      line_spacing=1.05):
    add_text(slide, MARGIN_L, top, width, 1.5, text,
             font_size=font_size, font=F_DISPLAY, italic=italic,
             color=color, line_spacing=line_spacing)


def add_subtitle(slide, text, *, top, width=11.0, font_size=14, color=INK_SOFT):
    add_text(slide, MARGIN_L, top, width, 0.5, text,
             font_size=font_size, font=F_BODY, italic=True, color=color)


def add_section_eyebrow(slide, text, left, top, *, width=8.0,
                        font_size=10, color=ACCENT):
    add_multitext(slide, left, top, width, 0.20, [
        {"text": text.upper(), "size": font_size, "color": color,
         "font": F_BODY_REG, "bold": True, "char_spacing": 1.5},
    ])


def add_number_marker(slide, text, left, top, *, width=0.6, height=0.6,
                      font_size=24, color=ACCENT):
    """Big numbered marker (01, 02, ...) — Cambria for elegant numerals."""
    add_text(slide, left, top, width, height, text,
             font_size=font_size, font=F_DISPLAY, italic=True,
             color=color, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)


def add_code_block(slide, left, top, width, height, content, *,
                   font_size=9, fill=CODE_BG):
    box = add_filled_rect(slide, left, top, width, height, fill=fill)
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.10)
    tf.margin_bottom = Inches(0.10)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    lines = content.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.05
        run = p.add_run()
        run.text = ln
        run.font.size = Pt(font_size)
        run.font.color.rgb = INK
        run.font.name = F_MONO
        p.space_after = Pt(0)
    return box


def add_footer(slide, slide_num, total):
    """Editorial footer — single hairline + small label/number."""
    add_hairline(slide, MARGIN_L, FOOTER_Y, SLIDE_W - MARGIN_L - MARGIN_R,
                 color=HAIRLINE, weight=0.5)
    add_text(slide, MARGIN_L, FOOTER_Y + 0.10, 6.0, 0.20,
             "Lab Software Story Generation  ·  Client Alignment",
             font_size=8, font=F_BODY, italic=True, color=QUIET)
    add_text(slide, SLIDE_W - MARGIN_R - 1.0, FOOTER_Y + 0.10, 1.0, 0.20,
             f"{slide_num:02d} / {total:02d}",
             font_size=9, font=F_DISPLAY, italic=True, color=ACCENT,
             align=PP_ALIGN.RIGHT)


# ============================================================================
# SLIDE 1 — Opening / tone-setter
# ============================================================================

def slide_one(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_eyebrow(slide, "Lab Software Story Generation")
    add_display_title(slide,
                      "From SOPs to a ready-to-use\ndev backlog — without\nan expert reviewing each item.",
                      top=0.95, font_size=34, line_spacing=1.05)
    add_subtitle(slide,
                 "A discipline-agnostic system for clinical-lab software development.",
                 top=3.05)
    add_hairline(slide, MARGIN_L, 3.55, SLIDE_W - MARGIN_L - MARGIN_R)

    # 3-step flow as numbered moments
    flow_y = 3.85
    step_x = [MARGIN_L, 5.20, 9.85]
    steps = [
        ("01", "SOPs", "Standard Operating Procedures"),
        ("02", "Our system", "reads, drafts, validates"),
        ("03", "Dev backlog", "Jira: Epics + Stories"),
    ]
    for i, (num, head, sub) in enumerate(steps):
        x = step_x[i]
        add_number_marker(slide, num, x, flow_y, font_size=22, color=ACCENT)
        add_text(slide, x + 0.7, flow_y + 0.05, 3.5, 0.45, head,
                 font_size=22, font=F_DISPLAY, italic=True, color=INK)
        add_text(slide, x + 0.7, flow_y + 0.55, 3.5, 0.40, sub,
                 font_size=11, font=F_BODY, italic=True, color=INK_SOFT)
        if i < 2:
            arrow_x = step_x[i] + 4.3
            add_arrow(slide, arrow_x, flow_y + 0.30,
                      step_x[i+1] - 0.10, flow_y + 0.30,
                      color=ACCENT, weight=1.5)

    # Why this matters — editorial pull
    add_section_eyebrow(slide, "Why this matters", MARGIN_L, 5.10)
    matters = [
        ("Today.",
         "Every new SOP → manual expert review → hand-written stories. Slow, costly, inconsistent."),
        ("With this system.",
         "SOPs in, structured stories out. Quality is built into the rules, not into one expert's review."),
        ("Designed to generalize.",
         "Same system, any clinical discipline. Each new discipline is a configuration change, not an engineering project."),
    ]
    for i, (head, body) in enumerate(matters):
        y = 5.40 + i * 0.50
        add_multitext(slide, MARGIN_L, y, 12.2, 0.45, [
            {"text": head + "  ", "size": 12, "font": F_BODY_REG,
             "bold": True, "color": INK},
            {"text": body, "size": 12, "font": F_BODY, "color": INK_SOFT},
        ], line_spacing=1.20)


# ============================================================================
# SLIDE 2 — What "prior discipline" means
# ============================================================================

def slide_two(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_eyebrow(slide, "The teaching reference")
    add_display_title(slide,
                      'What "prior discipline" really means.',
                      top=0.95, font_size=30)
    add_subtitle(slide,
                 "A discipline whose dev work has already been validated. For our project — Cytology.",
                 top=2.05)
    add_hairline(slide, MARGIN_L, 2.50, SLIDE_W - MARGIN_L - MARGIN_R)

    # 4 columns: 2 mandatory + 1 optional + 1 not-inherited
    col_y = 2.85
    cols = [
        ("01", "MANDATORY", "Theme starter set",
         "The vocabulary of workflow categories the new discipline classifies its content against. "
         "Pre-Analytic, Analytic, Reporting, QC, Compliance, …", ACCENT),
        ("02", "MANDATORY", "Epic starter set",
         "The high-level dev groupings the new discipline's draft epics match against. "
         "Specimen Receiving, Reporting, Billing, …", ACCENT),
        ("03", "OPTIONAL", "Story exemplars",
         "Validated Jira stories with SOP citations — used as few-shot examples by the Story Extractor. "
         "Improves style fidelity. The system runs without these.", QUIET),
    ]
    col_w = (SLIDE_W - MARGIN_L - MARGIN_R - 2 * 0.30) / 3
    for i, (num, tag, head, body, color) in enumerate(cols):
        x = MARGIN_L + i * (col_w + 0.30)
        add_number_marker(slide, num, x, col_y, font_size=22, color=color)
        add_section_eyebrow(slide, tag, x + 0.75, col_y + 0.10,
                            font_size=9, color=color)
        add_text(slide, x + 0.75, col_y + 0.32, col_w - 0.75, 0.55, head,
                 font_size=18, font=F_DISPLAY, italic=True, color=INK,
                 line_spacing=1.05)
        add_text(slide, x, col_y + 1.10, col_w, 1.50, body,
                 font_size=11, font=F_BODY, color=INK_SOFT, line_spacing=1.30)

    add_hairline(slide, MARGIN_L, 5.40, SLIDE_W - MARGIN_L - MARGIN_R)

    # Not inherited
    add_section_eyebrow(slide, "What's NOT inherited", MARGIN_L, 5.55)
    add_text(slide, MARGIN_L, 5.80, 12.2, 0.40,
             "Tests   ·   Roles   ·   The new discipline's own SOPs   ·   Tasks (out of scope)",
             font_size=12, font=F_BODY, italic=True, color=INK)

    # Sourcing footnote
    add_text(slide, MARGIN_L, 6.30, 12.2, 0.35,
             "Sourcing.  Theme + epic catalogs are snapshotted from Cytology's existing Connect build. "
             "Quality-check thresholds use defaults tuned by first-run telemetry — no holdout calibration needed.",
             font_size=10, font=F_BODY, italic=True, color=QUIET, line_spacing=1.30)


# ============================================================================
# SLIDE 3 — The agents
# ============================================================================

def slide_three(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_eyebrow(slide, "The cast")
    add_display_title(slide,
                      "Six agents do the work.",
                      top=0.95, font_size=32)
    add_subtitle(slide,
                 "Plus a panel of 5 reviewers that vote in parallel when admitting new categories.",
                 top=1.85)
    add_hairline(slide, MARGIN_L, 2.30, SLIDE_W - MARGIN_L - MARGIN_R)

    # Editorial agent list — left column number+name, right column description+timing
    agents = [
        ("01", "Theme Discovery",
         "Builds the new discipline's theme catalog by classifying sample SOPs against the prior discipline's themes; new themes emerge only with evidence.",
         "Once per new discipline"),
        ("02", "Epic Extractor",
         "Drafts high-level epics from each SOP; matches each draft against the prior discipline's epic catalog. Matches inherit; non-matches go to novelty review.",
         "Per SOP + batch"),
        ("03", "Story Extractor",
         "Drafts user stories from SOP chunks using schema, shape definitions, and closed-enum catalogs. Optionally retrieves examples from prior-discipline Jira.",
         "Per epic, per SOP"),
        ("04", "Validator",
         "Quality-checks every emitted story against shape-specific rules (MUST/SHALL for capabilities; concrete values for configurations; named artifact for cleanups). Failures auto-park.",
         "Twice per story"),
        ("05", "Cross-SOP Synthesis",
         "After all SOPs are processed, finds patterns recurring across ≥ 2 SOPs and lifts them into broader \"capability stories\" with parameters.",
         "Once per batch"),
        ("06", "Dependency Resolver",
         "Resolves dependencies between stories into a topological order — the ordered backlog the dev team consumes.",
         "Once per batch"),
    ]
    row_y = 2.55
    row_h = 0.62
    for i, (num, name, desc, when) in enumerate(agents):
        y = row_y + i * row_h
        # number marker
        add_text(slide, MARGIN_L, y + 0.08, 0.55, 0.40, num,
                 font_size=18, font=F_DISPLAY, italic=True, color=ACCENT)
        # name
        add_text(slide, MARGIN_L + 0.65, y + 0.10, 2.50, 0.30, name,
                 font_size=13, font=F_BODY_REG, bold=True, color=INK)
        # description
        add_text(slide, MARGIN_L + 3.30, y + 0.05, 7.50, 0.55, desc,
                 font_size=10, font=F_BODY, color=INK_SOFT, line_spacing=1.20)
        # when
        add_text(slide, MARGIN_L + 11.00, y + 0.10, 1.80, 0.30, when,
                 font_size=10, font=F_BODY, italic=True, color=ACCENT,
                 align=PP_ALIGN.RIGHT)
        # hairline between rows (except last)
        if i < len(agents) - 1:
            add_hairline(slide, MARGIN_L, y + row_h - 0.02,
                         SLIDE_W - MARGIN_L - MARGIN_R,
                         color=HAIRLINE, weight=0.4)

    # Quorum panel callout — full-width band at bottom
    cy = row_y + len(agents) * row_h + 0.10
    add_filled_rect(slide, MARGIN_L, cy, SLIDE_W - MARGIN_L - MARGIN_R, 0.70,
                    fill=BG_ALT)
    add_section_eyebrow(slide, "Quorum panel — 5 reviewers in parallel",
                        MARGIN_L + 0.20, cy + 0.10)
    add_text(slide, MARGIN_L + 0.20, cy + 0.32, 12.0, 0.35,
             "Used inside Theme Discovery and Epic Extractor when admitting new categories. Each reviewer asks a different question. Need 3 of 5 agreeing on the same action; their descriptions must be similar enough to count.",
             font_size=10, font=F_BODY, color=INK_SOFT, line_spacing=1.20)


# ============================================================================
# SLIDE 4 — Three phases
# ============================================================================

def slide_four(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_eyebrow(slide, "End to end")
    add_display_title(slide,
                      "Three phases.",
                      top=0.95, font_size=36)
    add_subtitle(slide,
                 "Setup once, run per SOP, combine across SOPs.",
                 top=1.95)
    add_hairline(slide, MARGIN_L, 2.45, SLIDE_W - MARGIN_L - MARGIN_R)

    phases = [
        ("01", "Setup",
         "Build theme & epic catalogs for the new discipline, conditioned on the prior.",
         "Theme Discovery, Epic Extractor, Quorum panel.",
         "Once per new discipline (re-runs only on alarm)."),
        ("02", "Per-SOP",
         "Read each SOP, draft stories, quality-check.",
         "Epic Extractor, Story Extractor, Validator.",
         "Runs for every new SOP."),
        ("03", "Combine + ship",
         "Look across all SOPs, lift recurring patterns, generate outputs.",
         "Cross-SOP Synthesis, Validator, Dependency Resolver.",
         "Runs after a batch of SOPs is processed."),
    ]
    px = MARGIN_L
    pw = (SLIDE_W - MARGIN_L - MARGIN_R - 2 * 0.40) / 3
    py = 2.80
    for i, (num, title, desc, agents, cadence) in enumerate(phases):
        x = px + i * (pw + 0.40)
        # giant number
        add_text(slide, x, py, 1.5, 1.0, num,
                 font_size=64, font=F_DISPLAY, italic=True, color=ACCENT)
        # title
        add_text(slide, x, py + 1.10, pw, 0.50, title,
                 font_size=22, font=F_DISPLAY, italic=True, color=INK)
        # description
        add_text(slide, x, py + 1.65, pw, 1.00, desc,
                 font_size=12, font=F_BODY, color=INK, line_spacing=1.30)
        # agents
        add_text(slide, x, py + 2.65, pw, 0.40, "Agents",
                 font_size=9, font=F_BODY_REG, bold=True, color=ACCENT,
                 char_spacing=1.5)
        add_text(slide, x, py + 2.85, pw, 0.50, agents,
                 font_size=10, font=F_BODY, italic=True, color=INK_SOFT,
                 line_spacing=1.20)
        # cadence
        add_text(slide, x, py + 3.40, pw, 0.40, "Cadence",
                 font_size=9, font=F_BODY_REG, bold=True, color=ACCENT,
                 char_spacing=1.5)
        add_text(slide, x, py + 3.60, pw, 0.50, cadence,
                 font_size=10, font=F_BODY, italic=True, color=INK_SOFT,
                 line_spacing=1.20)
        # vertical hairline between phases
        if i < 2:
            add_vrule(slide, x + pw + 0.20, py + 0.10, 4.0)

    # Continuous monitoring band
    cy = 6.40
    add_section_eyebrow(slide, "Continuous (alongside all phases)",
                        MARGIN_L, cy)
    add_text(slide, MARGIN_L, cy + 0.22, 12.2, 0.35,
             "Drift report   ·   alarm-driven re-runs when too many items go to the review pile.",
             font_size=11, font=F_BODY, italic=True, color=INK_SOFT)


# ============================================================================
# SLIDE 5 — Phase 1 setup
# ============================================================================

def slide_five(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_eyebrow(slide, "Phase 1 — Setup")
    add_display_title(slide,
                      "Building the theme catalog.",
                      top=0.85, font_size=28)
    add_subtitle(slide,
                 "Run once per new discipline. Result: a catalog mostly inherited from the prior + 0–N new themes (each with evidence).",
                 top=1.65)
    add_hairline(slide, MARGIN_L, 2.20, SLIDE_W - MARGIN_L - MARGIN_R)

    # ----- LEFT: editorial step list -----
    lx = MARGIN_L
    lw = 6.6
    add_section_eyebrow(slide, "How it runs", lx, 2.40)
    steps = [
        "Read ~20 sample SOPs from the new discipline.",
        "Tag each chunk with stage, test, role.",
        "Compare each chunk to the prior discipline's themes — match → inherit; no match → set aside as residual.",
        "Cluster the residual into novel-theme candidates.",
        "Check inherited themes that have very little new-discipline evidence — candidates for dropping.",
        "5-reviewer quorum panel votes on each candidate (admit / merge / set aside).",
        "Result: new discipline's theme catalog (mostly inherited + 0–N new themes).",
    ]
    sy = 2.65
    for i, s in enumerate(steps):
        y = sy + i * 0.42
        # numbered marker
        add_text(slide, lx, y, 0.40, 0.30, f"{i+1:02d}",
                 font_size=12, font=F_DISPLAY, italic=True, color=ACCENT)
        add_text(slide, lx + 0.50, y, lw - 0.50, 0.40, s,
                 font_size=10.5, font=F_BODY, color=INK, line_spacing=1.30)

    # ----- RIGHT: chunk distribution feature -----
    rx = lx + lw + 0.30
    rw = SLIDE_W - rx - MARGIN_R
    add_section_eyebrow(slide, "Pass 1 — what came back", rx, 2.40)
    add_text(slide, rx, 2.65, rw, 0.50,
             "600 chunks classified",
             font_size=22, font=F_DISPLAY, italic=True, color=INK)

    # bar bands: inherited
    bands_y = 3.25
    band_h = 0.60
    inh_x = rx
    inh_w = rw * 0.62
    res_x = inh_x + inh_w + 0.25
    res_w = rw - inh_w - 0.25

    # inherited band
    add_filled_rect(slide, inh_x, bands_y, inh_w, band_h, fill=BG_ALT)
    add_text(slide, inh_x, bands_y - 0.30, inh_w, 0.25,
             "INHERITED FROM CYTOLOGY", font_size=8, font=F_BODY_REG,
             bold=True, color=INK_SOFT, char_spacing=1.5)

    # split inherited into 8 sub-buckets visually
    inherited = [("G1", 145), ("G2", 88), ("G3", 60), ("G4", 52),
                 ("G5", 44), ("G6", 30), ("G7", 38), ("G8", 20)]
    sub_w = inh_w / 8
    for i, (name, cnt) in enumerate(inherited):
        bx = inh_x + i * sub_w
        if i > 0:
            add_vrule(slide, bx, bands_y + 0.05, band_h - 0.10,
                      color=BG_CREAM, weight=1.0)
        add_text(slide, bx, bands_y + 0.05, sub_w, 0.25, name,
                 font_size=10, font=F_BODY_REG, bold=True, color=INK,
                 align=PP_ALIGN.CENTER)
        add_text(slide, bx, bands_y + 0.30, sub_w, 0.25, str(cnt),
                 font_size=11, font=F_DISPLAY, italic=True, color=ACCENT,
                 align=PP_ALIGN.CENTER)

    add_text(slide, inh_x, bands_y + band_h + 0.05, inh_w, 0.20,
             "477 chunks", font_size=10, font=F_BODY, italic=True,
             color=INK_SOFT, align=PP_ALIGN.CENTER)

    # residual band
    add_filled_rect(slide, res_x, bands_y, res_w, band_h, fill=HIGHLIGHT_BG)
    add_text(slide, res_x, bands_y - 0.30, res_w, 0.25,
             "RESIDUAL → PASS 2", font_size=8, font=F_BODY_REG,
             bold=True, color=ACCENT, char_spacing=1.5)
    add_text(slide, res_x, bands_y + 0.05, res_w, 0.25, "—",
             font_size=10, font=F_BODY_REG, color=INK_SOFT,
             align=PP_ALIGN.CENTER)
    add_text(slide, res_x, bands_y + 0.20, res_w, 0.40, "123",
             font_size=22, font=F_DISPLAY, italic=True, color=ACCENT,
             align=PP_ALIGN.CENTER)
    add_text(slide, res_x, bands_y + band_h + 0.05, res_w, 0.20,
             "feeds Pass 2", font_size=10, font=F_BODY, italic=True,
             color=ACCENT, align=PP_ALIGN.CENTER)

    # output preview — what the catalog looks like
    out_y = 4.65
    add_section_eyebrow(slide, "Output — micro_v1 catalog", rx, out_y)
    yaml = """inherited:  G1, G2, G3, G4, G5, G6, G7, G8
admitted:   MI1 Susceptibility Testing
            MI2 Biosafety Containment
unclassified: 38 chunks (6.3%)"""
    add_code_block(slide, rx, out_y + 0.25, rw, 1.30, yaml, font_size=10)

    # multi-label footnote
    add_text(slide, MARGIN_L, 6.30, 12.2, 0.50,
             "Pass 1 here is single-label — each chunk goes to its best-matching theme above τ_match=0.65 (or to residual). Runtime story / chunk theme tags are multi-label per D3 (themes are soft tags).",
             font_size=9, font=F_BODY, italic=True, color=QUIET,
             line_spacing=1.30)


# ============================================================================
# SLIDE 6 — Phase 2 per-SOP
# ============================================================================

def slide_six(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_eyebrow(slide, "Phase 2 — Per-SOP")
    add_display_title(slide,
                      "Each SOP goes through the same pipeline.",
                      top=0.85, font_size=26)
    add_subtitle(slide,
                 "N SOPs run in parallel — N independent pipeline instances.",
                 top=1.75)
    add_hairline(slide, MARGIN_L, 2.30, SLIDE_W - MARGIN_L - MARGIN_R)

    # ----- LEFT: pipeline ForEach container -----
    lx = MARGIN_L
    lw = 6.0
    # ForEach container — subtle frame around the pipeline
    container = add_filled_rect(slide, lx, 2.55, lw, 4.30, fill=BG_ALT)
    container.line.color.rgb = HAIRLINE
    container.line.width = Pt(1.0)
    try:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        container.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    except Exception:
        pass

    add_text(slide, lx + 0.15, 2.45, 4.0, 0.22,
             "ForEach SOP — N instances run in parallel",
             font_size=9, font=F_BODY_REG, bold=True, color=ACCENT,
             char_spacing=1.5)

    # the pipeline steps
    pipe_steps = [
        ("SOP arrives", "start"),
        ("Tag chunks", "process"),
        ("Epic Extractor", "agent"),
        ("Story Extractor", "agent"),
        ("Validator", "agent"),
    ]
    sy = 2.70
    sh = 0.45
    sx = lx + 0.40
    sw = lw - 0.80
    for i, (label, kind) in enumerate(pipe_steps):
        y = sy + i * (sh + 0.10)
        if kind == "start":
            fill = BG_CREAM
            ink = INK
        elif kind == "agent":
            fill = HIGHLIGHT_BG
            ink = INK
        else:
            fill = WHITE
            ink = INK
        rect = add_rounded(slide, sx, y, sw, sh, fill=fill,
                           line=HAIRLINE, line_weight=0.5)
        _set_text(rect, label, font_size=12, font=F_BODY_REG, bold=True,
                  color=ink, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < len(pipe_steps) - 1:
            add_arrow(slide, sx + sw/2, y + sh + 0.02,
                      sx + sw/2, y + sh + 0.16, color=INK_SOFT, weight=1.0)

    # decision diamond + branches at the bottom
    dy = sy + len(pipe_steps) * (sh + 0.10) + 0.02
    dh = 0.45
    dw = 1.6
    dx = sx + (sw - dw) / 2
    add_diamond(slide, dx, dy, dw, dh)
    _set_text(slide.shapes[len(slide.shapes)-1],
              "Pass quality?", font_size=9, font=F_BODY_REG, bold=True,
              color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # branches: yes (down), no (right)
    end_y = dy + dh + 0.10
    end_h = 0.32
    add_rounded(slide, sx, end_y, sw * 0.55, end_h, fill=BG_CREAM,
                line=SUCCESS, line_weight=0.75)
    _set_text(slide.shapes[-1], "kept", font_size=10, font=F_BODY_REG,
              bold=True, color=SUCCESS, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE)
    rb_x = sx + sw * 0.55 + 0.20
    rb_w = sw - sw * 0.55 - 0.20
    add_rounded(slide, rb_x, end_y, rb_w, end_h, fill=BG_CREAM,
                line=FAILURE, line_weight=0.75)
    _set_text(slide.shapes[-1], "park (revise ≤2)", font_size=10,
              font=F_BODY_REG, bold=True, color=FAILURE,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ----- RIGHT: story output snippet + 4 quality-rule chips -----
    rx = lx + lw + 0.40
    rw = SLIDE_W - rx - MARGIN_R

    add_section_eyebrow(slide, "Output — a single story", rx, 2.45)
    yaml = """story:
  shape:    capability
  title:    "Validate received specimens
             against open orders"
  persona:  lims
  stage:    accessioning_verification
  tests:    [gram_stain, blood_culture]
  AC: "When specimen received with
       accession, system MUST match
       against open order"
  source:   SOP-MICRO-014, lines 23-31
  quality:  passed"""
    add_code_block(slide, rx, 2.70, rw, 2.50, yaml, font_size=9.5)

    # Quality rules — 4 chips
    add_section_eyebrow(slide, "The 4 quality rules — one per shape",
                        rx, 5.30)
    chip_y = 5.55
    chip_h = 0.65
    chips = [
        ("Capability",    "MUST/SHALL\n+ params"),
        ("Stage-split",   "Stage in title\n+ siblings"),
        ("Config-instance","Concrete typed values"),
        ("Cleanup",       "Named artifact\n+ before/after"),
    ]
    chip_w = (rw - 3 * 0.10) / 4
    for i, (head, body) in enumerate(chips):
        x = rx + i * (chip_w + 0.10)
        # header band
        add_filled_rect(slide, x, chip_y, chip_w, 0.22, fill=ACCENT_SOFT)
        add_text(slide, x, chip_y + 0.02, chip_w, 0.20, head,
                 font_size=9, font=F_BODY_REG, bold=True, color=INK,
                 align=PP_ALIGN.CENTER)
        # body
        add_text(slide, x + 0.05, chip_y + 0.27, chip_w - 0.10, chip_h - 0.27,
                 body, font_size=8.5, font=F_BODY, color=INK_SOFT,
                 align=PP_ALIGN.CENTER, line_spacing=1.20)


# ============================================================================
# SLIDE 7 — Phase 3 combine + monitor
# ============================================================================

def slide_seven(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_eyebrow(slide, "Phase 3 — Combine + monitor")
    add_display_title(slide,
                      "Patterns lift across SOPs.",
                      top=0.95, font_size=30)
    add_subtitle(slide,
                 "Patterns recurring in ≥ 2 SOPs become broader \"capability stories\" with parameters.",
                 top=1.80)
    add_hairline(slide, MARGIN_L, 2.25, SLIDE_W - MARGIN_L - MARGIN_R)

    # ----- LEFT: synthesis steps editorial style -----
    lx = MARGIN_L
    lw = 5.6
    add_section_eyebrow(slide, "What synthesis does", lx, 2.45)
    steps = [
        "Look across all kept stories from all SOPs.",
        "Cluster on (test, role, stage, shape, behavior).",
        "Patterns in ≥ 2 SOPs → lift to capability story.",
        "Validator gate 2 on the lifted capability.",
        "Dependency Resolver → final ordered backlog.",
    ]
    sy = 2.75
    for i, s in enumerate(steps):
        y = sy + i * 0.50
        add_text(slide, lx, y, 0.40, 0.30, f"{i+1:02d}",
                 font_size=14, font=F_DISPLAY, italic=True, color=ACCENT)
        add_text(slide, lx + 0.55, y, lw - 0.55, 0.50, s,
                 font_size=11, font=F_BODY, color=INK, line_spacing=1.30)

    # the ≥ 2 SOPs lift visual at the bottom of the left side
    lift_y = 5.60
    add_section_eyebrow(slide, "Worked example — ≥ 2 SOPs lift", lx, lift_y)
    box_h = 0.40
    src_w = 1.55
    add_rounded(slide, lx, lift_y + 0.30, src_w, box_h, fill=BG_ALT,
                line=HAIRLINE, line_weight=0.5)
    _set_text(slide.shapes[-1], "Story · SOP-007", font_size=9,
              font=F_BODY_REG, bold=True, color=INK,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rounded(slide, lx, lift_y + 0.80, src_w, box_h, fill=BG_ALT,
                line=HAIRLINE, line_weight=0.5)
    _set_text(slide.shapes[-1], "Story · SOP-014", font_size=9,
              font=F_BODY_REG, bold=True, color=INK,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # cluster
    cl_x = lx + src_w + 0.25
    cl_w = 1.7
    add_rounded(slide, cl_x, lift_y + 0.50, cl_w, box_h + 0.2, fill=HIGHLIGHT_BG,
                line=ACCENT, line_weight=0.75)
    _set_text(slide.shapes[-1], "Cluster: ≥ 2 SOPs", font_size=10,
              font=F_BODY_REG, bold=True, color=ACCENT,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # capability
    cp_x = cl_x + cl_w + 0.25
    cp_w = lw - (cp_x - lx)
    add_rounded(slide, cp_x, lift_y + 0.50, cp_w, box_h + 0.2, fill=BG_CREAM,
                line=SUCCESS, line_weight=0.75)
    _set_text(slide.shapes[-1], "Capability story\n+ parameters", font_size=10,
              font=F_BODY_REG, bold=True, color=SUCCESS,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # arrows
    add_arrow(slide, lx + src_w + 0.02, lift_y + 0.30 + box_h/2,
              cl_x - 0.02, lift_y + 0.50 + (box_h+0.2)*0.4,
              color=INK_SOFT, weight=1.0)
    add_arrow(slide, lx + src_w + 0.02, lift_y + 0.80 + box_h/2,
              cl_x - 0.02, lift_y + 0.50 + (box_h+0.2)*0.6,
              color=INK_SOFT, weight=1.0)
    add_arrow(slide, cl_x + cl_w + 0.02, lift_y + 0.50 + (box_h+0.2)/2,
              cp_x - 0.02, lift_y + 0.50 + (box_h+0.2)/2,
              color=ACCENT, weight=1.5)

    # ----- RIGHT: capability story example -----
    rx = lx + lw + 0.40
    rw = SLIDE_W - rx - MARGIN_R
    add_section_eyebrow(slide, "Capability story output", rx, 2.45)
    yaml = """capability_story:
  id:        STORY-MICRO-CAP-0007
  title:     "Configurable critical-result
              notification across tests"

  parameters:                # what varies
    - test:           enum
    - threshold:      number+units
    - persona_owner:  supervisor

  AC:
    - when:  "{test} flagged critical"
      then: "notify {persona_owner}
             within 1 hour"

  child_stories:            # the concrete cases
    - STORY-MICRO-0023
    - STORY-MICRO-0041
    - STORY-MICRO-0089
  source_sops: [SOP-014, SOP-019, SOP-022]"""
    add_code_block(slide, rx, 2.70, rw, 4.20, yaml, font_size=9)


# ============================================================================
# SLIDE 8 — 4 story shapes, one example each
# ============================================================================

def slide_eight(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_eyebrow(slide, "Story shapes")
    add_display_title(slide,
                      "Four shapes, one example each.",
                      top=0.85, font_size=28)
    add_subtitle(slide,
                 "Different shapes have different rules. Examples show what makes each shape that shape.",
                 top=1.75)
    add_hairline(slide, MARGIN_L, 2.25, SLIDE_W - MARGIN_L - MARGIN_R)

    grid_x = MARGIN_L
    grid_y = 2.50
    h_gap = 0.30
    v_gap = 0.30
    card_w = (SLIDE_W - MARGIN_L - MARGIN_R - h_gap) / 2
    card_h = (FOOTER_Y - grid_y - v_gap - 0.10) / 2

    examples = [
        ("01", "CAPABILITY",
         "MUST/SHALL  +  configurable parameters",
         """title:    "Validate received specimens
           against open orders"
persona:  lims
stage:    accessioning_verification
tests:    [gram_stain, blood_culture]

AC: When specimen received with
    accession, system MUST match
    against open order
AC: When no match within 24h, system
    MUST flag for review"""),
        ("02", "WORKFLOW-STAGE-SPLIT",
         "Stage in title  +  sibling stories enumerated",
         """title:    "PHI update — after results
           are reported"
persona:  supervisor
stage:    reporting_case_closure

AC: When PHI corrected on finalised
    report, system MUST regenerate
    with audit trail entry

cross_links: STORY-0076, STORY-0077"""),
        ("03", "CONFIGURATION-INSTANCE",
         "Concrete typed values  ·  no MUST/SHALL pretense",
         """title:    "Blood culture incubation
           — bottle parameters"
tests:    [blood_culture_incubation]
persona:  null    stage: null

AC: When bottle received, configure
    incubation
    expected_value:
      temperature_c:    35
      duration_h:       120
      agitation_rpm:    220"""),
        ("04", "CLEANUP",
         "Named artifact  +  before/after observable",
         """title:    "Remove obsolete 'Pending
           Pathologist Review' status
           from results screen"
persona:  null    stage: null

AC: When viewing the results screen,
    the option MUST NOT appear
AC: When existing case has the status,
    MUST migrate to 'Pending Review'"""),
    ]
    for i, (num, name, feature, body) in enumerate(examples):
        col = i % 2
        row = i // 2
        x = grid_x + col * (card_w + h_gap)
        y = grid_y + row * (card_h + v_gap)
        # number marker + name (eyebrow style)
        add_text(slide, x, y, 0.55, 0.30, num,
                 font_size=14, font=F_DISPLAY, italic=True, color=ACCENT)
        add_text(slide, x + 0.65, y + 0.05, card_w - 0.65, 0.30, name,
                 font_size=11, font=F_BODY_REG, bold=True, color=INK,
                 char_spacing=1.5)
        # feature subtitle
        add_text(slide, x, y + 0.35, card_w, 0.25, feature,
                 font_size=10, font=F_BODY, italic=True, color=INK_SOFT)
        # hairline
        add_hairline(slide, x, y + 0.65, card_w)
        # body code block
        add_code_block(slide, x, y + 0.75, card_w, card_h - 0.75, body,
                       font_size=8.5)


# ============================================================================
# SLIDE 9 — Inside Validator + Dependency Resolver
# ============================================================================

def slide_nine(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_eyebrow(slide, "Two key agents")
    add_display_title(slide,
                      "Inside the Validator and the Dependency Resolver.",
                      top=0.85, font_size=24)
    add_subtitle(slide,
                 "What these two agents actually do — beyond the one-line description on slide 3.",
                 top=1.85)
    add_hairline(slide, MARGIN_L, 2.30, SLIDE_W - MARGIN_L - MARGIN_R)

    # ===== LEFT: Validator =====
    lx = MARGIN_L
    lw = 6.20
    add_section_eyebrow(slide, "Validator", lx, 2.50, color=ACCENT)
    add_text(slide, lx, 2.70, lw, 0.40,
             "Quality-checks every story.",
             font_size=20, font=F_DISPLAY, italic=True, color=INK)
    add_text(slide, lx, 3.18, lw, 0.30,
             "Runs at gate 1 (after Story Extractor) and gate 2 (after Synthesis).",
             font_size=10, font=F_BODY, italic=True, color=INK_SOFT)

    v_steps = [
        "01", "Receives a draft story.",
        "02", "Closed-enum check  (tests, persona, stage in catalog?)",
        "03", "Shape verification  (declared shape matches content?)",
        "04", "Apply shape-specific rubric  (4 rules — see slide 8)",
    ]
    sy = 3.65
    for i in range(0, len(v_steps), 2):
        num = v_steps[i]
        body = v_steps[i + 1]
        y = sy + (i // 2) * 0.42
        add_text(slide, lx, y, 0.40, 0.30, num,
                 font_size=12, font=F_DISPLAY, italic=True, color=ACCENT)
        add_text(slide, lx + 0.50, y, lw - 0.50, 0.40, body,
                 font_size=11, font=F_BODY, color=INK, line_spacing=1.30)

    # decision: pass → kept; fail → revise/park
    dec_y = sy + 4 * 0.42 + 0.10
    add_text(slide, lx, dec_y, 0.40, 0.30, "→",
             font_size=14, font=F_DISPLAY, color=ACCENT)
    add_multitext(slide, lx + 0.50, dec_y + 0.02, lw - 0.50, 0.40, [
        {"text": "Pass?  ", "size": 11, "font": F_BODY_REG, "bold": True,
         "color": INK},
        {"text": "yes → kept (Output A)   ·   ", "size": 11,
         "font": F_BODY, "color": SUCCESS},
        {"text": "no → revise (≤2) then auto-park", "size": 11,
         "font": F_BODY, "color": FAILURE},
    ], line_spacing=1.30)

    # vertical hairline divider
    add_vrule(slide, lx + lw + 0.20, 2.50, 4.30)

    # ===== RIGHT: Dependency Resolver =====
    rx = lx + lw + 0.40
    rw = SLIDE_W - rx - MARGIN_R
    add_section_eyebrow(slide, "Dependency Resolver", rx, 2.50, color=ACCENT)
    add_text(slide, rx, 2.70, rw, 0.40,
             "Sorts the backlog.",
             font_size=20, font=F_DISPLAY, italic=True, color=INK)
    add_text(slide, rx, 3.18, rw, 0.55,
             "Each story has dependencies[] and cross_links[]. The resolver builds a DAG and produces a topological order — a sprint-ready backlog.",
             font_size=10, font=F_BODY, italic=True, color=INK_SOFT,
             line_spacing=1.30)

    # DAG
    s_w, s_h = 0.95, 0.36
    panel_cx = rx + rw / 2
    a_x = panel_cx - s_w / 2
    a_y = 3.95
    add_rounded(slide, a_x, a_y, s_w, s_h, fill=BG_ALT,
                line=HAIRLINE, line_weight=0.5)
    _set_text(slide.shapes[-1], "A", font_size=12, font=F_BODY_REG, bold=True,
              color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bc_y = a_y + s_h + 0.30
    b_x = rx + 0.4
    c_x = rx + rw - s_w - 0.4
    add_rounded(slide, b_x, bc_y, s_w, s_h, fill=BG_ALT,
                line=HAIRLINE, line_weight=0.5)
    _set_text(slide.shapes[-1], "B", font_size=12, font=F_BODY_REG, bold=True,
              color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rounded(slide, c_x, bc_y, s_w, s_h, fill=BG_ALT,
                line=HAIRLINE, line_weight=0.5)
    _set_text(slide.shapes[-1], "C", font_size=12, font=F_BODY_REG, bold=True,
              color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    d_y = bc_y + s_h + 0.30
    d_x = panel_cx - s_w / 2
    add_rounded(slide, d_x, d_y, s_w, s_h, fill=BG_ALT,
                line=HAIRLINE, line_weight=0.5)
    _set_text(slide.shapes[-1], "D", font_size=12, font=F_BODY_REG, bold=True,
              color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    add_arrow(slide, a_x + s_w/2, a_y + s_h, b_x + s_w/2, bc_y,
              color=INK_SOFT, weight=0.8)
    add_arrow(slide, a_x + s_w/2, a_y + s_h, c_x + s_w/2, bc_y,
              color=INK_SOFT, weight=0.8)
    add_arrow(slide, b_x + s_w/2, bc_y + s_h, d_x + s_w/2, d_y,
              color=INK_SOFT, weight=0.8)
    add_arrow(slide, c_x + s_w/2, bc_y + s_h, d_x + s_w/2, d_y,
              color=INK_SOFT, weight=0.8)

    # output
    out_y = d_y + s_h + 0.30
    add_filled_rect(slide, rx, out_y, rw, 0.35, fill=HIGHLIGHT_BG)
    add_text(slide, rx, out_y + 0.05, rw, 0.25,
             "Output:   1. A   →   2. B   →   3. C   →   4. D",
             font_size=11, font=F_BODY_REG, bold=True, color=ACCENT,
             align=PP_ALIGN.CENTER)


# ============================================================================
# SLIDE 10 — Validator at work
# ============================================================================

def slide_ten(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_eyebrow(slide, "Validator at work")
    add_display_title(slide,
                      "A sample story, checked twice.",
                      top=0.85, font_size=26)
    add_subtitle(slide,
                 "Failure on attempt 01 — pass on revise (attempt 02).",
                 top=1.75)
    add_hairline(slide, MARGIN_L, 2.25, SLIDE_W - MARGIN_L - MARGIN_R)

    grid_y = 2.55
    h_gap = 0.40
    card_w = (SLIDE_W - MARGIN_L - MARGIN_R - h_gap) / 2

    # ===== ATTEMPT 01 (left, fail) =====
    a1_x = MARGIN_L
    add_text(slide, a1_x, grid_y, 0.55, 0.40, "01",
             font_size=20, font=F_DISPLAY, italic=True, color=FAILURE)
    add_multitext(slide, a1_x + 0.65, grid_y + 0.10, card_w - 0.65, 0.30, [
        {"text": "ATTEMPT 1   ", "size": 11, "font": F_BODY_REG,
         "bold": True, "color": FAILURE, "char_spacing": 1.5},
        {"text": "draft", "size": 11, "font": F_BODY, "italic": True,
         "color": INK_SOFT},
    ])
    add_hairline(slide, a1_x, grid_y + 0.55, card_w, color=FAILURE,
                 weight=1.0)

    a1_yaml = """story (draft):
  shape:    capability
  title:    "Validate accession appropriately"
  AC: "When specimen received, system
       MUST validate appropriately"
  source:   SOP-MICRO-014, lines 23-31"""
    add_code_block(slide, a1_x, grid_y + 0.70, card_w, 1.55, a1_yaml,
                   font_size=9)

    # checks
    chk_y = grid_y + 2.40
    add_section_eyebrow(slide, "Validator checks", a1_x, chk_y)
    a1_checks = [
        ("✓", "1. Closed-enum check passes", SUCCESS),
        ("✓", "2. Shape verification passes", SUCCESS),
        ("✗", "3. Shape rubric — 'appropriately' is ambiguous", FAILURE),
    ]
    for i, (mark, text, color) in enumerate(a1_checks):
        y = chk_y + 0.28 + i * 0.28
        add_text(slide, a1_x, y, 0.30, 0.25, mark,
                 font_size=14, font=F_BODY_REG, bold=True, color=color)
        add_text(slide, a1_x + 0.35, y + 0.03, card_w - 0.35, 0.30, text,
                 font_size=10, font=F_BODY, color=INK)

    # verdict
    vy1 = chk_y + 0.28 + 3 * 0.28 + 0.10
    add_filled_rect(slide, a1_x, vy1, card_w, 0.40, fill=HIGHLIGHT_BG)
    add_text(slide, a1_x, vy1 + 0.08, card_w, 0.25,
             "→ Revise (1 of 2) — Story Extractor reworks the AC",
             font_size=11, font=F_BODY_REG, bold=True, color=FAILURE,
             align=PP_ALIGN.CENTER)

    # ===== ATTEMPT 02 (right, pass) =====
    a2_x = a1_x + card_w + h_gap
    add_text(slide, a2_x, grid_y, 0.55, 0.40, "02",
             font_size=20, font=F_DISPLAY, italic=True, color=SUCCESS)
    add_multitext(slide, a2_x + 0.65, grid_y + 0.10, card_w - 0.65, 0.30, [
        {"text": "ATTEMPT 2   ", "size": 11, "font": F_BODY_REG,
         "bold": True, "color": SUCCESS, "char_spacing": 1.5},
        {"text": "revised", "size": 11, "font": F_BODY, "italic": True,
         "color": INK_SOFT},
    ])
    add_hairline(slide, a2_x, grid_y + 0.55, card_w, color=SUCCESS,
                 weight=1.0)

    a2_yaml = """story (revised):
  shape:    capability
  title:    "Validate accession against open order"
  AC: "When specimen received with accession,
       system MUST match against open order;
       if no match within 24h, MUST flag for
       accessioning review"
  source:   SOP-MICRO-014, lines 23-31"""
    add_code_block(slide, a2_x, grid_y + 0.70, card_w, 1.55, a2_yaml,
                   font_size=9)

    add_section_eyebrow(slide, "Validator checks", a2_x, chk_y)
    a2_checks = [
        ("✓", "1. Closed-enum check passes", SUCCESS),
        ("✓", "2. Shape verification passes", SUCCESS),
        ("✓", "3. Shape rubric — concrete, observable, parameterized (24h)", SUCCESS),
    ]
    for i, (mark, text, color) in enumerate(a2_checks):
        y = chk_y + 0.30 + i * 0.32
        add_text(slide, a2_x, y, 0.30, 0.25, mark,
                 font_size=14, font=F_BODY_REG, bold=True, color=color)
        add_text(slide, a2_x + 0.35, y + 0.03, card_w - 0.35, 0.30, text,
                 font_size=10, font=F_BODY, color=INK)

    add_filled_rect(slide, a2_x, vy1, card_w, 0.40, fill=HIGHLIGHT_BG)
    add_text(slide, a2_x, vy1 + 0.08, card_w, 0.25,
             "→ Story kept (validated) → Output A",
             font_size=11, font=F_BODY_REG, bold=True, color=SUCCESS,
             align=PP_ALIGN.CENTER)

    # bottom note
    note_y = vy1 + 0.50
    add_text(slide, MARGIN_L, note_y, 12.2, 0.30,
             "Closed-enum violations (e.g. invalid persona) hard-reject without entering the revise loop. Only shape/AC issues consume revisions. Failures after 2 revisions auto-park (slide 11).",
             font_size=9, font=F_BODY, italic=True, color=QUIET,
             line_spacing=1.20)


# ============================================================================
# SLIDE 11 — Failure handling
# ============================================================================

def slide_eleven(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_eyebrow(slide, "Failure handling")
    add_display_title(slide,
                      "What happens when something fails.",
                      top=0.95, font_size=28)
    add_subtitle(slide,
                 "Failures at three levels — each has a deterministic fallback. No human intervention at runtime.",
                 top=1.85)
    add_hairline(slide, MARGIN_L, 2.40, SLIDE_W - MARGIN_L - MARGIN_R)

    # 3 columns — clean editorial style
    cy = 2.65
    ch = 3.40
    cw = (SLIDE_W - MARGIN_L - MARGIN_R - 2 * 0.30) / 3
    panels = [
        ("01", "Per-call",     "LLM error or timeout",
         "Retry with exponential backoff (3 attempts).",
         "If persistent, defer the item to the next batch; log to audit.",
         "Other items continue. No block."),
        ("02", "Per-story",    "Quality check fails",
         "Up to 2 revision attempts; the Validator drives revisions with the failed-checks list.",
         "If still failing, auto-park to the review pile (excluded from Output A in strict mode).",
         "Subsequent stories continue. No block."),
        ("03", "Per-batch",    "Drift / catalog issues",
         "Auto-rerun discovery with τ_match −0.05 (more permissive); re-cluster G0; quorum re-evaluates.",
         "If τ_match floor (0.50) hit, drift report flags for manual catalog re-curation.",
         "Pipeline continues. Drift report is informational."),
    ]
    for i, (num, head, sub, fall, persist, impact) in enumerate(panels):
        x = MARGIN_L + i * (cw + 0.30)
        # number
        add_text(slide, x, cy, 0.50, 0.40, num,
                 font_size=20, font=F_DISPLAY, italic=True, color=ACCENT)
        # header
        add_text(slide, x + 0.65, cy + 0.10, cw - 0.65, 0.30, head,
                 font_size=14, font=F_BODY_REG, bold=True, color=INK)
        # sub
        add_text(slide, x + 0.65, cy + 0.40, cw - 0.65, 0.25, sub,
                 font_size=9, font=F_BODY, italic=True, color=INK_SOFT,
                 char_spacing=1.0)
        # hairline
        add_hairline(slide, x, cy + 0.75, cw)
        # rows: Fallback / If persistent / Pipeline impact
        rows = [("Fallback", fall),
                ("If persistent", persist),
                ("Pipeline impact", impact)]
        for j, (label, body) in enumerate(rows):
            ry = cy + 0.90 + j * 0.85
            add_text(slide, x, ry, cw, 0.20, label.upper(),
                     font_size=8, font=F_BODY_REG, bold=True, color=ACCENT,
                     char_spacing=1.5)
            add_text(slide, x, ry + 0.22, cw, 0.55, body,
                     font_size=10, font=F_BODY, color=INK_SOFT,
                     line_spacing=1.30)
        # vertical divider between panels (except last)
        if i < 2:
            add_vrule(slide, x + cw + 0.15, cy + 0.10, ch - 0.20)

    # 3 queues banner at bottom
    qy = cy + ch + 0.20
    add_section_eyebrow(slide, "Three queues", MARGIN_L, qy)
    queue_items = [
        ("Review pile", "quality / rubric failures (parked stories)"),
        ("G0 / E0 buckets", "classification failures (low-confidence chunks/drafts)"),
        ("Audit trail", "LLM-call failures after retry exhaustion"),
    ]
    qx = MARGIN_L
    qw = (SLIDE_W - MARGIN_L - MARGIN_R - 2 * 0.20) / 3
    for i, (name, body) in enumerate(queue_items):
        x = qx + i * (qw + 0.20)
        add_multitext(slide, x, qy + 0.25, qw, 0.45, [
            {"text": name + "  ", "size": 11, "font": F_BODY_REG,
             "bold": True, "color": INK},
            {"text": body, "size": 10, "font": F_BODY, "italic": True,
             "color": INK_SOFT},
        ], line_spacing=1.30)


# ============================================================================
# SLIDE 12 — What ships + alignment ask
# ============================================================================

def slide_twelve(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_eyebrow(slide, "Closing")
    add_display_title(slide,
                      "What ships, and what we need.",
                      top=0.85, font_size=28)
    add_subtitle(slide,
                 "Three deliverables per discipline. The test catalog and sample SOPs are the only blocking inputs.",
                 top=1.75)
    add_hairline(slide, MARGIN_L, 2.30, SLIDE_W - MARGIN_L - MARGIN_R)

    # ===== Ships row — 3 deliverable cards =====
    sy = 2.50
    sh = 1.20
    sw = (SLIDE_W - MARGIN_L - MARGIN_R - 2 * 0.30) / 3
    deliverables = [
        ("Jira backlog",       "Epics + Stories. Ready for the dev team. Includes traceability to source SOPs."),
        ("Configuration files","Per-test settings (e.g. blood culture parameters). Deterministic structure per discipline."),
        ("Reference + audit",  "Theme + epic catalogs, ANALOGY map, review pile, decision log."),
    ]
    add_section_eyebrow(slide, "What ships", MARGIN_L, sy - 0.25)
    for i, (head, body) in enumerate(deliverables):
        x = MARGIN_L + i * (sw + 0.30)
        add_filled_rect(slide, x, sy, sw, sh, fill=BG_ALT)
        add_text(slide, x + 0.20, sy + 0.15, sw - 0.40, 0.40, head,
                 font_size=15, font=F_DISPLAY, italic=True, color=INK)
        add_text(slide, x + 0.20, sy + 0.55, sw - 0.40, 0.60, body,
                 font_size=10, font=F_BODY, color=INK_SOFT,
                 line_spacing=1.30)

    # ===== Inputs section — 3 categorized columns =====
    iy = sy + sh + 0.35
    ih = 2.10
    add_section_eyebrow(slide, "Inputs to the system", MARGIN_L, iy - 0.25)
    cw = (SLIDE_W - MARGIN_L - MARGIN_R - 2 * 0.20) / 3
    cols = [
        ("Already captured  ✓",
         "from your project documentation",
         [
             "Persona catalog  —  from the Roles & Process Map document",
             "Universal 6-stage enum  —  from the Process Map",
             "ANALOGY persona_links  —  from the cross-discipline role inventory",
         ],
         SUCCESS),
        ("From the seminal prior  ✓",
         "snapshotted from Cytology's Connect build",
         [
             "Theme catalog (cyto_v1)  —  warm-start prior for Theme Discovery",
             "Epic catalog (cyto_epic_v1)  —  used by the conditioned Epic Extractor",
         ],
         INK_SOFT),
        ("Still needed",
         "the only blocking inputs",
         [
             "Test catalog (micro_test_v1.yaml)  —  list of in-scope Microbiology tests",
             "Sample SOPs  (~20 representative procedure documents)",
         ],
         ACCENT),
    ]
    for i, (head, sub, items, color) in enumerate(cols):
        x = MARGIN_L + i * (cw + 0.20)
        add_text(slide, x, iy, cw, 0.30, head,
                 font_size=12, font=F_BODY_REG, bold=True, color=color)
        add_text(slide, x, iy + 0.30, cw, 0.25, sub,
                 font_size=9, font=F_BODY, italic=True, color=INK_SOFT)
        add_hairline(slide, x, iy + 0.60, cw, color=color, weight=0.75)
        for j, item in enumerate(items):
            iy_item = iy + 0.75 + j * 0.42
            add_multitext(slide, x, iy_item, cw, 0.40, [
                {"text": "·  ", "size": 11, "font": F_BODY_REG,
                 "color": color, "bold": True},
                {"text": item, "size": 10, "font": F_BODY,
                 "color": INK, "line_spacing": 1.30},
            ], line_spacing=1.30)

    # alignment ask — pulled-quote style
    ay = iy + ih + 0.05
    add_filled_rect(slide, MARGIN_L, ay, SLIDE_W - MARGIN_L - MARGIN_R, 0.55,
                    fill=HIGHLIGHT_BG)
    add_text(slide, MARGIN_L + 0.30, ay + 0.13, 12.0, 0.30,
             "Alignment ask:  the test catalog + sample SOPs are the gating inputs. Everything else is in hand. Confirm?",
             font_size=12, font=F_DISPLAY, italic=True, color=ACCENT,
             align=PP_ALIGN.CENTER)


# ============================================================================
# MAIN
# ============================================================================

def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    slide_one(prs)
    slide_two(prs)
    slide_three(prs)
    slide_four(prs)
    slide_five(prs)
    slide_six(prs)
    slide_seven(prs)
    slide_eight(prs)
    slide_nine(prs)
    slide_ten(prs)
    slide_eleven(prs)
    slide_twelve(prs)

    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, start=1):
        add_footer(slide, i, total)

    out = Path(__file__).parent / "client_alignment_deck_refined.pptx"
    prs.save(str(out))
    print(f"wrote: {out}  ({total} slides)")


if __name__ == "__main__":
    main()
